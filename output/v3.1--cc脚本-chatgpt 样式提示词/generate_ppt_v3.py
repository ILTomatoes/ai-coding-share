#!/usr/bin/env python3
"""PPT v3 生成器 — 极简白色主题
用法: python generate_ppt_v3.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os

# ══════════════════════════════════════════════════════
# 配色方案 — 极简白色主题
# ══════════════════════════════════════════════════════
BG_WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
BG_LIGHT   = RGBColor(0xF8, 0xF9, 0xFA)
PRIMARY    = RGBColor(0x1E, 0x29, 0x3B)     # 深蓝黑
ACCENT     = RGBColor(0x25, 0x63, 0xEB)     # 蓝色强调
ACCENT2    = RGBColor(0xEA, 0x58, 0x0C)     # 橙色强调
ACCENT3    = RGBColor(0x16, 0xA3, 0x4A)     # 绿色
ACCENT4    = RGBColor(0xDC, 0x26, 0x26)     # 红色
TEXT_P     = RGBColor(0x1E, 0x29, 0x3B)     # 正文
TEXT_S     = RGBColor(0x64, 0x74, 0x8B)     # 次要文字
TEXT_L     = RGBColor(0x94, 0xA3, 0xB8)     # 浅灰
DIVIDER    = RGBColor(0xE2, 0xE8, 0xF0)
CARD_BG    = RGBColor(0xF1, 0xF5, 0xF9)
CARD_BD    = RGBColor(0xCB, 0xD5, 0xE1)
FN = '微软雅黑'

# ══════════════════════════════════════════════════════
# 全局设置
# ══════════════════════════════════════════════════════
SLIDE_W, SLIDE_H = 13.333, 7.5
ML, MR, MT = 1.3, 1.3, 0.8  # 左/右/上边距 (>=10%)
CW = SLIDE_W - ML - MR       # 内容宽度 ~10.7
CH = SLIDE_H - MT - 0.8      # 内容高度

prs = Presentation()
prs.slide_width = Inches(SLIDE_W)
prs.slide_height = Inches(SLIDE_H)
_pn = [0]  # 页码计数器

# ══════════════════════════════════════════════════════
# 基础辅助函数
# ══════════════════════════════════════════════════════

def set_bg(slide, c=BG_WHITE):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = c

def tb(slide, l, t, w, h):
    return slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))

def p(tf, text, sz=20, c=TEXT_P, bold=False, align=PP_ALIGN.LEFT, sa=Pt(8), fn=FN, ls=None):
    pg = tf.paragraphs[-1] if len(tf.paragraphs)==1 and tf.paragraphs[0].text=='' else tf.add_paragraph()
    pg.text = text; pg.font.size=Pt(sz); pg.font.color.rgb=c
    pg.font.bold=bold; pg.font.name=fn; pg.alignment=align; pg.space_after=sa
    if ls: pg.line_spacing = ls
    return pg

def box(slide, l, t, w, h, fill=None, line=None, lw=1, rounded=False):
    st = MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
    s = slide.shapes.add_shape(st, Inches(l), Inches(t), Inches(w), Inches(h))
    if fill: s.fill.solid(); s.fill.fore_color.rgb = fill
    else: s.fill.background()
    if line: s.line.color.rgb=line; s.line.width=Pt(lw)
    else: s.line.fill.background()
    return s

def oval(slide, l, t, w, h, fill=None):
    s = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(l), Inches(t), Inches(w), Inches(h))
    if fill: s.fill.solid(); s.fill.fore_color.rgb=fill
    else: s.fill.background()
    s.line.fill.background()
    return s

def ft(slide):
    _pn[0] += 1
    box(slide, ML, 6.9, CW, 0.005, fill=DIVIDER)
    t = tb(slide, SLIDE_W/2-1, 6.95, 2, 0.35)
    p(t.text_frame, str(_pn[0]), sz=10, c=TEXT_L, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════
# 页面类型生成器
# ══════════════════════════════════════════════════════

def cover_slide(title, subtitle='', info=''):
    sl = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(sl)
    box(sl, 1.0, 2.5, 0.06, 2.5, fill=ACCENT)
    t = tb(sl, 1.5, 2.5, 10, 1.5); t.text_frame.word_wrap=True
    p(t.text_frame, title, sz=40, bold=True, c=PRIMARY)
    box(sl, 1.5, 4.1, 3.0, 0.04, fill=ACCENT)
    if subtitle:
        t2=tb(sl,1.5,4.4,10,0.6); t2.text_frame.word_wrap=True
        p(t2.text_frame, subtitle, sz=20, c=TEXT_S)
    if info:
        t3=tb(sl,1.5,5.2,10,0.4)
        p(t3.text_frame, info, sz=14, c=TEXT_L)
    ft(sl); return sl

def section_slide(label, title, subtitle=''):
    sl = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(sl, BG_LIGHT)
    box(sl, 1.3, 1.5, 0.06, 4.5, fill=ACCENT)
    if label:
        t0=tb(sl,1.8,2.0,8,0.4); p(t0.text_frame, label, sz=14, c=ACCENT, bold=True)
    t=tb(sl,1.8,2.6,9,1.2); t.text_frame.word_wrap=True
    p(t.text_frame, title, sz=36, bold=True, c=PRIMARY)
    box(sl, 1.8, 3.9, 4.0, 0.04, fill=ACCENT)
    if subtitle:
        t2=tb(sl,1.8,4.2,9,0.6); t2.text_frame.word_wrap=True
        p(t2.text_frame, subtitle, sz=18, c=TEXT_S)
    ft(sl); return sl

def content_slide(title, bullets):
    sl = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(sl)
    tt=tb(sl,ML,MT,CW,0.6); p(tt.text_frame,title,sz=32,bold=True,c=PRIMARY)
    box(sl, ML, MT+0.65, 2.5, 0.03, fill=ACCENT)
    t=tb(sl,ML,MT+1.0,CW,CH-1.2); t.text_frame.word_wrap=True
    for b in bullets:
        if isinstance(b, tuple):
            txt,lv=b[0], b[1] if len(b)>1 else 0
            clr=b[2] if len(b)>2 else TEXT_S
            ind='    '*lv; pre='• ' if lv==0 else '◦ '
            p(t.text_frame, ind+pre+txt, sz=20 if lv==0 else 18, c=clr, sa=Pt(6), ls=Pt(28))
        elif b=='':
            p(t.text_frame,'',sz=8,sa=Pt(4))
        elif b.startswith('##'):
            p(t.text_frame, b[2:].strip(), sz=22, bold=True, c=ACCENT, sa=Pt(6))
        elif b.startswith('>>'):
            p(t.text_frame, '▸ '+b[2:].strip(), sz=20, bold=True, c=ACCENT, sa=Pt(6))
        elif b.startswith('!!'):
            p(t.text_frame, '⚠ '+b[2:].strip(), sz=20, bold=True, c=ACCENT4, sa=Pt(6))
        else:
            p(t.text_frame, '• '+b, sz=20, c=TEXT_P, sa=Pt(6), ls=Pt(28))
    ft(sl); return sl

def comparison_slide(title, lt, li, rt, ri, lc=ACCENT, rc=ACCENT2):
    sl = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(sl)
    tt=tb(sl,ML,MT,CW,0.6); p(tt.text_frame,title,sz=32,bold=True,c=PRIMARY)
    box(sl,ML,MT+0.65,2.5,0.03,fill=ACCENT)
    cw=5.0; cy=1.8; ch=4.8
    box(sl,ML,cy,cw,ch,fill=CARD_BG,line=lc,lw=2,rounded=True)
    box(sl,ML,cy,cw,0.05,fill=lc)
    t1=tb(sl,ML+0.2,cy+0.2,cw-0.4,0.4); p(t1.text_frame,lt,sz=20,bold=True,c=lc)
    t2=tb(sl,ML+0.2,cy+0.7,cw-0.4,ch-1.0); t2.text_frame.word_wrap=True
    for x in li: p(t2.text_frame,'• '+x,sz=18,c=TEXT_P,sa=Pt(4))
    rx=ML+cw+0.7
    box(sl,rx,cy,cw,ch,fill=CARD_BG,line=rc,lw=2,rounded=True)
    box(sl,rx,cy,cw,0.05,fill=rc)
    t3=tb(sl,rx+0.2,cy+0.2,cw-0.4,0.4); p(t3.text_frame,rt,sz=20,bold=True,c=rc)
    t4=tb(sl,rx+0.2,cy+0.7,cw-0.4,ch-1.0); t4.text_frame.word_wrap=True
    for x in ri: p(t4.text_frame,'• '+x,sz=18,c=TEXT_P,sa=Pt(4))
    ft(sl); return sl

def emphasis_slide(big, sub='', bc=ACCENT, bs=48):
    sl = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(sl, BG_LIGHT)
    t=tb(sl,2.0,2.2,9.3,2.0); t.text_frame.word_wrap=True
    p(t.text_frame, big, sz=bs, bold=True, c=bc, align=PP_ALIGN.CENTER)
    if sub:
        t2=tb(sl,2.0,4.3,9.3,1.0); t2.text_frame.word_wrap=True
        p(t2.text_frame, sub, sz=20, c=TEXT_S, align=PP_ALIGN.CENTER)
    box(sl,5.5,5.5,2.3,0.03,fill=ACCENT)
    ft(sl); return sl

def data_slide(title, items):
    sl = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(sl)
    tt=tb(sl,ML,MT,CW,0.6); p(tt.text_frame,title,sz=32,bold=True,c=PRIMARY)
    box(sl,ML,MT+0.65,2.5,0.03,fill=ACCENT)
    n=len(items); cols=min(n,3); cw=3.2; ch=2.2; gap=0.5
    tw=cols*cw+(cols-1)*gap; sx=(SLIDE_W-tw)/2; sy=2.0
    for i,(v,l) in enumerate(items):
        r,c2=divmod(i,cols); x=sx+c2*(cw+gap); y=sy+r*(ch+0.4)
        box(sl,x,y,cw,ch,fill=CARD_BG,line=CARD_BD,lw=1,rounded=True)
        tv=tb(sl,x,y+0.3,cw,0.7); p(tv.text_frame,v,sz=36,bold=True,c=ACCENT,align=PP_ALIGN.CENTER)
        tl=tb(sl,x+0.1,y+1.2,cw-0.2,0.8); tl.text_frame.word_wrap=True
        for ln in l.split('\n'): p(tl.text_frame,ln,sz=14,c=TEXT_S,align=PP_ALIGN.CENTER,sa=Pt(2))
    ft(sl); return sl

def timeline_slide(title, items):
    sl = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(sl)
    tt=tb(sl,ML,MT,CW,0.6); p(tt.text_frame,title,sz=32,bold=True,c=PRIMARY)
    box(sl,ML,MT+0.65,2.5,0.03,fill=ACCENT)
    n=len(items); sp=CW/(n-1) if n>1 else 0; ly=3.3
    box(sl,ML,ly,CW,0.04,fill=ACCENT)
    for i,(lb,desc,det) in enumerate(items):
        x=ML+i*sp
        oval(sl,x-0.15,ly-0.15,0.3,0.3,fill=ACCENT)
        tu=tb(sl,x-1.0,1.7,2.2,1.0); tu.text_frame.word_wrap=True
        p(tu.text_frame,lb,sz=14,bold=True,c=ACCENT,align=PP_ALIGN.CENTER)
        for ln in desc.split('\n'): p(tu.text_frame,ln,sz=12,c=TEXT_S,align=PP_ALIGN.CENTER,sa=Pt(2))
        cww=min(2.4,sp*0.8 if sp>0 else 2.4); chh=2.3; cx=x-cww/2; cy2=ly+0.5
        box(sl,cx,cy2,cww,chh,fill=CARD_BG,line=CARD_BD,lw=1,rounded=True)
        td=tb(sl,cx+0.1,cy2+0.15,cww-0.2,chh-0.3); td.text_frame.word_wrap=True
        for ln in det.split('\n'): p(td.text_frame,ln,sz=13,c=TEXT_S,align=PP_ALIGN.CENTER,sa=Pt(2))
    ft(sl); return sl

def table_slide(title, headers, rows):
    sl = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(sl)
    tt=tb(sl,ML,MT,CW,0.6); p(tt.text_frame,title,sz=32,bold=True,c=PRIMARY)
    box(sl,ML,MT+0.65,2.5,0.03,fill=ACCENT)
    nc=len(headers); nr=len(rows)+1; rh=0.5
    tbl=sl.shapes.add_table(nr,nc,Inches(ML),Inches(1.8),Inches(CW),Inches(rh*nr)).table
    for i,h in enumerate(headers):
        c=tbl.cell(0,i); c.text=h
        for pg in c.text_frame.paragraphs:
            pg.font.size=Pt(16); pg.font.bold=True; pg.font.color.rgb=RGBColor(0xFF,0xFF,0xFF)
            pg.font.name=FN; pg.alignment=PP_ALIGN.LEFT
        c.fill.solid(); c.fill.fore_color.rgb=ACCENT
        c.margin_left=Inches(0.15); c.margin_top=Inches(0.08); c.margin_bottom=Inches(0.08)
    for r,row in enumerate(rows):
        for ci,val in enumerate(row):
            c=tbl.cell(r+1,ci); c.text=val
            for pg in c.text_frame.paragraphs:
                pg.font.size=Pt(14); pg.font.color.rgb=TEXT_P; pg.font.name=FN
                if ci==0: pg.font.bold=True; pg.font.color.rgb=ACCENT
            c.fill.solid(); c.fill.fore_color.rgb=BG_WHITE if r%2==0 else CARD_BG
            c.margin_left=Inches(0.15); c.margin_top=Inches(0.08); c.margin_bottom=Inches(0.08)
    ft(sl); return sl

def card_slide(title, cards):
    sl = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(sl)
    tt=tb(sl,ML,MT,CW,0.6); p(tt.text_frame,title,sz=32,bold=True,c=PRIMARY)
    box(sl,ML,MT+0.65,2.5,0.03,fill=ACCENT)
    n=len(cards); gap=0.3; cww=(CW-(n-1)*gap)/n; chh=4.5; sy=1.8
    for i,(ct,cc,clr) in enumerate(cards):
        x=ML+i*(cww+gap)
        box(sl,x,sy,cww,chh,fill=CARD_BG,line=clr,lw=2,rounded=True)
        box(sl,x,sy,cww,0.05,fill=clr)
        tct=tb(sl,x+0.15,sy+0.15,cww-0.3,0.4); tct.text_frame.word_wrap=True
        p(tct.text_frame,ct,sz=18,bold=True,c=clr)
        tcc=tb(sl,x+0.15,sy+0.7,cww-0.3,chh-1.0); tcc.text_frame.word_wrap=True
        for ln in cc.split('\n'):
            if ln.strip(): p(tcc.text_frame,ln.strip(),sz=14,c=TEXT_S,sa=Pt(3))
    ft(sl); return sl

def diagram_slide(title, code_text):
    sl = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(sl)
    tt=tb(sl,ML,MT,CW,0.6); p(tt.text_frame,title,sz=32,bold=True,c=PRIMARY)
    box(sl,ML,MT+0.65,2.5,0.03,fill=ACCENT)
    box(sl,ML,1.7,CW,5.0,fill=CARD_BG,line=CARD_BD,lw=1,rounded=True)
    t=tb(sl,ML+0.3,1.9,CW-0.6,4.6); t.text_frame.word_wrap=True
    for i,ln in enumerate(code_text.split('\n')):
        pg=t.text_frame.paragraphs[0] if i==0 else t.text_frame.add_paragraph()
        pg.text=ln; pg.font.size=Pt(12); pg.font.color.rgb=TEXT_P
        pg.font.name='Consolas'; pg.space_after=Pt(0); pg.line_spacing=Pt(14)
    ft(sl); return sl

# ══════════════════════════════════════════════════════
# 52 页幻灯片内容
# ══════════════════════════════════════════════════════

def create_slides():

    # ══════════════════════════════════════════
    # 封面 & 开场 (P1-P2)
    # ══════════════════════════════════════════

    # ── P1: 封面 ──
    cover_slide(
        'Java 开发中的 AI 协作编程实践',
        subtitle='从真实项目说起',
        info='演讲人 / 2026',
    )

    # ── P2: 开场 ──
    content_slide(
        '关于这次分享',
        [
            '三个真实 Java 项目的 AI 协作编程经历',
            '有成功的喜悦，也有翻车的教训',
            '>> 目标：听完之后，你能带走一些可以马上用的东西',
        ],
    )

    # ══════════════════════════════════════════
    # 第一部分：从真实项目说起 (P3-P14)
    # ══════════════════════════════════════════

    # ── P3: Part 1 标题页 ──
    section_slide(
        '第一部分',
        '从真实项目说起',
        subtitle='两个案例，两种结局',
    )

    # ── P4: 案例一 — aip-server 项目概览 ──
    content_slide(
        '案例一：aip-server 项目概览',
        [
            'AI 智能平台后端服务',
            ('企业级智能体管理、知识库、数据分析平台', 1),
            '100 个 Controller · 11 个业务模块 · 2394 次提交',
            '技术栈：Spring Boot 2.7 + MyBatis Plus + PostgreSQL + Java 17',
            '多数据源：PostgreSQL + Neo4j + Dremio',
        ],
    )

    # ── P5: 案例一 — 聚焦 logic/tool 技能工具模块 ──
    data_slide(
        '案例一：聚焦 logic/tool 技能工具模块',
        [
            ('5', '个 Controller\nLtCodeRepo / LtMcp / LtExtensionPlugin\nLtOpenapi / LtTaskflow'),
            ('~80', '个 API 端点'),
            ('8', '个核心实体'),
            ('16', '个数据库迁移版本\nv2.1.0 → v2.2.14'),
        ],
    )

    # ── P6: 案例一 — 协作模式 ──
    content_slide(
        '案例一：协作模式 — 人与 AI 的分工',
        [
            '>> 人设计原型 → 人设计表结构（AI 参与讨论）→ AI 编码 → 人评审 & 测试',
            '',
            '## 关键原则',
            '人负责"定方向"（设计、评审、验收）',
            'AI 负责"跑量"（编码实现）',
            '及时评审，小步验证',
        ],
    )

    # ── P7: 案例一 — 关键成果 ──
    emphasis_slide(
        '5 天 80 个 API',
        sub='29 次提交（3月4-11日）· 13 次 fix 占 45% · 效率提升 2-4 倍',
        bs=52,
    )

    # ── P8: 案例一 — 小结 ──
    emphasis_slide(
        '设计到位 + 边界清晰 = AI 效率拉满',
        sub='表结构设计清晰 → AI 有明确边界\n及时评审，小步验证\n人把控方向，AI 堆砌工作量',
    )

    # ── P9: 案例二 — A2A 适配器项目背景 ──
    content_slide(
        '案例二：A2A 适配器 — 项目背景',
        [
            '开发 A2A 协议适配器（aip-gateway 的 a2a-wrapper 插件）',
            ('相对陌生的领域，没有深入调研', 0, TEXT_S),
            '让 AI 根据 A2A 协议文档直接编码',
            '!! 关键问题：自己对 A2A 不了解，无法判断 AI 输出是否正确',
        ],
    )

    # ── P10: 案例二 — 发生了什么 ──
    timeline_slide(
        '案例二：发生了什么',
        [
            ('第1天', '让 AI 编码', '根据 A2A 协议文档\n直接编码'),
            ('第1-2周', '大量产出', '46 次提交 · 16 天\n无法验证正确性'),
            ('转折', '发现官方 SDK', 'a2aproject/a2a-java\n有 Spring WebFlux 示例'),
            ('推倒重来', 'SDK 3 天搞定', '141,026 行删除\n1,786 行新增\n代码缩减 98.7%'),
        ],
    )

    # ── P11: 案例二 — 翻车现场（一）造轮子 ──
    comparison_slide(
        '案例二：翻车现场（一）— 造轮子',
        'AI 自研实现',
        [
            '从零编码，看似完整',
            '代码量大',
            '无法验证正确性',
        ],
        '官方 SDK',
        [
            '3 天完成',
            '代码少一个数量级',
            '稳定性有保障',
        ],
    )

    # ── P12: 案例二 — 翻车现场（二）更多教训 ──
    card_slide(
        '案例二：翻车现场（二）— 更多教训',
        [
            (
                'SQL 版本号错误',
                '生成 V2.2.14 而非 V2.2.13\n没有先查看现有版本文件\n数据库迁移执行失败',
                ACCENT,
            ),
            (
                '修改已冻结的 Flyway 脚本',
                '创建补丁 SQL 时\nAI 修改了已冻结的 init.sql\n已部署环境报错',
                ACCENT2,
            ),
            (
                '权限过滤理解偏差',
                'AI 保留不该有的后端逻辑\n"前端可见性控制权限"策略\n理解不深',
                ACCENT4,
            ),
        ],
    )

    # ── P13: 案例二 — 小结 ──
    emphasis_slide(
        '方向错了，效率越高代价越大',
        sub='对领域不了解 → 没有先建立认知 → 盲目信任 AI → AI 高效跑偏',
        bc=ACCENT4,
    )

    # ── P14: Part 1 收尾 — 两个案例对比 ──
    comparison_slide(
        'Part 1 收尾：两个案例的对比',
        '✅ aip-server',
        [
            '设计清晰 + 边界明确',
            '人定方向，AI 跑量',
            '一周搞定，质量可靠',
        ],
        '❌ A2A 适配器',
        [
            '方向不明 + 盲目信任',
            '没有调研就动手',
            '17 天返工',
        ],
        lc=ACCENT3,
        rc=ACCENT4,
    )

    # ══════════════════════════════════════════
    # 第二部分：反思 (P15-P21)
    # ══════════════════════════════════════════

    # ── P15: Part 2 标题页 ──
    section_slide(
        '第二部分',
        '反思——从案例到洞察',
    )

    # ── P16: 从案例中看到了什么 ──
    content_slide(
        '从案例中看到了什么',
        [
            'aip-server 成功：人给了清晰的设计和边界',
            'A2A 失败：人没有给方向',
            '>>核心观察：AI 的效率是确定的，不确定的是人的引导',
            '同一个 AI，同一个开发者，结果天壤之别',
        ],
    )

    # ── P17: AI 的潜力与边界 ──
    comparison_slide(
        'AI 的潜力与边界',
        '潜力 ✅',
        [
            '编码效率提升 2-4 倍',
            '擅长重复性工作、模式化代码',
            '趋势明确，只会越来越强',
        ],
        '边界 ⚠️',
        [
            '不了解领域背景时容易跑偏',
            '缺乏全局视角，可能局部最优全局灾难',
            '不会主动质疑你的设计',
        ],
    )

    # ── P18: 核心洞察：把最擅长的事交给 AI ──
    content_slide(
        '核心洞察：把最擅长的事交给 AI',
        [
            '!!常见想法：先把不擅长的交给 AI',
            '>>实际经验：先把最擅长的交给 AI',
            '擅长 = 你能准确验收 → 质量有保障',
            '这些工作对个人无成长价值，属于纯工作量堆叠',
            '释放时间 → 投入到研究、设计、架构等创造性工作',
        ],
    )

    # ── P19: 角色转变 ──
    content_slide(
        '角色转变',
        [
            '过去：设计 → 编码 → 测试（全自己做）',
            '现在：设计 → AI 编码 → 评审测试',
            '>>从"写代码的人"变成"设计系统的人"',
        ],
    )

    # ── P20: 陷阱：AI 会"PUA"你 ──
    content_slide(
        '陷阱：AI 会"PUA"你',
        [
            '永远自信满满，即使方向错误',
            '语气笃定、逻辑自洽，很容易被"说服"',
            '不会犹豫、不会说"我不确定"',
            'A2A 翻车：AI 全程自信地造了一个轮子',
            '权限过滤翻车：AI 的理由听起来"完全合理"',
        ],
    )

    # ── P21: 关键问题引出 ──
    emphasis_slide(
        '我们需要一套系统化的方法',
        sub='把擅长的事交给 AI → 但 AI 会 PUA 你 → 靠个人判断力够吗？\n单靠经验不可持续、不可复制 → 团队协作呢？',
        bc=ACCENT2,
    )

    # ══════════════════════════════════════════
    # 第三部分：SDD & Harness 实践 (P22-P41)
    # ══════════════════════════════════════════

    # ── P22: Part 3 标题页 ──
    section_slide('第三部分', '系统化提升 — SDD 与 Harness 实践', '')

    # ── P23: 什么是 SDD ──
    content_slide('什么是 SDD（规范驱动开发）', [
        'SDD = Specification-Driven Development',
        '传统开发：规范是文档，人看文档写代码',
        'AI 编程：规范是"AI 的行为边界"，直接决定输出质量',
        '没有规范 = 新员工自由发挥（不可控）',
        '>>AI 就是那个"超级新员工"',
    ])

    # ── P24: SDD 三层体现 ──
    card_slide('SDD 在 AI 编程中的三层体现', [
        ('项目级规范', 'CLAUDE.md + rules/\n所有会话共享，每次自动加载\n= 员工手册', ACCENT),
        ('任务级规范', 'Plan、PRD\n特定功能的约束\n= 任务说明书', ACCENT2),
        ('会话级引导', '对话中的约束和纠正\n实时交互中补充\n= 日常工作指导', ACCENT3),
    ])

    # ── P25: aip-portal 全貌 ──
    data_slide('案例：智能体门户项目（aip-portal）', [
        ('259次', '总提交数'),
        ('72天', '活跃开发'),
        ('281个', 'Java 文件'),
        ('7个', '业务模块'),
        ('4人', '协作开发'),
        ('24个', 'Worktree 并行'),
    ])

    # ── P26: Rules 演化总览 ──
    diagram_slide('Rules 演化总览：从混乱到精炼', r"""
264KB  |          ╱╲
~200KB |         ╱  \
       |        ╱    \
       |       ╱      \   ← paths: **.java 全部加载
       |      ╱        \
 ~50KB |     ╱          \
       |    ╱            \
 272行 |___╱______________╲___________
       P1   P2    P2.5  P3    P4

Phase 1: 什么都放（5487 行）
Phase 2: 越加越多（24 个文件，~264KB）
Phase 2.5: 发现根因（paths 导致全部加载）
Phase 3: 痛定思痛（参考文档移至 docs/）
Phase 4: 精炼到位（272 行，压缩 95%）
""")

    # ── P27: Rules Phase 1-2 ──
    comparison_slide(
        'Rules 演化：Phase 1-2 — 混乱与膨胀',
        'Phase 1（5487 行 20 个文件）',
        ['把所有规范都塞进去', 'Token 消耗巨大', 'AI 抓不住重点'],
        'Phase 2（~6000+ 行）',
        ['踩坑就加规则，越加越多', 'Token 爆炸', '恶性循环：犯错→加规则→再犯错'],
        lc=ACCENT4, rc=ACCENT2,
    )

    # ── P28: Rules Phase 3-4 ──
    comparison_slide(
        'Rules 演化：Phase 3-4 — 清理与精炼',
        'Phase 3（删掉 2714 行）',
        ['意识到规范不是越多越好', 'rules/ 是 token 预算管理', '详细文档移到 docs/'],
        'Phase 4（272 行，压缩 95%）',
        ['只保留 AI 必须知道的约束', '稳定后维护成本极低'],
        lc=ACCENT, rc=ACCENT3,
    )

    # ── P29: Rules 最终形态对比 ──
    comparison_slide(
        'Rules 最终形态对比',
        'Before（Phase 1）',
        ['所有规范堆在一个目录', '每次会话全部加载', '5487 行 20 个文件'],
        'After（Phase 4）',
        ['核心约束 272 行自动加载', '详细规范 docs/ 按需引用', 'ARCHITECTURE.md(124行)\n+ CODE-STYLE.md(14行)\n+ testing.md(134行)'],
    )

    # ── P30: Rules 演化关键经验 ──
    content_slide('Rules 演化的关键经验', [
        '>>rules/ 是 token 预算管理，不是 wiki',
        '规范不是越多越好，而是越精炼越好',
        '每次踩坑都是完善规范的机会',
        '仅 9 天完成重构（04-01 至 04-09），之后趋于稳定',
    ])

    # ── P31: SDD 落地 CLAUDE.md ──
    content_slide('SDD 落地：CLAUDE.md — 极简导航', [
        'aip-portal 的 CLAUDE.md 只有 23 行',
        '项目定位（一句话）',
        '文档索引（指向 rules/ 和 docs/）',
        '常用 Maven 命令',
        '>>不在 CLAUDE.md 里堆砌内容——它是导航页，不是百科全书',
    ])

    # ── P32: SDD 落地 ARCHITECTURE.md ──
    content_slide('SDD 落地：ARCHITECTURE.md — 架构约束', [
        '124 行，包含核心架构约定',
        'Maven 模块划分',
        'SQL 版本管理策略',
        '分页查询统一方案',
        'AI 每次会话都知道项目的核心架构约定',
    ])

    # ── P33: SDD 落地 CODE-STYLE.md ──
    content_slide('SDD 落地：CODE-STYLE.md — 编码规范索引', [
        '仅 14 行，不在 rules 里放编码细节',
        '索引到 docs/code-style/ 下的 8 个专题文档',
        '核心约束 → 自动加载（每次会话）',
        '详细规范 → 按需加载（AI 自行查阅）',
    ])

    # ── P34: 什么是 Harness Engineering ──
    card_slide('什么是 Harness Engineering', [
        ('规范文件', 'CLAUDE.md\nrules/', ACCENT),
        ('权限配置', 'settings.json\n自动化授权', ACCENT2),
        ('自定义技能', 'Skills\n工作流配置', ACCENT3),
        ('工作流', 'Plan、Worktree\nAgent Team', TEXT_S),
    ])

    # ── P35: Harness 三大支柱 ──
    content_slide('Harness 的三大支柱', [
        '>>上下文工程：给 AI 正确的信息',
        '>>架构约束：让 AI 在框架内工作',
        '>>熵管理：防止混乱累积',
        '目标：让 AI 在你的项目里工作得又快又准',
    ])

    # ── P36: 混乱期表现 ──
    content_slide('混乱期的表现（aip-portal 教训）', [
        'AI 每次会话都在"重新认识项目"，重复犯错',
        'rules 膨胀到 6000+ 行，AI 反而更"懵"',
        '犯错→加规则→再犯错→恶性循环',
        'settings.local.json 累积 58 条允许规则',
        '!!没有 Harness 的 AI 编程，靠的是个人"手感"',
    ])

    # ── P37: 转折的四个关键动作 ──
    card_slide('转折的四个关键动作', [
        ('规范精简', '6000+ 行 → 272 行\n大清理', ACCENT),
        ('分层架构', 'rules/ 核心约束\ndocs/ 按需引用', ACCENT),
        ('复盘机制', '踩坑不只是修 bug\n而是完善规范', ACCENT3),
        ('权限收敛', '只保留必要权限\n减少误操作风险', ACCENT2),
    ])

    # ── P38: 知识体系三层架构 ──
    diagram_slide('知识体系三层架构', r"""
🌐 全局层 (~/.claude/)
   └─ 通用规则、偏好设置、MCP 工具配置
   └─ 所有项目共享

📁 项目层 (项目 .claude/)
   ├─ CLAUDE.md (23行) — 项目导航
   ├─ rules/ (272行) — 核心约束
   │   ├─ ARCHITECTURE.md (124行)
   │   ├─ CODE-STYLE.md (14行)
   │   └─ testing.md (134行)
   └─ settings.local.json — 权限配置

💬 会话层 (对话中)
   └─ 临时约束、特定任务指导
   └─ 一次性使用
""")

    # ── P39: 系统化之后的变化 ──
    content_slide('系统化之后的变化', [
        'AI 编码准确率稳定提升',
        '重复性错误逐步减少',
        '新会话快速上手（知识已沉淀在 Harness 中）',
        'rules 相关提交集中在 4 月初（仅 9 天）',
        '>>Harness 到了稳定期后，维护成本很低',
    ])

    # ── P40: 补充案例 aip-gateway ──
    content_slide('补充案例：aip-gateway — 系统化的开发流程', [
        '59 个 Java 文件',
        'AMG 核心模块仅 2 天完成',
        '设计文档迭代 9 轮',
        '总提交 209 次',
        '>>Day 1 就搭建好 Harness，全程无混乱期',
    ])

    # ── P41: 项目对比 ──
    table_slide(
        '项目对比：被动建设 vs 主动建设',
        ['维度', 'aip-portal（被动）', 'aip-gateway（主动）'],
        [
            ['Harness 建设', '先踩坑再补规范', 'Day 1 就建规范'],
            ['混乱期', '约 2 个月', '基本没有'],
            ['配置风格', '分层（rules + docs）', '集中（单文件 CLAUDE.md）'],
            ['提交数', '259 次', '209 次'],
        ],
    )

    # ══════════════════════════════════════════
    # 第四部分：工具链与实战技巧 (P42-P52)
    # ══════════════════════════════════════════

    # ── P42: 第四部分标题页 ──
    section_slide('第四部分', '工具链与实战技巧', '')

    # ── P43: Claude Code 工具介绍 ──
    content_slide('Claude Code — AI 编程终端工具', [
        '直接在终端工作，与项目代码零距离',
        '支持读取文件、编辑代码、执行命令，全流程闭环',
        '自动加载项目 CLAUDE.md 和 rules/，理解项目上下文',
        '支持 Plan 模式、Worktree 并行、Agent Team 协作',
    ])

    # ── P44: cc-switch 工具介绍 ──
    content_slide('cc-switch — 配置切换工具', [
        '快速切换不同的 Claude Code 配置',
        '不同项目需要不同配置',
        '复杂任务用强模型，简单任务用快模型',
        '团队成员共享配置',
    ])

    # ── P45: 技巧一 ──
    content_slide('技巧一：小步快跑，持续验证', [
        '设计 → 实现 → 测试，每完成一个模块就验证',
        'AI 生成快 → 错误积累也快 → 一次贪多反而更慢',
        '遇到问题再"磨刀"，不提前过度设计',
        '!!aip-portal 早期一次做大模块导致批量返工',
    ])

    # ── P46: 技巧二 ──
    comparison_slide(
        '技巧二：先建立认知，再让 AI 干活',
        '2 小时调研', [
            '先调研 A2A 官方仓库',
            '发现官方 SDK',
            '3 天搞定',
        ],
        '17 天返工', [
            '不了解协议就动手',
            'AI 造了一个轮子',
            '推倒重来',
        ],
        lc=ACCENT3, rc=ACCENT4,
    )

    # ── P47: 协作模式选择矩阵 ──
    table_slide(
        '协作模式选择矩阵',
        ['场景', '推荐模式', '说明'],
        [
            ['明确小任务', '直接对话', '给清楚需求，AI 直接做'],
            ['复杂功能开发', 'Plan 模式', '先让 AI 规划，确认后再执行'],
            ['多任务并行', 'Worktree + Team', '隔离开发，互不干扰'],
            ['新领域探索', '调研先行', '先让 AI 帮你建认知'],
            ['代码评审', '只读模式', 'AI 读代码，人做决策'],
        ],
    )

    # ── P48: 新手快速行动建议 ──
    content_slide('快速行动建议 — 新手起步', [
        '##Step 1：从一个小项目或小模块开始',
        '##Step 2：用好 CLAUDE.md（3 行就够）',
        ('项目是什么 + 技术栈 + 核心约定', 1),
        '##Step 3：养成"设计 → AI 实现 → 验证"的习惯',
    ])

    # ── P49: 进阶者快速行动建议 ──
    content_slide('快速行动建议 — 进阶提升', [
        '##搭建项目级 Harness',
        ('CLAUDE.md + rules/ARCHITECTURE.md + docs/', 1),
        '##尝试 Worktree 并行开发',
        ('aip-portal 用了 24 个 worktree 并行', 1),
        '##建立复盘习惯：踩坑 → 补规范 → 定期精简',
    ])

    # ── P50: 观望者快速行动建议 ──
    card_slide('快速行动建议 — 观望者入门', [
        ('让 AI 写单元测试', '最小场景起步\n感受效率提升', ACCENT3),
        ('让 AI 做代码评审', '不需要写代码\n也能体验 AI 协作', ACCENT),
        ('让 AI 写 CRUD', '最常见的开发任务\n最容易上手', ACCENT2),
    ])

    # ── P51: 一页总结 ──
    content_slide('总结', [
        '>>AI 效率确定，不确定的是人的引导',
        '>>把最擅长的事交给 AI，释放时间做更有价值的事',
        '>>用 SDD + Harness 系统化地驾驭 AI',
        '>>从"写代码的人"变成"设计系统的人"',
    ])

    # ── P52: Q&A 结束页 ──
    emphasis_slide('感谢聆听', 'Q & A 交流时间', bc=ACCENT)


# ══════════════════════════════════════════════════════
# 主入口
# ══════════════════════════════════════════════════════

if __name__ == '__main__':
    create_slides()
    d = os.path.dirname(os.path.abspath(__file__))
    fp = os.path.join(d, 'AI协作编程实践v3.pptx')
    prs.save(fp)
    print(f'PPT saved: {fp}')
    print(f'Total slides: {len(prs.slides)}')
