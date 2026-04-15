"""
Part 1: Real Projects (P3-P14)
"""

from pptx.util import Cm, Pt
from pptx.enum.text import PP_ALIGN
from gen_helpers import *

def add_slides(prs):
    # ══════════════════════════════════════════════════════════
    # P3 - Part 1 Section Divider
    # ══════════════════════════════════════════════════════════
    slide = add_slide(prs)
    add_textbox(slide, Cm(3), Cm(5.0), SLIDE_W - Cm(6), Cm(1.2),
                "Part 01", 20, ACCENT_BLUE, align=PP_ALIGN.CENTER)
    add_textbox(slide, Cm(3), Cm(7.0), SLIDE_W - Cm(6), Cm(2.0),
                "\u4ece\u771f\u5b9e\u9879\u76ee\u8bf4\u8d77", 40, TEXT_WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_textbox(slide, Cm(3), Cm(10.0), SLIDE_W - Cm(6), Cm(1.2),
                "\u4e24\u4e2a\u6848\u4f8b\uff0c\u4e24\u79cd\u7ed3\u5c40", 18, TEXT_GREY, align=PP_ALIGN.CENTER)

    # ══════════════════════════════════════════════════════════
    # P4 - Case 1: aip-server Overview
    # ══════════════════════════════════════════════════════════
    slide = add_slide(prs)
    add_slide_title(slide, "\u6848\u4f8b\u4e00\uff1aaip-server \u9879\u76ee\u6982\u89c8")

    # Left side (40%) - project info
    left_w = Cm(12)
    add_textbox(slide, MARGIN_L, Cm(2.8), left_w, Cm(0.8),
                "aip-server", 20, ACCENT_BLUE, bold=True)
    add_multiline_textbox(slide, MARGIN_L, Cm(3.8), left_w, Cm(3.0),
        ["AI \u667a\u80fd\u5e73\u53f0\u540e\u7aef\u670d\u52a1",
         "\u4f01\u4e1a\u7ea7\u667a\u80fd\u4f53\u7ba1\u7406\u3001\u77e5\u8bc6\u5e93\u3001\u6570\u636e\u5206\u6790\u5e73\u53f0"],
        14, TEXT_WHITE)

    # Right side (55%) - three data cards
    right_x = Cm(15)
    right_w = SLIDE_W - Cm(15) - MARGIN_R
    card_h = Cm(3.2)

    # Card 1: Key numbers
    add_card(slide, right_x, Cm(2.8), right_w, card_h,
             "100 \u4e2a Controller | 11 \u4e2a\u4e1a\u52a1\u6a21\u5757 | 2394 \u6b21\u63d0\u4ea4",
             [], accent_color=ACCENT_BLUE, title_size=12)

    # Card 2: Tech stack
    add_card(slide, right_x, Cm(6.5), right_w, card_h,
             "\u6280\u672f\u6808",
             ["Spring Boot 2.7.18 + MyBatis Plus 3.4.2",
              "+ PostgreSQL + Java 17"],
             accent_color=ACCENT_BLUE, title_size=12, desc_size=11)

    # Card 3: Multi-datasource
    add_card(slide, right_x, Cm(10.2), right_w, card_h,
             "\u591a\u6570\u636e\u6e90",
             ["PostgreSQL + Neo4j + Dremio",
              "Feign \u5fae\u670d\u52a1\u8c03\u7528"],
             accent_color=ACCENT_BLUE, title_size=12, desc_size=11)

    add_page_number(slide, 4)

    # ══════════════════════════════════════════════════════════
    # P5 - Case 1: Logic/Tool Module Focus
    # ══════════════════════════════════════════════════════════
    slide = add_slide(prs)
    add_slide_title(slide, "\u805a\u7126\u6280\u80fd\u5de5\u5177\u6a21\u5757\uff08logic/tool\uff09")

    add_subtitle(slide, "5 \u4e2a Controller\uff0c\u7ea6 80 \u4e2a API \u7aef\u70b9\uff0c8 \u4e2a\u6838\u5fc3\u5b9e\u4f53", Cm(2.4), 13)

    # Five card columns
    controllers = [
        ("\u4ee3\u7801\u5e93\u7ba1\u7406", "LtCodeRepo\nController"),
        ("MCP \u670d\u52a1", "LtMcp\nController"),
        ("\u6269\u5c55\u63d2\u4ef6", "LtExtensionPlugin\nController"),
        ("OpenAPI", "LtOpenapi\nController"),
        ("\u5de5\u4f5c\u6d41", "LtTaskflow\nController"),
    ]
    card_w = Cm(5.5)
    gap = Cm(0.6)
    start_x = MARGIN_L
    for i, (title, cls) in enumerate(controllers):
        x = start_x + i * (card_w + gap)
        add_card(slide, x, Cm(4.5), card_w, Cm(5.0), title,
                 cls.split("\n"), title_size=13, desc_size=9,
                 accent_color=ACCENT_BLUE)

    add_page_number(slide, 5)

    # ══════════════════════════════════════════════════════════
    # P6 - Case 1: Collaboration Mode
    # ══════════════════════════════════════════════════════════
    slide = add_slide(prs)
    add_slide_title(slide, "\u534f\u4f5c\u6a21\u5f0f \u2014 \u4eba\u4e0e AI \u7684\u5206\u5de5")

    # Four flow nodes horizontally
    nodes = [
        ("\u4eba\u8bbe\u8ba1\u539f\u578b", "\U0001f9d1", ACCENT_BLUE),
        ("\u4eba\u8bbe\u8ba1\u8868\u7ed3\u6784", "\U0001f9d1", ACCENT_BLUE),
        ("AI \u7f16\u7801", "\U0001f916", GREEN),
        ("\u4eba\u8bc4\u5ba1 & \u6d4b\u8bd5", "\U0001f9d1", ACCENT_BLUE),
    ]
    node_w = Cm(6.0)
    node_h = Cm(2.2)
    arrow_w = Cm(1.2)
    start_x = Cm(1.5)
    y_flow = Cm(3.5)
    for i, (text, emoji, color) in enumerate(nodes):
        x = start_x + i * (node_w + arrow_w)
        border = color
        fill = CARD_BG if color != GREEN else RGBColor(0x0a, 0x2e, 0x14)
        add_flow_node(slide, x, y_flow, node_w, node_h, text, emoji,
                      fill_color=fill, border_color=border, text_size=12)
        if i < 3:
            add_arrow(slide, x + node_w + Cm(0.1), y_flow + Cm(0.7), Cm(0.9), Cm(0.7), TEXT_GREY)

    # Annotation under node 2
    add_textbox(slide, start_x + 1 * (node_w + arrow_w) + Cm(0.3), y_flow + node_h + Cm(0.2),
                node_w - Cm(0.6), Cm(0.6), "AI\u53c2\u4e0e\u8ba8\u8bba", 9, TEXT_GREY, align=PP_ALIGN.CENTER)

    # Three principles at bottom
    principles = [
        "\u2705 \u4eba\u8d1f\u8d23\u201c\u5b9a\u65b9\u5411\u201d\uff08\u8bbe\u8ba1\u3001\u8bc4\u5ba1\u3001\u9a8c\u6536\uff09",
        "\u2705 AI \u8d1f\u8d23\u201c\u8dd1\u91cf\u201d\uff08\u7f16\u7801\u5b9e\u73b0\uff09",
        "\u2705 \u53ca\u65f6\u8bc4\u5ba1\uff0c\u5c0f\u6b65\u9a8c\u8bc1",
    ]
    y = Cm(8.5)
    for p in principles:
        add_textbox(slide, MARGIN_L + Cm(2), y, SLIDE_W - MARGIN_L - MARGIN_R - Cm(4), Cm(0.7),
                    p, 13, TEXT_WHITE)
        y += Cm(1.2)

    add_page_number(slide, 6)

    # ══════════════════════════════════════════════════════════
    # P7 - Case 1: Key Results
    # ══════════════════════════════════════════════════════════
    slide = add_slide(prs)
    add_slide_title(slide, "\u5173\u952e\u6210\u679c")

    # Top: two big number cards
    num_w = Cm(14)
    add_rect(slide, MARGIN_L, Cm(2.8), num_w, Cm(2.5), CARD_BG, ACCENT_BLUE)
    add_textbox(slide, MARGIN_L, Cm(3.0), num_w, Cm(1.2),
                "5\u5929 \u5b8c\u6210 ~80\u4e2a API \u7aef\u70b9", 20, ACCENT_BLUE, bold=True, align=PP_ALIGN.CENTER)
    add_textbox(slide, MARGIN_L, Cm(4.2), num_w, Cm(0.6),
                "29\u6b21\u63d0\u4ea4\uff0c3\u67084-11\u65e5", 11, TEXT_GREY, align=PP_ALIGN.CENTER)

    add_rect(slide, MARGIN_L + num_w + Cm(1), Cm(2.8), num_w, Cm(2.5), CARD_BG, ACCENT_BLUE)
    add_textbox(slide, MARGIN_L + num_w + Cm(1), Cm(3.0), num_w, Cm(1.2),
                "\u6548\u7387\u63d0\u5347 2-4\u500d", 20, ACCENT_BLUE, bold=True, align=PP_ALIGN.CENTER)
    add_textbox(slide, MARGIN_L + num_w + Cm(1), Cm(4.2), num_w, Cm(0.6),
                "\u5bf9\u6bd4\u4f20\u7edf\u5f00\u53d1\u4f30\u8ba1\u503c", 11, TEXT_GREY, align=PP_ALIGN.CENTER)

    # Middle: fix commits
    add_textbox(slide, MARGIN_L, Cm(6.0), SLIDE_W - MARGIN_L - MARGIN_R, Cm(0.8),
                "13\u6b21 fix \u63d0\u4ea4\uff08\u536245%\uff09 \u2014 AI\u4ee3\u7801\u6709\u5c0f\u9519\u6f0f\uff0c\u4f46\u4eba\u5728\u8bc4\u5ba1\u4e2d\u53ca\u65f6\u6355\u83b7",
                12, TEXT_WHITE)

    # Bottom: success foundations (green)
    successes = [
        "\u2705 16 \u4e2a\u6570\u636e\u5e93\u8fc1\u79fb\u7248\u672c\uff0c\u8868\u7ed3\u6784\u6e05\u6670",
        "\u2705 8 \u4e2a\u6838\u5fc3\u5b9e\u4f53\uff0c\u8fb9\u754c\u660e\u786e",
        "\u2705 \u8bc4\u5ba1\u540e\u66f4\u65b0 CLAUDE.md\uff0c\u6301\u7eed\u6539\u8fdb",
    ]
    y = Cm(7.8)
    for s in successes:
        add_textbox(slide, MARGIN_L + Cm(1), y, SLIDE_W - MARGIN_L - MARGIN_R - Cm(2), Cm(0.7),
                    s, 13, GREEN)
        y += Cm(1.2)

    add_page_number(slide, 7)

    # ══════════════════════════════════════════════════════════
    # P8 - Case 1: Summary
    # ══════════════════════════════════════════════════════════
    slide = add_slide(prs)
    # Center big text
    add_textbox(slide, Cm(3), Cm(4.0), SLIDE_W - Cm(6), Cm(2.5),
                "\u8bbe\u8ba1\u5230\u4f4d + \u8fb9\u754c\u6e05\u6670 = AI \u6548\u7387\u62c9\u6ee1",
                32, ACCENT_BLUE, bold=True, align=PP_ALIGN.CENTER)

    # Three supports
    supports = [
        "\u2705 \u8868\u7ed3\u6784\u8bbe\u8ba1\u6e05\u6670 \u2192 AI \u6709\u660e\u786e\u8fb9\u754c",
        "\u2705 \u53ca\u65f6\u8bc4\u5ba1\uff0c\u5c0f\u6b65\u9a8c\u8bc1",
        "\u2705 \u4eba\u628a\u63a7\u65b9\u5411\uff0cAI \u5806\u780c\u5de5\u4f5c\u91cf",
    ]
    y = Cm(8.0)
    for s in supports:
        add_textbox(slide, Cm(4), y, SLIDE_W - Cm(8), Cm(0.8),
                    s, 16, GREEN, align=PP_ALIGN.CENTER)
        y += Cm(1.6)

    add_page_number(slide, 8)

    # ══════════════════════════════════════════════════════════
    # P9 - Case 2: A2A Adapter Background
    # ══════════════════════════════════════════════════════════
    slide = add_slide(prs)
    add_slide_title(slide, "\u6848\u4f8b\u4e8c\uff1aA2A \u9002\u914d\u5668\uff08\u4e00\u573a\u4ee3\u4ef7\u60e8\u75db\u7684\u8fd4\u5de5\uff09")

    # Left (50%): background info
    left_w = Cm(14.5)
    add_textbox(slide, MARGIN_L, Cm(3.0), left_w, Cm(0.8),
                "\u80cc\u666f\u4fe1\u606f", 16, TEXT_WHITE, bold=True)
    add_multiline_textbox(slide, MARGIN_L, Cm(4.0), left_w, Cm(4.0),
        ["\u2022 \u9700\u6c42\uff1a\u5f00\u53d1 A2A\uff08Agent-to-Agent\uff09\u534f\u8bae\u9002\u914d\u5668",
         "\u2022 \u6240\u5728\u9879\u76ee\uff1aaip-gateway \u7684 a2a-wrapper \u63d2\u4ef6",
         "\u2022 \u80cc\u666f\uff1a\u76f8\u5bf9\u964c\u751f\u7684\u9886\u57df\uff0c\u6ca1\u6709\u6df1\u5165\u8c03\u7814"],
        13, TEXT_WHITE)

    # Right (45%): red warning card
    right_x = Cm(17)
    right_w = SLIDE_W - right_x - MARGIN_R
    add_card(slide, right_x, Cm(3.0), right_w, Cm(7.0),
             "\u26a0\ufe0f \u5173\u952e\u95ee\u9898",
             ["\u534f\u4f5c\u65b9\u5f0f\uff1a\u8ba9 AI \u6839\u636e\u534f\u8bae\u6587\u6863\u76f4\u63a5\u7f16\u7801",
              "",
              "\u81ea\u5df1\u5bf9 A2A \u4e0d\u4e86\u89e3\uff0c",
              "\u65e0\u6cd5\u5224\u65ad AI \u8f93\u51fa\u662f\u5426\u6b63\u786e"],
             title_size=14, desc_size=12, title_color=RED_ORANGE,
             accent_color=RED_ORANGE, desc_color=TEXT_WHITE)

    add_page_number(slide, 9)

    # ══════════════════════════════════════════════════════════
    # P10 - Case 2: What Happened (Timeline)
    # ══════════════════════════════════════════════════════════
    slide = add_slide(prs)
    add_slide_title(slide, "\u53d1\u751f\u4e86\u4ec0\u4e48")

    # Vertical timeline (left 65%)
    timeline = [
        ("\U0001f7e2 \u7b2c1\u5929\uff1a\u8ba9 AI \u6839\u636e A2A \u534f\u8bae\u6587\u6863\u76f4\u63a5\u7f16\u7801", GREEN),
        ("\U0001f534 \u7b2c1-2\u5468\uff1aAI \u201c\u9ad8\u6548\u201d\u4ea7\u51fa\u5927\u91cf\u4ee3\u7801\uff0846\u6b21\u63d0\u4ea4\uff0c16\u5929\uff09", RED_ORANGE),
        ("\U0001f534 \u5173\u952e\u95ee\u9898\uff1a\u65e0\u6cd5\u5224\u65ad AI \u8f93\u51fa\u662f\u5426\u7b26\u5408\u534f\u8bae", RED_ORANGE),
        ("\U0001f534 \u6df1\u5165\u8c03\u7814\uff1a\u53d1\u73b0\u6709\u5b98\u65b9 Java SDK\uff08a2aproject/a2a-java\uff09", RED_ORANGE),
        ("\U0001f7e2 \u63a8\u5012\u91cd\u6765\uff1a\u7528 SDK 3\u5929\u641e\u5b9a", GREEN),
    ]
    y = Cm(3.0)
    for text, color in timeline:
        add_textbox(slide, MARGIN_L, y, Cm(20), Cm(0.8), text, 12, color)
        y += Cm(1.8)

    # Right side: big number comparison
    right_x = Cm(22)
    add_textbox(slide, right_x, Cm(4.0), Cm(9), Cm(2.0),
                "141,026\u884c \u2192 1,786\u884c", 28, RED_ORANGE, bold=True, align=PP_ALIGN.CENTER)
    add_textbox(slide, right_x, Cm(6.5), Cm(9), Cm(1.0),
                "\u4ee3\u7801\u91cf\u7f29\u51cf 98.7%", 18, RED_ORANGE, align=PP_ALIGN.CENTER)

    add_page_number(slide, 10)

    # ══════════════════════════════════════════════════════════
    # P11 - Case 2: Crash Scene 1 - Reinventing the Wheel
    # ══════════════════════════════════════════════════════════
    slide = add_slide(prs)
    add_slide_title(slide, "\u7ffb\u8f66\u73b0\u573a\uff08\u4e00\uff09\u2014 AI \u9020\u4e86\u4e2a\u8f6e\u5b50")

    # Left: self-built (red border)
    lw = Cm(14)
    add_card(slide, MARGIN_L, Cm(3.0), lw, Cm(5.5),
             "\u274c \u81ea\u7814\u65b9\u6848",
             ["\u4ece\u96f6\u5f00\u59cb\u7f16\u7801\uff0c\u5168\u7a0b\u81ea\u4fe1\u8f93\u51fa",
              "16\u5929\uff0c46\u6b21\u63d0\u4ea4\uff0c141,026\u884c\u4ee3\u7801",
              "\u770b\u4f3c\u5b8c\u6574\uff0c\u5b9e\u5219\u8d70\u5f2f\u8def"],
             title_size=14, desc_size=12, title_color=RED_ORANGE,
             accent_color=RED_ORANGE)

    # Right: official SDK (green border)
    rx = Cm(16.5)
    rw = SLIDE_W - rx - MARGIN_R
    add_card(slide, rx, Cm(3.0), rw, Cm(5.5),
             "\u2705 \u5b98\u65b9 SDK \u65b9\u6848",
             ["\u53d1\u73b0 a2aproject/a2a-java",
              "3\u5929\u5b8c\u6210\uff0c\u4ec5 1,786\u884c",
              "\u7a33\u5b9a\u6027\u3001\u517c\u5bb9\u6027\u6709\u4fdd\u969c"],
             title_size=14, desc_size=12, title_color=GREEN,
             accent_color=GREEN)

    # Bottom warning
    add_plain_rect(slide, MARGIN_L, Cm(9.5), SLIDE_W - MARGIN_L - MARGIN_R, Cm(2.0),
                   RGBColor(0x2d, 0x1a, 0x0a), YELLOW)
    add_textbox(slide, MARGIN_L + Cm(0.5), Cm(9.7), SLIDE_W - MARGIN_L - MARGIN_R - Cm(1), Cm(1.5),
                "\u26a0\ufe0f AI \u4e0d\u4f1a\u4e3b\u52a8\u544a\u8bc9\u4f60\u201c\u6709\u66f4\u597d\u7684\u65b9\u6848\u201d\uff0c\u65b9\u5411\u662f\u4f60\u5b9a\u7684\uff0cAI \u4e0d\u8d1f\u8d23\u7ea0\u504f",
                14, YELLOW, bold=True, align=PP_ALIGN.CENTER)

    add_page_number(slide, 11)

    # ══════════════════════════════════════════════════════════
    # P12 - Case 2: Crash Scene 2 - More Lessons
    # ══════════════════════════════════════════════════════════
    slide = add_slide(prs)
    add_slide_title(slide, "\u7ffb\u8f66\u73b0\u573a\uff08\u4e8c\uff09\u2014 \u4e09\u4e2a\u5c0f\u6848\u4f8b")

    # Three cards side by side
    card_w = Cm(9.5)
    gap = Cm(0.5)
    cases = [
        ("SQL \u7248\u672c\u53f7\u9519\u8bef",
         ["\u751f\u6210\u4e86V2.2.14\u800c\u975eV2.2.13",
          "\u6ca1\u6709\u5148\u67e5\u73b0\u6709\u7248\u672c\u6587\u4ef6",
          "\u6570\u636e\u5e93\u8fc1\u79fb\u6267\u884c\u5931\u8d25"]),
        ("\u4fee\u6539\u5df2\u51bb\u7ed3\u811a\u672c",
         ["\u4fee\u6539\u4e86\u5df2\u51bb\u7ed3\u7684V2.2.13__init.sql",
          "\u5df2\u90e8\u7f72\u73af\u5883\u6587\u4ef6\u6458\u8981\u4e0d\u4e00\u81f4",
          "\u590d\u76d8\uff1aFlyway\u5df2\u51bb\u7ed3\u811a\u672c\u4e25\u7981\u4fee\u6539"]),
        ("\u6743\u9650\u8fc7\u6ee4\u504f\u5dee",
         ["\u628a\u5e94\u4ea4\u7ed9\u524d\u7aef\u7684\u903b\u8f91\u4fdd\u7559\u5728\u540e\u7aef",
          "\u540e\u7aef\u591a\u4e86\u4e0d\u5fc5\u8981\u7684\u6743\u9650\u68c0\u67e5",
          "\u590d\u76d8\uff1a\u7528\u6237\u7ea7\u6743\u9650\u5224\u65ad\u4ea4\u7ed9\u524d\u7aef"]),
    ]
    for i, (title, lines) in enumerate(cases):
        x = MARGIN_L + i * (card_w + gap)
        add_card(slide, x, Cm(3.2), card_w, Cm(8.0), title, lines,
                 title_size=13, desc_size=11, accent_color=RED_ORANGE, desc_color=TEXT_WHITE)

    add_page_number(slide, 12)

    # ══════════════════════════════════════════════════════════
    # P13 - Case 2: Summary
    # ══════════════════════════════════════════════════════════
    slide = add_slide(prs)

    # Center big text (red-orange)
    add_textbox(slide, Cm(3), Cm(3.5), SLIDE_W - Cm(6), Cm(2.0),
                "\u65b9\u5411\u9519\u4e86\uff0c\u6548\u7387\u8d8a\u9ad8\u4ee3\u4ef7\u8d8a\u5927",
                32, RED_ORANGE, bold=True, align=PP_ALIGN.CENTER)

    # Root cause chain (4 steps with arrows)
    steps = [
        "\u5bf9\u9886\u57df\u4e0d\u4e86\u89e3",
        "\u6ca1\u6709\u5148\u5efa\u7acb\u8ba4\u77e5",
        "\u76f2\u76ee\u4fe1\u4efb AI",
        "AI \u9ad8\u6548\u8dd1\u504f",
    ]
    node_w = Cm(6.0)
    arrow_w = Cm(1.5)
    total_w = len(steps) * node_w + (len(steps) - 1) * arrow_w
    start_x = (SLIDE_W - total_w) // 2
    y_chain = Cm(8.0)
    for i, step in enumerate(steps):
        x = start_x + i * (node_w + arrow_w)
        colors = [YELLOW, YELLOW, RED_ORANGE, RED_ORANGE]
        add_flow_node(slide, x, y_chain, node_w, Cm(1.8), step, "",
                      fill_color=CARD_BG, border_color=colors[i],
                      text_color=colors[i], text_size=12)
        if i < 3:
            add_arrow(slide, x + node_w + Cm(0.15), y_chain + Cm(0.5), Cm(1.2), Cm(0.6), TEXT_GREY)

    add_page_number(slide, 13)

    # ══════════════════════════════════════════════════════════
    # P14 - Part 1 Wrap: Comparison
    # ══════════════════════════════════════════════════════════
    slide = add_slide(prs)
    add_slide_title(slide, "\u4e24\u4e2a\u6848\u4f8b\uff0c\u4e24\u79cd\u7ed3\u5c40")

    # Left: aip-server (green)
    lw = Cm(14)
    add_card(slide, MARGIN_L, Cm(3.0), lw, Cm(6.5),
             "\u2705 aip-server",
             ["\u8bbe\u8ba1\u6e05\u6670 + \u8fb9\u754c\u660e\u786e",
              "\u4eba\u5b9a\u65b9\u5411\uff0cAI \u8dd1\u91cf",
              "5\u5929\uff0c80\u4e2a\u63a5\u53e3\uff0c\u8d28\u91cf\u53ef\u9760"],
             title_size=16, desc_size=13, title_color=GREEN,
             accent_color=GREEN, desc_color=TEXT_WHITE)

    # Right: A2A (red)
    rx = Cm(16.5)
    rw = SLIDE_W - rx - MARGIN_R
    add_card(slide, rx, Cm(3.0), rw, Cm(6.5),
             "\u274c A2A \u9002\u914d\u5668",
             ["\u65b9\u5411\u4e0d\u660e + \u76f2\u76ee\u4fe1\u4efb",
              "\u6ca1\u6709\u8c03\u7814\u5c31\u52a8\u624b",
              "16\u5929\u8fd4\u5de5\uff0c98.7%\u4ee3\u7801\u5e9f\u5f03"],
             title_size=16, desc_size=13, title_color=RED_ORANGE,
             accent_color=RED_ORANGE, desc_color=TEXT_WHITE)

    # Bottom question (blue)
    add_textbox(slide, MARGIN_L, Cm(11.5), SLIDE_W - MARGIN_L - MARGIN_R, Cm(1.5),
                "\u4ec0\u4e48\u65f6\u5019\u8be5\u4fe1\u4efb AI\uff1f\u5982\u4f55\u7cfb\u7edf\u5730\u8ba9 AI \u7f16\u7a0b\u8d8a\u6765\u8d8a\u9ad8\u6548\uff1f",
                18, ACCENT_BLUE, bold=True, align=PP_ALIGN.CENTER)

    add_page_number(slide, 14)
