import sys, io
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')

from pptx import Presentation

prs = Presentation('D:/GitCode/self-dev/ai-coding-share/output/v5.2/slides/output/ai_coding_practice.pptx')

DARK_COLORS = {'000000', '0D1117', '161B22', '30363D'}
ISSUE_COUNT = 0

for i, slide in enumerate(prs.slides):
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                text_preview = run.text[:60].strip()
                if not text_preview:
                    continue
                font = run.font
                color_str = None
                try:
                    if font.color.rgb is not None:
                        color_str = str(font.color.rgb)
                except:
                    pass

                if color_str and color_str in DARK_COLORS:
                    ISSUE_COUNT += 1
                    # Avoid emoji encoding issues
                    safe_text = text_preview.encode('ascii', 'replace').decode('ascii')
                    print(f'  Slide {i+1}: color=#{color_str} text="{safe_text}"')

print(f'Total slides: {len(prs.slides)}')
print(f'Dark text issues found: {ISSUE_COUNT}')