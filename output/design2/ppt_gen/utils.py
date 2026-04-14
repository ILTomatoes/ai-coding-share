"""PPT 生成工具函数库"""
from pptx import Presentation
from pptx.util import Cm, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ==================== 颜色常量 ====================
BG_COLOR = RGBColor(0x0D, 0x11, 0x17)
ACCENT_BLUE = RGBColor(0x2F, 0x81, 0xF7)
ACCENT_GREEN = RGBColor(0x3F, 0xB9, 0x50)
ACCENT_RED = RGBColor(0xF8, 0x51, 0x49)
ACCENT_YELLOW = RGBColor(0xE3, 0xB3, 0x41)
TEXT_LIGHT = RGBColor(0xE6, 0xED, 0xF3)
TEXT_GRAY = RGBColor(0x8B, 0x94, 0x9E)
CARD_BG = RGBColor(0x16, 0x1B, 0x22)
CARD_BORDER = RGBColor(0x30, 0x36, 0x3D)
CODE_GREEN = RGBColor(0x7E, 0xE7, 0x87)
CODE_BLUE = RGBColor(0x79, 0xC0, 0xFF)
DARKER_BG = RGBColor(0x08, 0x0C, 0x12)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

# ==================== 字体常量 ====================
FONT_CN = '微软雅黑'
FONT_EN = 'Calibri'
FONT_CODE = 'Courier New'

# ==================== 尺寸常量 (cm) ====================
SLIDE_W = 33.87
SLIDE_H = 19.05
MARGIN_L = 1.8
MARGIN_R = 1.8
MARGIN_T = 1.5
MARGIN_B = 1.2
CONTENT_W = SLIDE_W - MARGIN_L - MARGIN_R  # 30.27
CONTENT_H = SLIDE_H - MARGIN_T - MARGIN_B  # 16.35
CONTENT_RIGHT = SLIDE_W - MARGIN_R          # 32.07
CONTENT_BOTTOM = SLIDE_H - MARGIN_B         # 17.85

# 无页码的页面
NO_PAGE_NUM = {1, 3, 15, 22, 42}

# 对齐映射
ALIGN_MAP = {'left': PP_ALIGN.LEFT, 'center': PP_ALIGN.CENTER, 'right': PP_ALIGN.RIGHT}

# ==================== 基础函数 ====================

def init_presentation():
    """创建并配置 Presentation"""
    prs = Presentation()
    prs.slide_width = Cm(SLIDE_W)
    prs.slide_height = Cm(SLIDE_H)
    return prs

def new_slide(prs):
    """创建深色背景的新幻灯片"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = BG_COLOR
    return slide

def add_text(slide, left, top, width, height, text, font_size=14,
             color=TEXT_LIGHT, bold=False, alignment='left', font_name=None):
    """添加文本框"""
    txBox = slide.shapes.add_textbox(Cm(left), Cm(top), Cm(width), Cm(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = str(text)
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name or FONT_CN
    p.alignment = ALIGN_MAP.get(alignment, PP_ALIGN.LEFT)
    return txBox

def add_multiline(slide, left, top, width, height, lines, font_size=14,
                  color=TEXT_LIGHT, bold=False, alignment='left', font_name=None,
                  line_spacing=1.2, bullet_char=None):
    """添加多行文本"""
    txBox = slide.shapes.add_textbox(Cm(left), Cm(top), Cm(width), Cm(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        txt = f"{bullet_char} {line}" if bullet_char else line
        p.text = txt
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.bold = bold
        p.font.name = font_name or FONT_CN
        p.alignment = ALIGN_MAP.get(alignment, PP_ALIGN.LEFT)
        p.space_after = Pt(font_size * (line_spacing - 1))
    return txBox

def add_page_title(slide, text, top=1.0, font_size=26, color=TEXT_LIGHT):
    """添加页面标题"""
    return add_text(slide, MARGIN_L, top, CONTENT_W, 1.5, text,
                    font_size=font_size, color=color, bold=True)

def add_page_num(slide, num):
    """添加页码（跳过封面和章节分隔页）"""
    if num in NO_PAGE_NUM:
        return
    add_text(slide, 30.5, 17.8, 2.5, 0.8, f"P{num}",
             font_size=10, color=TEXT_GRAY, alignment='right')

# ==================== 形状函数 ====================

def add_rect(slide, left, top, width, height, fill_color=CARD_BG,
             border_color=CARD_BORDER, border_width=1):
    """添加矩形"""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Cm(left), Cm(top), Cm(width), Cm(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(border_width)
    else:
        shape.line.fill.background()
    return shape

def add_rounded_rect(slide, left, top, width, height, fill_color=CARD_BG,
                     border_color=CARD_BORDER, border_width=1):
    """添加圆角矩形"""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Cm(left), Cm(top), Cm(width), Cm(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(border_width)
    else:
        shape.line.fill.background()
    return shape

def add_card(slide, left, top, width, height, title, content_lines=None,
             title_color=ACCENT_BLUE, content_color=TEXT_LIGHT,
             card_bg=CARD_BG, border_color=CARD_BORDER, border_top_color=None,
             title_size=14, content_size=12):
    """添加带标题和内容的卡片"""
    bc = border_top_color or border_color
    bw = 2 if border_top_color else 1
    card = add_rounded_rect(slide, left, top, width, height, card_bg, bc, bw)
    add_text(slide, left + 0.3, top + 0.2, width - 0.6, 0.8, title,
             font_size=title_size, color=title_color, bold=True)
    if content_lines:
        add_multiline(slide, left + 0.3, top + 1.0, width - 0.6, height - 1.2,
                      content_lines, font_size=content_size, color=content_color)
    return card

def add_big_number(slide, left, top, number, label, num_color=ACCENT_BLUE,
                   num_size=36, label_size=12, label_color=TEXT_GRAY, width=8):
    """添加大数字 + 标签"""
    add_text(slide, left, top, width, 1.5, str(number),
             font_size=num_size, color=num_color, bold=True, alignment='center')
    add_text(slide, left, top + 1.5, width, 0.8, label,
             font_size=label_size, color=label_color, alignment='center')

def add_arrow_right(slide, left, top, width=1.0, height=0.6, color=ACCENT_BLUE):
    """添加右箭头"""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RIGHT_ARROW, Cm(left), Cm(top), Cm(width), Cm(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_arrow_down(slide, left, top, width=0.6, height=0.8, color=ACCENT_BLUE):
    """添加下箭头"""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.DOWN_ARROW, Cm(left), Cm(top), Cm(width), Cm(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_flow_node(slide, left, top, width, height, text, fill_color=ACCENT_BLUE,
                  text_color=None, font_size=12, bold=True):
    """添加流程节点"""
    text_color = text_color or WHITE
    shape = add_rounded_rect(slide, left, top, width, height,
                             fill_color=fill_color, border_color=None)
    txBox = add_text(slide, left, top + 0.1, width, height - 0.2, text,
                     font_size=font_size, color=text_color, bold=bold, alignment='center')
    return shape, txBox

# ==================== 特殊页面函数 ====================

def add_cover_page(prs):
    """创建封面页 (P1)"""
    slide = new_slide(prs)
    add_text(slide, 3, 5.5, 28, 3, 'Java 开发中的 AI 协作编程实践',
             font_size=44, color=ACCENT_BLUE, bold=True, alignment='center')
    add_text(slide, 3, 8.5, 28, 1.5, '从真实项目说起',
             font_size=20, color=TEXT_GRAY, alignment='center')
    add_text(slide, 3, 14, 28, 1, '演讲人 / 2026年4月',
             font_size=14, color=TEXT_GRAY, alignment='center')
    add_text(slide, 20, 17, 12, 0.8, '内部交流分享',
             font_size=10, color=TEXT_GRAY, alignment='right')
    return slide

def add_section_divider(prs, part_label, title, subtitle=''):
    """创建章节分隔页"""
    slide = new_slide(prs)
    add_text(slide, 5, 6, 24, 1.5, part_label,
             font_size=18, color=ACCENT_BLUE, alignment='center')
    add_text(slide, 5, 7.5, 24, 2.5, title,
             font_size=36, color=TEXT_LIGHT, bold=True, alignment='center')
    if subtitle:
        add_text(slide, 5, 10.5, 24, 1.5, subtitle,
                 font_size=18, color=TEXT_GRAY, alignment='center')
    return slide

def add_table(slide, left, top, width, headers, rows, col_widths=None, row_h=1.2):
    """添加样式表格"""
    num_rows = len(rows) + 1
    num_cols = len(headers)
    height = row_h * num_rows
    table_shape = slide.shapes.add_table(
        num_rows, num_cols, Cm(left), Cm(top), Cm(width), Cm(height))
    table = table_shape.table
    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = Cm(w)
    # 表头
    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = ACCENT_BLUE
        for p in cell.text_frame.paragraphs:
            p.font.color.rgb = WHITE
            p.font.size = Pt(12)
            p.font.bold = True
            p.font.name = FONT_CN
    # 数据行
    for r, row_data in enumerate(rows):
        for c, val in enumerate(row_data):
            cell = table.cell(r + 1, c)
            cell.text = str(val)
            bg = CARD_BG if r % 2 == 0 else DARKER_BG
            cell.fill.solid()
            cell.fill.fore_color.rgb = bg
            for p in cell.text_frame.paragraphs:
                p.font.color.rgb = TEXT_LIGHT
                p.font.size = Pt(11)
                p.font.name = FONT_CN
    return table_shape

def add_code_block(slide, left, top, width, height, code_text, font_size=10):
    """添加代码块"""
    add_rounded_rect(slide, left, top, width, height, DARKER_BG, CARD_BORDER)
    lines = code_text.split('\n') if isinstance(code_text, str) else code_text
    add_multiline(slide, left + 0.3, top + 0.2, width - 0.6, height - 0.4,
                  lines, font_size=font_size, color=CODE_GREEN,
                  font_name=FONT_CODE, line_spacing=1.0)
