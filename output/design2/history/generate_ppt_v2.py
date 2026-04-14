#!/usr/bin/env python3
"""生成《Java开发中的 AI 协作编程实战》PPT v2.0 — 样式优化版"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
import os

# ── 配色方案 ──
BG_DARK = RGBColor(0x12, 0x15, 0x2A)        # 深蓝黑背景（更纯净）
BG_MID = RGBColor(0x1C, 0x25, 0x41)          # 中深背景（与 BG_DARK 拉开差距）
ACCENT = RGBColor(0x00, 0xD2, 0xFF)          # 青蓝强调色
ACCENT2 = RGBColor(0xFF, 0x6B, 0x6B)         # 红色强调
ACCENT3 = RGBColor(0x51, 0xCF, 0x66)         # 绿色强调
ACCENT4 = RGBColor(0xFF, 0xD4, 0x3B)         # 黄色强调
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT = RGBColor(0xBB, 0xBB, 0xCC)           # 浅灰文字（与白色拉开对比）
DARK_LIGHT = RGBColor(0x88, 0x88, 0x99)      # 更暗的灰色
CODE_BG = RGBColor(0x0A, 0x0E, 0x1A)         # 代码块背景（更深）
TABLE_HEADER_BG = RGBColor(0x00, 0x6B, 0x8A)
TABLE_ROW_EVEN = RGBColor(0x14, 0x1C, 0x34)
TABLE_ROW_ODD = RGBColor(0x10, 0x16, 0x2C)
CARD_BG = RGBColor(0x1A, 0x24, 0x42)         # 卡片背景
HEADER_BG = RGBColor(0x0E, 0x12, 0x22)       # 标题栏背景
SIDE_ACCENT = RGBColor(0x00, 0x6B, 0x8A)     # 侧边装饰色
DECO_BG = RGBColor(0x1E, 0x2D, 0x50)         # 装饰元素背景

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# 项目根目录
BASE_DIR = os.path.dirname(os.path.dirname(os.path.dirname(os.path.abspath(__file__))))

# 页码计数器
_page_counter = [0]


# ════════════════════════════════════════════════════════════
# 基础辅助函数
# ════════════════════════════════════════════════════════════

def set_slide_bg(slide, color=BG_DARK):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, left, top, width, height):
    return slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))


def add_para(tf, text, size=18, color=WHITE, bold=False, alignment=PP_ALIGN.LEFT,
             space_after=Pt(8), font_name='微软雅黑', line_spacing=None):
    p = tf.paragraphs[-1] if len(tf.paragraphs) == 1 and tf.paragraphs[0].text == '' else tf.add_paragraph()
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    p.space_after = space_after
    if line_spacing is not None:
        p.line_spacing = line_spacing
    return p


def add_shape(slide, shape_type, left, top, width, height, fill_color=None, line_color=None, line_width=None):
    """添加形状的通用函数"""
    shape = slide.shapes.add_shape(shape_type, Inches(left), Inches(top), Inches(width), Inches(height))
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        if line_width:
            shape.line.width = Pt(line_width)
    else:
        shape.line.fill.background()
    return shape


def add_side_decoration(slide):
    """右上角 + 右下角几何装饰"""
    # 右上角小方块
    add_shape(slide, MSO_SHAPE.RECTANGLE, 12.5, 0.2, 0.4, 0.4, fill_color=DECO_BG)
    add_shape(slide, MSO_SHAPE.RECTANGLE, 12.9, 0.5, 0.2, 0.2, fill_color=SIDE_ACCENT)


def add_footer(slide):
    """添加页码和底部装饰线"""
    _page_counter[0] += 1
    page_num = _page_counter[0]
    # 底部细线
    add_shape(slide, MSO_SHAPE.RECTANGLE, 0.8, 7.0, 11.7, 0.01, fill_color=RGBColor(0x33, 0x33, 0x44))
    # 页码
    tb = add_textbox(slide, 5.5, 7.05, 2.3, 0.35)
    tf = tb.text_frame
    add_para(tf, f'{page_num}', size=10, color=DARK_LIGHT, alignment=PP_ALIGN.CENTER)


def add_header_bar(slide, title):
    """统一的标题栏：深色背景条 + 底部强调线"""
    # 标题栏背景
    add_shape(slide, MSO_SHAPE.RECTANGLE, 0, 0, 13.333, 1.25, fill_color=HEADER_BG)
    # 强调色底线
    add_shape(slide, MSO_SHAPE.RECTANGLE, 0, 1.25, 13.333, 0.04, fill_color=ACCENT)
    # 标题文字
    tb = add_textbox(slide, 0.8, 0.25, 11.7, 0.85)
    tf = tb.text_frame
    tf.word_wrap = True
    add_para(tf, title, size=26, bold=True, color=WHITE)


# ════════════════════════════════════════════════════════════
# Slide 生成函数
# ════════════════════════════════════════════════════════════

def make_title_slide(title, subtitle=''):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    # 左侧大色条
    add_shape(slide, MSO_SHAPE.RECTANGLE, 0, 0, 0.12, 7.5, fill_color=ACCENT)
    # 左侧窄装饰条
    add_shape(slide, MSO_SHAPE.RECTANGLE, 0.22, 2.0, 0.04, 3.5, fill_color=SIDE_ACCENT)

    # 右上角装饰
    add_shape(slide, MSO_SHAPE.RECTANGLE, 11.5, 0.3, 1.2, 0.06, fill_color=SIDE_ACCENT)
    add_shape(slide, MSO_SHAPE.RECTANGLE, 12.2, 0.3, 0.06, 0.8, fill_color=SIDE_ACCENT)

    # 右下角装饰
    add_shape(slide, MSO_SHAPE.RECTANGLE, 11.5, 7.1, 1.2, 0.06, fill_color=DECO_BG)
    add_shape(slide, MSO_SHAPE.RECTANGLE, 12.7, 6.5, 0.06, 0.66, fill_color=DECO_BG)

    # 标题
    tb = add_textbox(slide, 0.8, 2.3, 11.3, 1.5)
    tf = tb.text_frame
    add_para(tf, title, size=42, bold=True, color=WHITE)

    # 强调线
    add_shape(slide, MSO_SHAPE.RECTANGLE, 0.8, 4.0, 4.0, 0.05, fill_color=ACCENT)

    if subtitle:
        tb2 = add_textbox(slide, 0.8, 4.3, 11.3, 1.0)
        tf2 = tb2.text_frame
        tf2.word_wrap = True
        add_para(tf2, subtitle, size=20, color=LIGHT)

    add_footer(slide)
    return slide


def make_section_slide(part_num, title, subtitle=''):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_MID)
    add_side_decoration(slide)

    if part_num:
        # 大号半透明背景编号
        tb_bg = add_textbox(slide, 0.5, 0.8, 5.0, 2.5)
        tf_bg = tb_bg.text_frame
        add_para(tf_bg, f'第{part_num}部分', size=14, color=ACCENT, bold=True)

        # 巨大的背景装饰数字
        add_shape(slide, MSO_SHAPE.RECTANGLE, 0.5, 0.6, 0.08, 2.8, fill_color=ACCENT)

    # 标题
    tb2 = add_textbox(slide, 1.2, 2.5, 10.5, 1.5)
    tf2 = tb2.text_frame
    tf2.word_wrap = True
    add_para(tf2, title, size=36, bold=True, color=WHITE)

    # 强调线
    add_shape(slide, MSO_SHAPE.RECTANGLE, 1.2, 4.1, 5.0, 0.05, fill_color=ACCENT)

    if subtitle:
        tb3 = add_textbox(slide, 1.2, 4.4, 10.5, 0.8)
        tf3 = tb3.text_frame
        add_para(tf3, subtitle, size=18, color=LIGHT)

    add_footer(slide)
    return slide


def make_content_slide(title, bullets):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_header_bar(slide, title)
    add_side_decoration(slide)

    tb2 = add_textbox(slide, 0.8, 1.6, 11.7, 5.2)
    tf2 = tb2.text_frame
    tf2.word_wrap = True

    for bullet in bullets:
        if isinstance(bullet, tuple):
            text, level, color = bullet
            prefix = '    ' * level + ('• ' if level == 0 else '◦ ' if level == 1 else '- ')
            add_para(tf2, prefix + text, size=17 if level == 0 else 15, color=color or WHITE,
                     space_after=Pt(6))
        elif isinstance(bullet, str):
            if bullet == '':
                # 空行作为分隔
                add_para(tf2, '', size=10, space_after=Pt(4))
            elif bullet.startswith('##'):
                # 小标题 - 前面加色点
                add_para(tf2, '  ◆  ' + bullet[2:].strip(), size=19, bold=True,
                         color=ACCENT, space_after=Pt(6))
            elif bullet.startswith('>>'):
                add_para(tf2, '  ▶  ' + bullet[2:].strip(), size=17, bold=True,
                         color=ACCENT4, space_after=Pt(6))
            elif bullet.startswith('!!'):
                add_para(tf2, '  ⚠  ' + bullet[2:].strip(), size=17, bold=True,
                         color=ACCENT2, space_after=Pt(6))
            else:
                add_para(tf2, '• ' + bullet, size=17, color=WHITE, space_after=Pt(6))

    add_footer(slide)
    return slide


def make_two_col_slide(title, left_title, left_items, right_title, right_items):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_header_bar(slide, title)
    add_side_decoration(slide)

    col_width = 5.6
    col_top = 1.6

    # 左栏卡片背景
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 0.6, col_top, col_width + 0.2, 5.2,
              fill_color=CARD_BG, line_color=RGBColor(0x2A, 0x3A, 0x5A), line_width=1)
    # 左栏标题色条
    add_shape(slide, MSO_SHAPE.RECTANGLE, 0.6, col_top, col_width + 0.2, 0.06, fill_color=ACCENT)

    # 左栏标题
    tb_lt = add_textbox(slide, 0.8, col_top + 0.2, col_width, 0.45)
    tf_lt = tb_lt.text_frame
    add_para(tf_lt, left_title, size=19, bold=True, color=ACCENT)
    # 左栏内容
    tb_l = add_textbox(slide, 0.8, col_top + 0.7, col_width, 4.3)
    tf_l = tb_l.text_frame
    tf_l.word_wrap = True
    for item in left_items:
        add_para(tf_l, '• ' + item, size=15, color=WHITE, space_after=Pt(5))

    # 分隔线
    add_shape(slide, MSO_SHAPE.RECTANGLE, 6.55, col_top + 0.3, 0.02, 4.6,
              fill_color=RGBColor(0x33, 0x44, 0x66))

    # 右栏卡片背景
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 6.8, col_top, col_width + 0.2, 5.2,
              fill_color=CARD_BG, line_color=RGBColor(0x2A, 0x3A, 0x5A), line_width=1)
    # 右栏标题色条
    add_shape(slide, MSO_SHAPE.RECTANGLE, 6.8, col_top, col_width + 0.2, 0.06, fill_color=ACCENT3)

    # 右栏标题
    tb_rt = add_textbox(slide, 7.0, col_top + 0.2, col_width, 0.45)
    tf_rt = tb_rt.text_frame
    add_para(tf_rt, right_title, size=19, bold=True, color=ACCENT3)
    # 右栏内容
    tb_r = add_textbox(slide, 7.0, col_top + 0.7, col_width, 4.3)
    tf_r = tb_r.text_frame
    tf_r.word_wrap = True
    for item in right_items:
        add_para(tf_r, '• ' + item, size=15, color=WHITE, space_after=Pt(5))

    add_footer(slide)
    return slide


def make_quote_slide(quote, source=''):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_MID)
    add_side_decoration(slide)

    # 大号装饰引号
    tb_deco = add_textbox(slide, 1.0, 1.0, 3.0, 2.0)
    tf_deco = tb_deco.text_frame
    add_para(tf_deco, '\u201C', size=120, bold=True, color=RGBColor(0x1E, 0x2D, 0x50),
             alignment=PP_ALIGN.LEFT)

    # 引号内容
    tb = add_textbox(slide, 1.8, 2.5, 9.7, 2.5)
    tf = tb.text_frame
    tf.word_wrap = True
    add_para(tf, quote, size=28, bold=True, color=ACCENT, alignment=PP_ALIGN.CENTER,
             space_after=Pt(10))

    if source:
        tb2 = add_textbox(slide, 1.8, 5.2, 9.7, 0.5)
        tf2 = tb2.text_frame
        add_para(tf2, f'—— {source}', size=16, color=LIGHT, alignment=PP_ALIGN.CENTER)

    # 底部装饰线
    add_shape(slide, MSO_SHAPE.RECTANGLE, 5.5, 6.2, 2.3, 0.04, fill_color=ACCENT)

    add_footer(slide)
    return slide


def make_code_slide(title, code_text):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_header_bar(slide, title)
    add_side_decoration(slide)

    # 代码区 - 圆角矩形
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.6), Inches(1.55), Inches(12.1), Inches(5.2)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = CODE_BG
    shape.line.color.rgb = RGBColor(0x22, 0x33, 0x55)
    shape.line.width = Pt(1)

    # 代码顶部标签条
    add_shape(slide, MSO_SHAPE.RECTANGLE, 0.6, 1.55, 12.1, 0.35, fill_color=RGBColor(0x15, 0x1A, 0x2A))
    tb_label = add_textbox(slide, 0.8, 1.55, 2.0, 0.35)
    tf_label = tb_label.text_frame
    add_para(tf_label, 'DIAGRAM', size=10, bold=True, color=DARK_LIGHT, font_name='Consolas')

    # 代码内容
    tf_code = shape.text_frame
    tf_code.word_wrap = True
    tf_code.margin_left = Inches(0.4)
    tf_code.margin_right = Inches(0.3)
    tf_code.margin_top = Inches(0.4)

    for i, line in enumerate(code_text.split('\n')):
        p = tf_code.paragraphs[0] if i == 0 else tf_code.add_paragraph()
        p.text = line
        p.font.size = Pt(12)
        p.font.color.rgb = ACCENT
        p.font.name = 'Consolas'
        p.space_after = Pt(0)
        p.space_before = Pt(0)
        p.line_spacing = Pt(14)

    add_footer(slide)
    return slide


def make_table_slide(title, headers, rows):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_header_bar(slide, title)
    add_side_decoration(slide)

    num_cols = len(headers)
    num_rows = len(rows) + 1
    row_h = 0.55
    table_shape = slide.shapes.add_table(
        num_rows, num_cols,
        Inches(0.6), Inches(1.6),
        Inches(12.1), Inches(row_h * num_rows)
    )
    table = table_shape.table

    # 表头
    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = h
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(15)
            p.font.bold = True
            p.font.color.rgb = WHITE
            p.font.name = '微软雅黑'
            p.alignment = PP_ALIGN.LEFT
        cell.fill.solid()
        cell.fill.fore_color.rgb = TABLE_HEADER_BG
        cell.margin_left = Inches(0.15)
        cell.margin_right = Inches(0.1)
        cell.margin_top = Inches(0.08)
        cell.margin_bottom = Inches(0.08)

    # 数据行
    for r_idx, row in enumerate(rows):
        for c_idx, val in enumerate(row):
            cell = table.cell(r_idx + 1, c_idx)
            cell.text = val
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(14)
                p.font.color.rgb = WHITE
                p.font.name = '微软雅黑'
                # 首列加粗变色
                if c_idx == 0:
                    p.font.bold = True
                    p.font.color.rgb = ACCENT
            cell.fill.solid()
            cell.fill.fore_color.rgb = TABLE_ROW_EVEN if r_idx % 2 == 0 else TABLE_ROW_ODD
            cell.margin_left = Inches(0.15)
            cell.margin_right = Inches(0.1)
            cell.margin_top = Inches(0.08)
            cell.margin_bottom = Inches(0.08)

    add_footer(slide)
    return slide


def make_timeline_slide(title, phases):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_header_bar(slide, title)
    add_side_decoration(slide)

    num = len(phases)
    spacing = 11.3 / (num - 1) if num > 1 else 0
    line_y = 3.3

    # 时间线横线（更粗）
    add_shape(slide, MSO_SHAPE.RECTANGLE, 0.8, line_y, 11.7, 0.06, fill_color=ACCENT)

    for i, (phase_title, phase_desc, phase_detail) in enumerate(phases):
        x = 1.0 + i * spacing

        # 节点圆（更大）
        dot_size = 0.36
        add_shape(slide, MSO_SHAPE.OVAL, x - dot_size / 2, line_y - dot_size / 2 + 0.03,
                  dot_size, dot_size, fill_color=ACCENT)
        # 内部小圆
        add_shape(slide, MSO_SHAPE.OVAL, x - 0.08, line_y - 0.08 + 0.03,
                  0.16, 0.16, fill_color=BG_DARK)

        # 上方标签（阶段名 + 简述）
        tb_up = add_textbox(slide, x - 1.0, 1.6, 2.2, 1.2)
        tf_up = tb_up.text_frame
        tf_up.word_wrap = True
        add_para(tf_up, phase_title, size=15, bold=True, color=ACCENT, alignment=PP_ALIGN.CENTER)
        for line in phase_desc.split('\n'):
            add_para(tf_up, line, size=12, color=LIGHT, alignment=PP_ALIGN.CENTER, space_after=Pt(2))

        # 下方详情卡片
        card_w = 2.4
        card_h = 2.5
        card_x = x - card_w / 2
        card_y = line_y + 0.5
        add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, card_x, card_y, card_w, card_h,
                  fill_color=CARD_BG, line_color=RGBColor(0x2A, 0x3A, 0x5A), line_width=1)
        # 卡片顶部小色条
        add_shape(slide, MSO_SHAPE.RECTANGLE, card_x + 0.2, card_y, card_w - 0.4, 0.04,
                  fill_color=ACCENT)

        tb_dn = add_textbox(slide, card_x + 0.15, card_y + 0.15, card_w - 0.3, card_h - 0.3)
        tf_dn = tb_dn.text_frame
        tf_dn.word_wrap = True
        for line in phase_detail.split('\n'):
            add_para(tf_dn, line, size=13, color=LIGHT, alignment=PP_ALIGN.CENTER, space_after=Pt(2))

    add_footer(slide)
    return slide


def make_number_slide(title, numbers):
    """数据冲击页：大数字卡片展示"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_header_bar(slide, title)
    add_side_decoration(slide)

    num_items = len(numbers)
    if num_items == 0:
        add_footer(slide)
        return slide

    cols = min(num_items, 3)
    card_w = 3.5
    card_h = 2.0
    gap = 0.5
    total_w = cols * card_w + (cols - 1) * gap
    start_x = (13.333 - total_w) / 2
    start_y = 2.2

    for idx, (num_val, num_desc) in enumerate(numbers):
        row = idx // cols
        col = idx % cols
        x = start_x + col * (card_w + gap)
        y = start_y + row * (card_h + 0.5)

        # 卡片
        card = add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, card_w, card_h,
                         fill_color=CARD_BG, line_color=RGBColor(0x2A, 0x3A, 0x5A), line_width=1)
        # 顶部色条
        add_shape(slide, MSO_SHAPE.RECTANGLE, x + 0.1, y, card_w - 0.2, 0.05, fill_color=ACCENT)

        # 数字
        tb_num = add_textbox(slide, x + 0.1, y + 0.2, card_w - 0.2, 0.7)
        tf_num = tb_num.text_frame
        add_para(tf_num, num_val, size=36, bold=True, color=ACCENT, alignment=PP_ALIGN.CENTER)

        # 描述
        tb_desc = add_textbox(slide, x + 0.15, y + 1.0, card_w - 0.3, 0.8)
        tf_desc = tb_desc.text_frame
        tf_desc.word_wrap = True
        for line in num_desc.split('\n'):
            add_para(tf_desc, line, size=13, color=LIGHT, alignment=PP_ALIGN.CENTER, space_after=Pt(2))

    add_footer(slide)
    return slide


def make_card_slide(title, cards):
    """卡片式布局：多张卡片并排展示"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_header_bar(slide, title)
    add_side_decoration(slide)

    num_cards = len(cards)
    if num_cards == 0:
        add_footer(slide)
        return slide

    gap = 0.35
    card_w = (11.7 - (num_cards - 1) * gap) / num_cards
    card_h = 4.8
    start_x = 0.8
    start_y = 1.6

    for idx, (card_title, card_content, card_color) in enumerate(cards):
        x = start_x + idx * (card_w + gap)

        # 卡片背景
        add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, start_y, card_w, card_h,
                  fill_color=CARD_BG, line_color=card_color, line_width=2)
        # 顶部色条
        add_shape(slide, MSO_SHAPE.RECTANGLE, x + 0.05, start_y, card_w - 0.1, 0.06,
                  fill_color=card_color)

        # 卡片标题
        tb_ct = add_textbox(slide, x + 0.2, start_y + 0.2, card_w - 0.4, 0.5)
        tf_ct = tb_ct.text_frame
        tf_ct.word_wrap = True
        add_para(tf_ct, card_title, size=17, bold=True, color=card_color)

        # 标题下分隔线
        add_shape(slide, MSO_SHAPE.RECTANGLE, x + 0.2, start_y + 0.7, card_w - 0.4, 0.02,
                  fill_color=RGBColor(0x33, 0x44, 0x66))

        # 卡片内容
        tb_cc = add_textbox(slide, x + 0.2, start_y + 0.85, card_w - 0.4, card_h - 1.1)
        tf_cc = tb_cc.text_frame
        tf_cc.word_wrap = True
        for line in card_content.split('\n'):
            if line.strip():
                add_para(tf_cc, line.strip(), size=14, color=LIGHT, space_after=Pt(3))

    add_footer(slide)
    return slide


def make_image_slide(title, image_path, caption=''):
    """图片展示页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_header_bar(slide, title)
    add_side_decoration(slide)

    if os.path.exists(image_path):
        slide.shapes.add_picture(
            image_path,
            Inches(1.5), Inches(1.7),
            Inches(10.3), Inches(4.8)
        )
    else:
        tb_missing = add_textbox(slide, 1.5, 3.0, 10.3, 1.0)
        tf_missing = tb_missing.text_frame
        add_para(tf_missing, f'[图片未找到: {image_path}]', size=16, color=ACCENT2,
                 alignment=PP_ALIGN.CENTER)

    if caption:
        tb_cap = add_textbox(slide, 0.8, 6.6, 11.7, 0.4)
        tf_cap = tb_cap.text_frame
        add_para(tf_cap, caption, size=13, color=LIGHT, alignment=PP_ALIGN.CENTER)

    add_footer(slide)
    return slide


# ════════════════════════════════════════════════════════════
# 42 页 PPT 内容
# ════════════════════════════════════════════════════════════

# ── P0: 封面 ──
make_title_slide(
    'Java 开发中的 AI 协作编程实战',
    '从"氛围编程"到"驾驭工程" —— 来自 793+ 个真实会话的实践经验'
)

# ── P1: 真实数据冲击 ──
make_number_slide('我的 AI 编程实践——真实数据', [
    ('793+', 'Claude Code 会话'),
    ('50+', '活跃项目'),
    ('1,000~6,000', '日均消息量'),
    ('256', 'aip-portal 提交数'),
    ('24', 'Worktree 并行开发'),
    ('95%', 'Rules 压缩率'),
])

# ── P2: 核心观点 ──
make_content_slide('核心观点', [
    '>>AI 编程不是"让 AI 替你写代码"，而是"人掌舵，AI 执行"',
    '',
    '##角色转变',
    ('从"代码作者"到"系统设计师"', 1, LIGHT),
    ('从"氛围编程(Vibe Coding)"到"驾驭工程(Harness Engineering)"', 1, LIGHT),
    '',
    '##今天的分享',
    ('我在 AI 编程中踩过的坑', 1, LIGHT),
    ('以及如何建立起有效的协作体系', 1, ACCENT3),
])

# ── P3: 氛围编程 vs 驾驭工程 ──
make_two_col_slide(
    '从"氛围编程"到"驾驭工程"',
    '❌ 氛围编程 (Vibe Coding)',
    [
        '告诉 AI 需求，让它自由发挥',
        '不做技术选型，让 AI 自己决定',
        '缺少审查，直接采纳大量代码',
        'Git 记录：连续的 fix fix fix',
        '结果：快速生成 → 大量返工',
        '典型：A2A 适配器 17 天返工',
    ],
    '✅ 驾驭工程 (Harness Engineering)',
    [
        '人做技术选型和架构设计',
        '通过规则约束 AI 的工作范围',
        '小批量推进，逐步验证',
        'Git 记录：清晰的阶段划分',
        '结果：方向正确，返工减少',
        '典型：aip-portal 24 个模块并行',
    ]
)

# ── P4: 金句 ──
make_quote_slide(
    'AI 编程的核心不是"生成"，而是"驾驭"',
    '核心观点'
)

# ── P5: 第一部分标题 ──
make_section_slide(
    '一', '驾驭 AI——从混乱到有序',
    '我踩过的坑，以及怎么建立起有效的协作体系'
)

# ── P6: 蜜月期与翻车期 ──
make_two_col_slide(
    '蜜月期与翻车期',
    '兴奋的蜜月期',
    [
        'AI 快速生成大量代码',
        'A2A 适配器：一天写了 3840 行',
        '代码结构完整，编译能过',
        '感觉找到了"银弹"',
    ],
    '残酷的翻车现实',
    [
        'A2A 案例：不了解协议就开干',
        '手写 7600 行 → 发现官方 SDK → 推倒重来',
        '早期 fix fix fix 的 git 记录',
        '缺少审查导致大量返工',
        '关键转折：AI 编程的核心不是"生成"，而是"驾驭"',
    ]
)

# ── P7: Harness 马具比喻 ──
make_content_slide('Harness Engineering：给 AI 套上马具', [
    '##马具比喻',
    ('AI 模型 = 烈马（强大但难以控制方向）', 1, LIGHT),
    ('Harness = 马具（缰绳、鞍具）', 1, ACCENT),
    ('人 = 骑手（决定方向，控制节奏）', 1, ACCENT3),
    '',
    '>>缰绳不增加马的力量，但让力量用在对的方向',
    '',
    '##Harness 工程的三大支柱',
    ('1. 上下文工程（Context Engineering）：确保 AI 看到正确的信息', 0, WHITE),
    ('2. 架构约束（Architectural Constraints）：用工具机械地强制执行', 0, WHITE),
    ('3. 熵管理（Entropy Management）：控制代码库的混乱度', 0, WHITE),
])

# ── P8: 三大支柱详解 ──
make_table_slide(
    '三大支柱详解',
    ['支柱', '核心思想', '实践方式'],
    [
        ['上下文工程', 'AI 只能看到你给它的信息', '三层知识体系：CLAUDE.md + rules/ + memory/'],
        ['架构约束', '用工具强制执行，不指望 AI 记住', 'paths frontmatter、Hook 自动化、编译验证'],
        ['熵管理', '控制代码库混乱度', 'rules/ 精简（5487→272行）、定期清理'],
    ]
)

# ── P9: 行业佐证 ──
make_number_slide('行业佐证：Harness 的威力', [
    ('100 万行', 'OpenAI Codex\n3 名工程师 5 个月构建\n零手写'),
    ('Top 30 → Top 5', 'LangChain\n不改模型只改 Harness\n排名飙升'),
    ('256 次提交', '我的 aip-portal 项目\n24 个 worktree 并行'),
])

# ── P10: 三层知识体系 section ──
make_section_slide(
    '', '三层知识体系：我的实践架构',
    '给 AI 传递知识的最佳方式'
)

# ── P11: 三层架构图 ──（修复对齐）
make_code_slide(
    '三层知识体系架构',
    '+-------------------------------------------------------------+\n'
    '|  全局层 (~/.claude/CLAUDE.md)                                |\n'
    '|  > 跨项目的个人偏好（Git 规范、Shell 规范）                   |\n'
    '|  > 所有项目共享，长期稳定                                     |\n'
    '+-------------------------------------------------------------+\n'
    '|  项目层 (CLAUDE.md + .claude/rules/)                          |\n'
    '|  > 项目定位、技术栈、架构、编码规范                            |\n'
    '|  > ⚠ 每次对话全部加载，每行都在消耗 token                     |\n'
    '+-------------------------------------------------------------+\n'
    '|  会话层 (memory/)                                             |\n'
    '|  > 复盘经验、项目上下文、外部引用                             |\n'
    '|  > 按需检索，不自动全部加载                                   |\n'
    '|  > 自由积累经验，不担心占用 token                             |\n'
    '+-------------------------------------------------------------+'
)

# ── P12: rules/ 演变时间线 ──
make_timeline_slide(
    'rules/ 演变时间线——好的结构是"长"出来的',
    [
        ('Phase 1', '初始搭建\n5,487 行', '想给 AI 所有知识\n20 个文件一次创建'),
        ('Phase 2', '膨胀期\n~6,000+ 行', '功能驱动增长\n设计文档、执行计划堆积\ntoken 消耗爆炸'),
        ('Phase 3', '大清理\n-2,714 行', '核心转折：\n"rules 不是 wiki"\n移除人看的文档'),
        ('Phase 4', '持续精简\n最终 272 行', '压缩率 95%\n只保留 AI 必看的\n架构+规范索引+测试'),
    ]
)

# ── P13: 第二部分标题 ──
make_section_slide(
    '二', '实战工作流——设计、实现、测试、复盘闭环',
    'SDD（规范驱动开发）在 AI 编程中的落地'
)

# ── P14: PRD 质量决定成败 ──
make_two_col_slide(
    'PRD 质量决定成败',
    '清晰 PRD → AI 事半功倍',
    [
        '明确的功能描述和验收标准',
        '清晰的数据模型和接口定义',
        'Claude Code 输出实施计划',
        '人工审阅后启动实现',
        '方向正确，返工少',
    ],
    '模糊 PRD → AI 飞速跑偏',
    [
        '需求不明确，AI 自己发挥',
        '没有验收标准，不知道对不对',
        '无法判断输出质量',
        '越改越偏，返工频繁',
        '比手写代码更危险——有"已完成"错觉',
    ]
)

# ── P15: SDD 四阶段循环 ──（修复对齐）
make_code_slide(
    'SDD（规范驱动开发）四阶段循环',
    '  +----------------------------------------------------+\n'
    '  |                SDD 四阶段循环                       |\n'
    '  |                                                    |\n'
    '  |  1. Specify（写规范）                               |\n'
    '  |     -> IDE 写设计 Markdown                         |\n'
    '  |     -> 明确需求、数据模型、验收标准                 |\n'
    '  |           |                                        |\n'
    '  |  2. Plan（出计划）                                  |\n'
    '  |     -> Claude Code 输出实施计划                     |\n'
    '  |     -> 人工审阅确认                                 |\n'
    '  |           |                                        |\n'
    '  |  3. Task（拆任务）                                  |\n'
    '  |     -> 按计划拆分为独立任务                         |\n'
    '  |     -> 每个任务有明确产出物                         |\n'
    '  |           |                                        |\n'
    '  |  4. Implement（做实现）                             |\n'
    '  |     -> 按任务逐个实现                               |\n'
    '  |     -> 编译验证 + 接口测试                          |\n'
    '  |           |                                        |\n'
    '  |     回到 1，下一个功能                              |\n'
    '  +----------------------------------------------------+\n'
    '\n'
    '  规范文档是 AI 和人的"唯一真理来源"'
)

# ── P16: 设计文档示例 ──
make_content_slide('设计文档实践', [
    '##aip-portal 项目中的设计文档',
    ('14 个设计文档存放在 docs/design/', 1, LIGHT),
    ('每个功能先写设计，再让 AI 实现', 1, LIGHT),
    '',
    '##典型设计文档',
    ('SQL 版本管理 Starter 设计 —— 547 行', 0, WHITE),
    ('Karate 接口实践 —— 668 行', 0, WHITE),
    ('消息管理重构设计 —— 283 行', 0, WHITE),
    '',
    '>>实践流程：IDE 写设计 MD → Claude Code 输出实施计划 → 人工审阅',
])

# ── P17: 实现阶段 section ──
make_section_slide(
    '', '实现阶段：多种协作模式',
    '按任务复杂度选择合适的模式'
)

# ── P18: 协作模式选择矩阵 ──
make_table_slide(
    '协作模式选择矩阵——按复杂度选模式',
    ['任务复杂度', '推荐模式', '适用场景', '关键特点'],
    [
        ['简单修改', '直接对话', '单文件修改、小 bug', '最快，即改即验证'],
        ['中等功能', 'Plan 模式', '需探索代码再实现', '先探索 → 确认 → 再实现'],
        ['复杂功能', 'Sub Agent', '需拆分子任务', '控制上下文，避免跑偏'],
        ['多模块功能', 'Agent Team', '多 Agent 并行开发', '5人团队，分工协作'],
        ['多任务并行', 'Git Worktree', '多个需求同时推进', '物理隔离，独立上下文'],
    ]
)

# ── P19: Plan 模式工作流 ──
make_content_slide('Plan 模式工作流', [
    '##工作流程',
    ('用户需求 → Plan 模式启动', 0, WHITE),
    ('       → 探索代码库（Explore Agent）', 1, LIGHT),
    ('       → 制定实施计划', 1, LIGHT),
    ('       → 用户审阅确认 ← 关键环节', 1, ACCENT4),
    ('       → 执行实现', 1, LIGHT),
    '',
    '##实际案例：搜索权限过滤逻辑',
    ('同时搜索 !isAdmin 变量名', 1, LIGHT),
    ('搜索中文注释（"权限"、"角色"）', 1, LIGHT),
    ('搜索 isUserWorkspaceAdmin 方法名', 1, LIGHT),
    ('交叉验证，避免单一模式遗漏', 1, ACCENT3),
])

# ── P20: Git Worktree 并行开发 ──
image_path = os.path.join(BASE_DIR, 'my-ideas', 'images', 'worktree', '多个worktree并行.png')
make_image_slide(
    'Git Worktree：24 个并行开发',
    image_path,
    '每个 Worktree 有独立工作目录和 Claude Code 会话，互不干扰'
)

# ── P21: Agent Team 5人团队 ──（修复对齐）
make_code_slide(
    'Agent Team：5 人团队协作',
    '                     +------------------+\n'
    '                     |    Team Lead     |\n'
    '                     |    (你/用户)     |\n'
    '                     +--------+---------+\n'
    '                              |\n'
    '              +---------------+---------------+\n'
    '              |               |               |\n'
    '              v               v               v\n'
    '     +----------------+ +----------------+ +----------------+\n'
    '     |  需求分析师     | |  表结构设计师   | |  接口设计师     |\n'
    '     +----------------+ +----------------+ +----------------+\n'
    '              |               |               |\n'
    '              +---------------+---------------+\n'
    '                              |\n'
    '                              v\n'
    '                     +------------------+\n'
    '                     |   计划工程师      |\n'
    '                     +--------+---------+\n'
    '                              |\n'
    '                              v\n'
    '                     +------------------+\n'
    '                     |   评审验收        |\n'
    '                     +------------------+'
)

# ── P22: 测试阶段 section ──
make_section_slide(
    '', '测试阶段',
    '单元测试 + Karate 业务接口测试'
)

# ── P23: 测试策略 ──
make_content_slide('测试策略', [
    '##单元测试',
    ('AI 生成 + 人工审查', 1, LIGHT),
    ('注意维护成本：AI 生成的测试有时过于脆弱', 1, ACCENT2),
    '',
    '##Karate 业务接口测试（aip-portal 实践）',
    ('三层数据隔离策略', 1, LIGHT),
    ('标签体系：@smoke / @p0 / @p1 / @p2 / @negative', 1, LIGHT),
    ('执行：mvn test -Dtest=KarateRunner -Dkarate.tags=@smoke', 1, LIGHT),
    '',
    '##TDD 在 AI 编程中的价值',
    ('测试通过是任务完成的硬性门槛', 1, ACCENT3),
    ('避免"看起来完成了"但实际有问题的陷阱', 1, ACCENT3),
])

# ── P24: 复盘驱动改进 section ──
make_section_slide(
    '', '复盘驱动改进',
    '/retrospect 复盘三问法'
)

# ── P25: /retrospect 复盘 ──
make_content_slide('/retrospect 复盘三问法', [
    '##三问法',
    ('1. 做对了什么？→ 强化正确的做法', 0, ACCENT3),
    ('2. 做错了什么？→ 识别错误和根因', 0, ACCENT2),
    ('3. 如果重来会怎么做？→ 提炼可复用规则', 0, ACCENT),
    '',
    '##经验固化流程',
    ('完成任务 → 触发复盘 → 三问反思 → 提炼经验', 0, WHITE),
    ('                                    → 保存到 memory/', 1, LIGHT),
    ('                                    → 更新 rules/', 1, LIGHT),
    ('                                    → 更新 CLAUDE.md', 1, LIGHT),
    '',
    '##三大翻车案例',
    ('SQL 迁移版本号错误：假设验证失败', 1, ACCENT2),
    ('权限过滤逻辑误删：职责理解不深入', 1, ACCENT2),
    ('修改已冻结 Flyway 脚本：违反不可变原则', 1, ACCENT2),
])

# ── P26: 第三部分标题 ──
make_section_slide(
    '三', '深度案例——A2A 适配器的教训',
    '用一个完整案例串联所有经验'
)

# ── P27: A2A 时间线 ──
make_timeline_slide(
    'A2A 适配器时间线：从手写到 SDK',
    [
        ('Day 1', '自信开干\n3840 行', '不了解协议\n让 AI 自由发挥\n编译能过但不对'),
        ('Day 2-8', '不断修补\n膨胀到 7600', '发现问题→修补\n越补越偏\n规范符合性修复'),
        ('Day 9', '转折！\n发现 SDK', '终于看官方仓库\n发现 a2a-java SDK\n有 Spring WebFlux 示例'),
        ('Day 10-17', 'SDK 重写\n2400 行', '基于 SDK 重写\n代码量减少 68%\n但已浪费 17 天'),
    ]
)

# ── P28: 返工代价量化 ──
make_table_slide(
    '返工代价量化对比',
    ['维度', '旧方案（手写）', '新方案（基于 SDK）'],
    [
        ['代码量', '~7,600 行', '~2,400 行（减少 68%）'],
        ['开发时间', '9 天（含返工）', '9 天（含学习 SDK）'],
        ['协议正确性', '多处不一致', '由 SDK 保证'],
        ['维护成本', '高（需跟进规范变化）', '低（SDK 更新即可）'],
        ['核心实现', '自己写的 JsonRpcHandler', 'SDK 的 RestHandler'],
    ]
)

# ── P29: 根因分析 ──（修复对齐）
make_code_slide(
    '根因分析链',
    '不了解 A2A 协议\n'
    '  -> 没有阅读官方文档\n'
    '    -> 不知道有官方 SDK\n'
    '      -> 让 AI "自由发挥"实现协议\n'
    '        -> AI 生成的代码结构完整但细节不对\n'
    '          -> 不断修补\n'
    '            -> 越补越偏\n'
    '              -> 最终推倒重来\n'
    '\n'
    '核心问题：\n'
    '  X AI 不会主动告诉你"有官方 SDK 可以直接用"\n'
    '  X 对不了解的领域，无法判断 AI 输出的正确性\n'
    '  X AI 的输出看起来越合理，就越危险'
)

# ── P30: 三个关键教训 ──
make_card_slide(
    '三个关键教训',
    [
        ('教训 1：先调研再动手',
         '官方有 SDK 吗？\n有示例吗？\n社区有方案吗？\n\n如果先花 2 小时调研\n就能避免 17 天返工',
         ACCENT),
        ('教训 2：AI 不会替你做技术选型',
         'AI 不会告诉你\n"有官方 SDK 可以用"\n\nAI 不会质疑你的方案\n是否是最优解\n\n技术选型是人的责任',
         ACCENT2),
        ('教训 3：不了解的领域更危险',
         '熟悉的领域：\nAI 输出质量高\n\n不了解的领域：\n无法判断正确性\n\n越不了解，越要先建认知',
         ACCENT4),
    ]
)

# ── P31: 如果重来 ──
make_content_slide('如果重来会怎么做？', [
    '>>正确做法：总计 3 天，而不是 17 天',
    '',
    '##Day 1 上午（2 小时）：调研',
    ('阅读 A2A 官方 GitHub 仓库', 1, LIGHT),
    ('发现官方有 Java SDK + Spring WebFlux 示例', 1, ACCENT3),
    '',
    '##Day 1 下午 ~ Day 2：学习',
    ('学习 SDK 核心接口：AgentExecutor、RestHandler、TaskStore', 1, LIGHT),
    ('运行官方示例代码', 1, LIGHT),
    '',
    '##Day 3：适配',
    ('基于 SDK 实现 ShenYu 网关集成', 1, LIGHT),
    ('参考官方示例调通', 1, LIGHT),
    '',
    '!!17 天的返工，如果先花 2 小时调研，可能 3 天就搞定了',
])

# ── P32: 第四部分标题 ──
make_section_slide(
    '四', '工具链与进阶技巧',
    '日常开发 / 系统调研 / 学习研究'
)

# ── P33: 工具组合矩阵 ──
make_table_slide(
    '工具组合矩阵',
    ['场景', '工具组合', '说明'],
    [
        ['日常开发', 'Claude Code + cc-switch\n+ jdtls-lsp + Apifox MCP', '主力编码、模型切换\nJava 语言服务、API 文档查询'],
        ['系统调研', 'Claude Code + Playwright MCP\n+ Web Reader', '浏览器自动化、自动截图\n网页内容抓取、生成报告'],
        ['学习研究', 'Claude Code + Web Search\n+ guided-learning-tutor', '苏格拉底式引导学习\n搜索资料、深度阅读'],
    ]
)

# ── P34: MCP 扩展 + 自定义技能 ──
make_content_slide('MCP 扩展与自定义技能', [
    '##MCP 服务',
    ('Apifox MCP：开发接口时直接查询 API 定义', 1, LIGHT),
    ('Playwright MCP：浏览器自动化、系统调研、自动截图', 1, LIGHT),
    ('Web Reader：抓取在线文档、技术资料', 1, LIGHT),
    '',
    '##自定义技能',
    ('guided-learning-tutor：苏格拉底式学习引导', 1, ACCENT),
    ('  → 学生分类、三级掌握标准、五步交互', 1, DARK_LIGHT),
    ('retrospect：复盘技能，三问法 + 经验固化', 1, ACCENT),
    ('  → 完成任务 → 三问反思 → 提炼规则 → 保存到 memory', 1, DARK_LIGHT),
    ('commit-and-push：自动提交推送流程', 1, ACCENT),
    ('  → 分析变更 → 生成 message → 排除敏感文件 → 提交', 1, DARK_LIGHT),
])

# ── P35: 模型选择策略 ──
make_table_slide(
    '模型选择策略',
    ['场景', '推荐模型', '原因'],
    [
        ['日常编码', 'Sonnet', '性价比最高，速度最快'],
        ['系统设计、架构规划', 'Opus', '推理能力更强，适合复杂决策'],
        ['反复尝试无果时', '切换到 Opus', '更强模型可能一次性解决'],
        ['简单搜索、快速查询', 'Haiku', '速度快，成本低'],
    ]
)

# ── P36: AI 其他应用 ──
make_content_slide('AI 的更多应用场景', [
    '##解读开源项目',
    ('AG-UI 协议研究：21 篇深度文档', 1, LIGHT),
    ('从简介到 Java SDK 实战案例，系统化输出', 1, LIGHT),
    '',
    '##AI + Playwright 深入解读系统',
    ('自动截图、填表、点击', 1, LIGHT),
    ('生成用户手册和系统分析报告', 1, LIGHT),
    '',
    '##AI 自学新技术',
    ('Python 完整学习路径：16 模块，60+ 主题', 1, LIGHT),
    ('苏格拉底式引导，主动验证理解', 1, LIGHT),
])

# ── P37: 起步三步走 ──
make_content_slide('给同事的实操建议——起步三步走', [
    '##第一天：创建 CLAUDE.md，写三行',
    ('项目是什么', 1, ACCENT3),
    ('用什么技术栈', 1, ACCENT3),
    ('怎样算好的代码', 1, ACCENT3),
    '',
    '##第一周：遇到问题就加一条 rule',
    ('不要提前设计完美的规则体系', 1, ACCENT4),
    ('每一条规则都应有一个具体的"犯错事件"作为来源', 1, ACCENT4),
    '',
    '##第一个月：回顾 rules/，做一次大清理',
    ('删除可从代码推导的信息', 1, ACCENT),
    ('移除"人看"的文档到 docs/', 1, ACCENT),
    ('确认每条规则都有存在的价值', 1, ACCENT),
])

# ── P38: Harness 生长循环 ──（修复对齐）
make_code_slide(
    'Harness 生长循环',
    '  +----------------------------------------------------+\n'
    '  |                                                    |\n'
    '  |  1. 让 AI 做任务（尽量少约束）                     |\n'
    '  |         |                                          |\n'
    '  |  2. 观察 AI 的输出                                 |\n'
    '  |         |                                          |\n'
    '  |  3. 发现不符合预期的地方                            |\n'
    '  |         |                                          |\n'
    '  |  4. 分析根因（缺上下文？缺规则？）                  |\n'
    '  |         |                                          |\n'
    '  |  5. 添加最小约束（只解决当前问题）                  |\n'
    '  |         |                                          |\n'
    '  |  6. 验证效果                                       |\n'
    '  |         |                                          |\n'
    '  |     回到 1，循环往复                               |\n'
    '  |                                                    |\n'
    '  +----------------------------------------------------+\n'
    '\n'
    '  Harness 不是设计出来的，而是被经验喂养出来的\n'
    '  每次只添加解决当前问题的最小约束'
)

# ── P39: 心态转变 ──
make_table_slide(
    '心态转变',
    ['从 ❌', '到 ✅'],
    [
        ['代码作者', '系统设计师'],
        ['追求一步到位', '先有再优化'],
        ['自己记住经验', '让项目记住经验'],
        ['文档越全越好', '信息密度越高越好'],
    ]
)

# ── P40: 避坑指南 ──
make_card_slide(
    '避坑指南',
    [
        ('警惕"Markdown 怪兽"',
         '规范要短小、原子化\n可验证\n\n人看的放 docs/\nAI 必须看的放 rules/',
         ACCENT),
        ('不要过度微观管理',
         '明确意图和约束后\n给 AI 发挥空间\n\n约束优于提示\n用工具强制执行',
         ACCENT4),
        ('两轮失败原则',
         '修正超过两次还没过\n果断 /clear 重来\n\n不要在低效循环里死磕\n考虑切换更强模型',
         ACCENT2),
        ('先建认知再辅助',
         '越不了解的领域\n越要先建立认知\n\n不了解的东西\nAI 输出看起来越合理越危险',
         ACCENT3),
    ]
)

# ── P41: 结束页 ──
slide = make_quote_slide(
    '规范先行，代码自动。\n从"写代码的人"变成"设计系统的人"。',
    ''
)
# 添加底部文字
tb_end = add_textbox(slide, 1.5, 6.2, 10.3, 0.5)
tf_end = tb_end.text_frame
add_para(tf_end, '感谢聆听  |  Q&A 交流时间', size=16, color=LIGHT, alignment=PP_ALIGN.CENTER)


# ── 保存文件 ──
output_dir = os.path.join(BASE_DIR, 'output', 'design2')
os.makedirs(output_dir, exist_ok=True)
output_path = os.path.join(output_dir, 'AI协作编程实战v2.pptx')
prs.save(output_path)
print(f'PPT 已生成: {output_path}')
print(f'共 {len(prs.slides)} 页')
