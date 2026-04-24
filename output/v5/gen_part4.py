"""
Part 4: Tools & Tips + Closing (P38-P48)
"""

from pptx.util import Cm, Pt
from pptx.enum.text import PP_ALIGN
from gen_helpers import *

def add_slides(prs):
    # ══════════════════════════════════════════════════════════
    # P38 - Part 4 Section Divider
    # ══════════════════════════════════════════════════════════
    slide = add_slide(prs)
    add_textbox(slide, Cm(3), Cm(5.0), SLIDE_W - Cm(6), Cm(1.2),
                "Part 04", 20, ACCENT_BLUE, align=PP_ALIGN.CENTER)
    add_textbox(slide, Cm(3), Cm(7.0), SLIDE_W - Cm(6), Cm(2.0),
                "\u5de5\u5177\u94fe\u4e0e\u5b9e\u6218\u6280\u5de7", 40, TEXT_WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_textbox(slide, Cm(3), Cm(10.0), SLIDE_W - Cm(6), Cm(1.2),
                "\u53ef\u4ee5\u7acb\u5373\u5e26\u8d70\u7684\u4e1c\u897f", 18, TEXT_GREY, align=PP_ALIGN.CENTER)

    # ══════════════════════════════════════════════════════════
    # P39 - Tool: Claude Code
    # ══════════════════════════════════════════════════════════
    slide = add_slide(prs)
    add_slide_title(slide, "Claude Code \u2014 AI \u7f16\u7a0b\u7684\u7ec8\u7aef\u5de5\u5177")

    features = [
        ("\U0001f5a5\ufe0f \u7ec8\u7aef\u96f6\u8ddd\u79bb", "\u76f4\u63a5\u5728\u7ec8\u7aef\u5de5\u4f5c\uff0c\u4e0e\u9879\u76ee\u4ee3\u7801\u96f6\u8ddd\u79bb"),
        ("\U0001f504 \u5168\u6d41\u7a0b\u95ed\u73af", "\u652f\u6301\u8bfb\u53d6\u3001\u7f16\u8f91\u3001\u6267\u884c\uff0c\u5168\u6d41\u7a0b\u95ed\u73af"),
        ("\U0001f4cb \u81ea\u52a8\u52a0\u8f7d\u4e0a\u4e0b\u6587", "\u81ea\u52a8\u52a0\u8f7d CLAUDE.md \u548c rules/\uff0c\u7406\u89e3\u9879\u76ee\u4e0a\u4e0b\u6587"),
        ("\U0001f680 \u9ad8\u7ea7\u529f\u80fd", "Plan\u6a21\u5f0f / Worktree\u5e76\u884c / Agent Team \u534f\u4f5c\u7b49"),
    ]
    # 2x2 layout
    card_w = Cm(14.5)
    card_h = Cm(3.5)
    positions = [(MARGIN_L, Cm(3.2)), (MARGIN_L + card_w + Cm(0.8), Cm(3.2)),
                 (MARGIN_L, Cm(7.5)), (MARGIN_L + card_w + Cm(0.8), Cm(7.5))]
    for i, (title, desc) in enumerate(features):
        x, y = positions[i]
        add_card(slide, x, y, card_w, card_h, title, [desc],
                 title_size=14, desc_size=12, accent_color=ACCENT_BLUE)

    add_page_number(slide, 39)

    # ══════════════════════════════════════════════════════════
    # P40 - Tool: cc-switch
    # ══════════════════════════════════════════════════════════
    slide = add_slide(prs)
    add_slide_title(slide, "cc-switch \u2014 Claude Code \u914d\u7f6e\u5207\u6362\u5de5\u5177")

    add_subtitle(slide, "\u5feb\u901f\u5207\u6362\u4e0d\u540c\u7684 Claude Code \u914d\u7f6e\uff08\u6a21\u578b\u9009\u62e9\u3001\u6743\u9650\u8bbe\u7f6e\u7b49\uff09", Cm(2.4), 13)

    # Three scenario cards
    scenarios = [
        ("\U0001f4c1 \u4e0d\u540c\u9879\u76ee\u9700\u8981\u4e0d\u540c\u914d\u7f6e", ACCENT_BLUE),
        ("\U0001f9e9 \u590d\u6742\u4efb\u52a1\u7528\u5f3a\u6a21\u578b\uff0c\u7b80\u5355\u4efb\u52a1\u7528\u5feb\u6a21\u578b", GREEN),
        ("\U0001f465 \u56e2\u961f\u6210\u5458\u5171\u4eab\u914d\u7f6e", YELLOW),
    ]
    y = Cm(4.5)
    for text, color in scenarios:
        add_card(slide, MARGIN_L + Cm(2), y, SLIDE_W - MARGIN_L - MARGIN_R - Cm(4), Cm(2.2),
                 text, [], title_size=14, accent_color=color, title_color=color)
        y += Cm(2.8)

    add_page_number(slide, 40)

    # ══════════════════════════════════════════════════════════
    # P41 - Tip 1: Small Steps, Continuous Verification
    # ══════════════════════════════════════════════════════════
    slide = add_slide(prs)
    add_slide_title(slide, "\u6280\u5de7\u4e00\uff1a\u5c0f\u6b65\u5feb\u8dd1\uff0c\u6301\u7eed\u9a8c\u8bc1")

    # Cycle diagram: Design -> Implement -> Test -> Design
    cycle_items = [
        ("\u8bbe\u8ba1", ACCENT_BLUE),
        ("\u5b9e\u73b0", GREEN),
        ("\u6d4b\u8bd5", YELLOW),
    ]
    node_w = Cm(7)
    arrow_w = Cm(1.5)
    sx = Cm(5)
    y_cyc = Cm(3.5)
    for i, (text, color) in enumerate(cycle_items):
        x = sx + i * (node_w + arrow_w)
        add_flow_node(slide, x, y_cyc, node_w, Cm(2.0), text, "",
                      fill_color=CARD_BG, border_color=color,
                      text_color=color, text_size=14)
        if i < 2:
            add_arrow(slide, x + node_w + Cm(0.15), y_cyc + Cm(0.6), Cm(1.2), Cm(0.6), TEXT_GREY)
    # loop back arrow text
    add_textbox(slide, sx + 2 * (node_w + arrow_w) + node_w + Cm(0.3), y_cyc + Cm(0.3),
                Cm(3), Cm(1.0), "\u21a9 \u5faa\u73af", 12, TEXT_GREY)

    # Three key points
    points = [
        "\u2022 AI \u751f\u6210\u5feb \u2192 \u9519\u8bef\u79ef\u7d2f\u4e5f\u5feb \u2192 \u4e00\u6b21\u8d2a\u591a\u53cd\u800c\u66f4\u6162",
        "\u2022 \u6bcf\u5b8c\u6210\u4e00\u4e2a\u6a21\u5757\u5c31\u9a8c\u8bc1\uff0c\u800c\u4e0d\u662f\u5168\u90e8\u5b8c\u6210\u540e\u518d\u6d4b",
        "\u2022 \u9047\u5230\u95ee\u9898\u518d\u201c\u78e8\u5200\u201d\uff0c\u4e0d\u63d0\u524d\u8fc7\u5ea6\u8bbe\u8ba1",
    ]
    y = Cm(7.0)
    for p in points:
        add_textbox(slide, MARGIN_L + Cm(1), y, SLIDE_W - MARGIN_L - MARGIN_R - Cm(2), Cm(0.8),
                    p, 13, TEXT_WHITE)
        y += Cm(1.4)

    # Bottom lesson
    add_textbox(slide, MARGIN_L, Cm(12.0), SLIDE_W - MARGIN_L - MARGIN_R, Cm(0.7),
                "aip-portal \u65e9\u671f\u4e00\u6b21\u505a\u5927\u6a21\u5757\u5bfc\u81f4\u6279\u91cf\u8fd4\u5de5",
                11, TEXT_GREY, align=PP_ALIGN.CENTER)

    add_page_number(slide, 41)

    # ══════════════════════════════════════════════════════════
    # P42 - Tip 2: Build Cognition First
    # ══════════════════════════════════════════════════════════
    slide = add_slide(prs)
    add_slide_title(slide, "\u6280\u5de7\u4e8c\uff1a\u5148\u5efa\u7acb\u8ba4\u77e5\uff0c\u518d\u8ba9 AI \u5e72\u6d3b")

    # Core principle
    add_textbox(slide, Cm(2), Cm(3.0), SLIDE_W - Cm(4), Cm(1.0),
                "\u4f60\u7684\u8ba4\u77e5\u8fb9\u754c = AI \u6548\u7387\u4e0a\u9650",
                22, ACCENT_BLUE, bold=True, align=PP_ALIGN.CENTER)

    # Two column comparison
    lw = Cm(14)
    add_card(slide, MARGIN_L, Cm(5.0), lw, Cm(4.0),
             "\u274c \u76f4\u63a5\u8ba9AI\u5e72\uff08\u4e0d\u4e86\u89e3\u9886\u57df\uff09",
             ["A2A\u6848\u4f8b / 16\u5929 / 141,026\u884c\u5e9f\u5f03"],
             title_size=14, desc_size=12, title_color=RED_ORANGE,
             accent_color=RED_ORANGE, desc_color=TEXT_WHITE)
    rx = Cm(16.5)
    rw = SLIDE_W - rx - MARGIN_R
    add_card(slide, rx, Cm(5.0), rw, Cm(4.0),
             "\u2705 \u5148\u8c03\u7814\u5efa\u7acb\u8ba4\u77e5\uff082\u5c0f\u65f6\uff09",
             ["\u7528SDK / 3\u5929 / \u4ec51,786\u884c"],
             title_size=14, desc_size=12, title_color=GREEN,
             accent_color=GREEN, desc_color=TEXT_WHITE)

    # Bottom advice
    add_textbox(slide, MARGIN_L, Cm(11.0), SLIDE_W - MARGIN_L - MARGIN_R, Cm(0.8),
                "\u65b0\u9886\u57df\u5148\u8ba9 AI \u5e2e\u4f60\u8c03\u7814\u3001\u68b3\u7406\u8ba4\u77e5\u6846\u67b6\uff0c\u518d\u8fdb\u5165\u7f16\u7801",
                13, TEXT_WHITE, align=PP_ALIGN.CENTER)

    add_page_number(slide, 42)

    # ══════════════════════════════════════════════════════════
    # P43 - Tip 3: Collaboration Mode Matrix
    # ══════════════════════════════════════════════════════════
    slide = add_slide(prs)
    add_slide_title(slide, "\u6280\u5de7\u4e09\uff1a\u534f\u4f5c\u6a21\u5f0f\u9009\u62e9\u77e9\u9635")

    headers = ["\u573a\u666f", "\u63a8\u8350\u6a21\u5f0f", "\u8bf4\u660e"]
    rows = [
        ["\u660e\u786e\u7684\u5c0f\u4efb\u52a1", "\u76f4\u63a5\u5bf9\u8bdd", "\u7ed9\u6e05\u695a\u9700\u6c42\uff0cAI\u76f4\u63a5\u505a"],
        ["\u590d\u6742\u529f\u80fd\u5f00\u53d1", "Plan \u6a21\u5f0f", "\u5148\u8ba9AI\u89c4\u5212\uff0c\u786e\u8ba4\u540e\u6267\u884c"],
        ["\u591a\u4efb\u52a1\u5e76\u884c", "Worktree + Agent Team", "\u9694\u79bb\u5f00\u53d1\uff0c\u4e92\u4e0d\u5e72\u6270"],
        ["\u65b0\u9886\u57df\u63a2\u7d22", "\u8c03\u7814\u5148\u884c", "\u5148\u8ba9AI\u5e2e\u4f60\u5efa\u8ba4\u77e5"],
        ["\u4ee3\u7801\u8bc4\u5ba1", "\u53ea\u8bfb\u6a21\u5f0f", "AI\u8bfb\u4ee3\u7801\uff0c\u4eba\u505a\u51b3\u7b56"],
    ]
    add_table(slide, MARGIN_L, Cm(3.0), SLIDE_W - MARGIN_L - MARGIN_R,
              [Cm(7), Cm(10), Cm(13)], headers, rows)

    add_page_number(slide, 43)

    # ══════════════════════════════════════════════════════════
    # P44 - Quick Start for Beginners
    # ══════════════════════════════════════════════════════════
    slide = add_slide(prs)
    add_slide_title(slide, "\u7acb\u5373\u53ef\u7528\uff1a\u65b0\u624b\u5165\u95e8\u4e09\u6b65")

    # Three step cards
    steps = [
        ("Step 1", "\u4ece\u4e00\u4e2a\u5c0f\u6a21\u5757\u6216\u5c0f\u9879\u76ee\u5f00\u59cb", ACCENT_BLUE),
        ("Step 2", "\u5199\u597d\u4f60\u7684\u7b2c\u4e00\u4e2a CLAUDE.md", YELLOW),
        ("Step 3", "\u517b\u6210\u201c\u8bbe\u8ba1 \u2192 AI\u5b9e\u73b0 \u2192 \u9a8c\u8bc1\u201d\u7684\u4e60\u60ef", GREEN),
    ]
    card_w = Cm(9.5)
    gap = Cm(0.5)
    for i, (step, desc, color) in enumerate(steps):
        x = MARGIN_L + i * (card_w + gap)
        add_card(slide, x, Cm(3.2), card_w, Cm(3.5), step, [desc],
                 title_size=16, desc_size=12, accent_color=color, title_color=color)

    # Step 2 template example (code block)
    add_textbox(slide, MARGIN_L, Cm(8.0), SLIDE_W - MARGIN_L - MARGIN_R, Cm(0.8),
                "CLAUDE.md \u6700\u5c0f\u6a21\u677f\u793a\u4f8b\uff1a", 12, TEXT_GREY)
    code_lines = [
        "\u6280\u672f\u6808\uff1aJava 17 + Spring Boot 2.7 + PostgreSQL",
        "\u7ea6\u5b9a\uff1a\u7edf\u4e00\u5206\u9875 TypedSearchParam + PageResult",
        "API\u8fd4\u56de\uff1aResult<T> \u5305\u88c5",
    ]
    add_code_block(slide, MARGIN_L, Cm(9.0), SLIDE_W - MARGIN_L - MARGIN_R, Cm(2.5), code_lines, 11)

    add_page_number(slide, 44)

    # ══════════════════════════════════════════════════════════
    # P45 - Quick Start for Advanced Users
    # ══════════════════════════════════════════════════════════
    slide = add_slide(prs)
    add_slide_title(slide, "\u8fdb\u9636\uff1a\u642d\u5efa\u9879\u76ee\u7ea7 Harness")

    tips = [
        ("1", "\u642d\u5efa Harness \u5206\u5c42\u7ed3\u6784",
         ["CLAUDE.md\uff08\u5bfc\u822a\uff0c20-30\u884c\uff09",
          "rules/ARCHITECTURE.md\uff08\u67b6\u6784\u7ea6\u675f\uff09",
          "docs/\uff08\u8be6\u7ec6\u89c4\u8303\uff0c\u6309\u9700\u5f15\u7528\uff09"]),
        ("2", "\u5c1d\u8bd5 Worktree \u5e76\u884c\u5f00\u53d1",
         ["aip-portal \u7528\u4e8624\u4e2aworktree"]),
        ("3", "\u5efa\u7acb\u590d\u76d8\u4e60\u60ef",
         ["\u8e29\u5751 \u2192 \u8865\u89c4\u8303 \u2192 \u5b9a\u671f\u7cbe\u7b80"]),
    ]
    y = Cm(3.0)
    for num, title, lines in tips:
        card_h = Cm(2.0) + len(lines) * Cm(0.7)
        add_card(slide, MARGIN_L + Cm(1), y, SLIDE_W - MARGIN_L - MARGIN_R - Cm(2), card_h,
                 f"  {num}. {title}", lines,
                 title_size=14, accent_color=ACCENT_BLUE, title_color=ACCENT_BLUE,
                 desc_size=12, desc_color=TEXT_WHITE)
        y += card_h + Cm(0.8)

    add_page_number(slide, 45)

    # ══════════════════════════════════════════════════════════
    # P46 - Quick Start for Observers
    # ══════════════════════════════════════════════════════════
    slide = add_slide(prs)
    add_slide_title(slide, "\u7ed9\u89c2\u671b\u8005\uff1a\u4ece\u4e00\u4e2a\u6700\u5c0f\u573a\u666f\u5f00\u59cb")

    # Three scenario cards
    scenarios = [
        ("\U0001f9ea \u573a\u666f1", "\u8ba9 AI \u5199\u5355\u5143\u6d4b\u8bd5", "\u6700\u4f4e\u98ce\u9669"),
        ("\U0001f50d \u573a\u666f2", "\u8ba9 AI \u505a\u4ee3\u7801\u8bc4\u5ba1", "\u53ea\u8bfb\uff0c\u98ce\u9669\u4e3a\u96f6"),
        ("\u26a1 \u573a\u666f3", "\u8ba9 AI \u5199 CRUD \u63a5\u53e3", "\u5feb\u901f\u4f53\u9a8c\u6548\u7387\u63d0\u5347"),
    ]
    card_w = Cm(9.5)
    gap = Cm(0.5)
    for i, (title, desc, risk) in enumerate(scenarios):
        x = MARGIN_L + i * (card_w + gap)
        add_card(slide, x, Cm(3.2), card_w, Cm(4.5), title, [desc, risk],
                 title_size=14, desc_size=12, accent_color=GREEN, title_color=GREEN,
                 desc_color=TEXT_WHITE)

    # Bottom note
    add_textbox(slide, MARGIN_L, Cm(10.0), SLIDE_W - MARGIN_L - MARGIN_R, Cm(0.8),
                "\u611f\u53d7\u5230\u6548\u7387\u63d0\u5347\u540e\uff0c\u81ea\u7136\u4f1a\u6269\u5927\u4f7f\u7528\u8303\u56f4",
                13, TEXT_GREY, align=PP_ALIGN.CENTER)

    add_page_number(slide, 46)

    # ══════════════════════════════════════════════════════════
    # P47 - One Page Summary
    # ══════════════════════════════════════════════════════════
    slide = add_slide(prs)
    add_slide_title(slide, "\u603b\u7ed3")

    summaries = [
        "1. AI \u6548\u7387\u786e\u5b9a\uff0c\u4e0d\u786e\u5b9a\u7684\u662f\u4eba\u7684\u5f15\u5bfc",
        "2. \u628a\u6700\u64c5\u957f\u7684\u4e8b\u4ea4\u7ed9 AI\uff0c\u91ca\u653e\u65f6\u95f4\u505a\u66f4\u6709\u4ef7\u503c\u7684\u4e8b",
        "3. \u7528 SDD + Harness \u7cfb\u7edf\u5316\u5730\u9a7e\u9a6d AI",
        "4. \u4ece\u201c\u5199\u4ee3\u7801\u7684\u4eba\u201d\u53d8\u6210\u201c\u8bbe\u8ba1\u7cfb\u7edf\u7684\u4eba\u201d",
    ]
    y = Cm(3.5)
    for s in summaries:
        add_textbox(slide, MARGIN_L + Cm(1), y, SLIDE_W - MARGIN_L - MARGIN_R - Cm(2), Cm(1.2),
                    s, 20, ACCENT_BLUE, bold=True)
        y += Cm(2.5)

    add_page_number(slide, 47)

    # ══════════════════════════════════════════════════════════
    # P48 - Q&A
    # ══════════════════════════════════════════════════════════
    slide = add_slide(prs)

    # Big centered text
    add_textbox(slide, Cm(3), Cm(5.0), SLIDE_W - Cm(6), Cm(2.5),
                "\u611f\u8c22\u804a\u542c", 40, TEXT_WHITE, bold=True, align=PP_ALIGN.CENTER)

    add_textbox(slide, Cm(3), Cm(8.0), SLIDE_W - Cm(6), Cm(1.5),
                "Q & A", 28, ACCENT_BLUE, bold=True, align=PP_ALIGN.CENTER)

    # Bottom contact placeholder
    add_textbox(slide, Cm(3), SLIDE_H - Cm(2.5), SLIDE_W - Cm(6), Cm(0.8),
                "\u8054\u7cfb\u65b9\u5f0f\u5360\u4f4d\u7b26", 12, TEXT_GREY, align=PP_ALIGN.CENTER)
