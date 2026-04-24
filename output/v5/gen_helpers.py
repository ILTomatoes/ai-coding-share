"""
Shared design helpers for PPT generation.
All color constants, font settings, shape builders.
"""

from pptx import Presentation
from pptx.util import Cm, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import copy

# ─── Color constants ───
BG_DARK      = RGBColor(0x0D, 0x11, 0x17)  # deep blue-black background
CARD_BG      = RGBColor(0x16, 0x1B, 0x22)  # card background
BORDER_COLOR = RGBColor(0x30, 0x36, 0x3D)  # border grey
ACCENT_BLUE  = RGBColor(0x2F, 0x81, 0xF7)  # bright blue accent
GREEN        = RGBColor(0x3F, 0xB9, 0x50)  # success green
RED_ORANGE   = RGBColor(0xF8, 0x51, 0x49)  # warning red-orange
YELLOW       = RGBColor(0xE3, 0xB3, 0x41)  # caution yellow
TEXT_WHITE    = RGBColor(0xE6, 0xED, 0xF3)  # body text
TEXT_GREY     = RGBColor(0x8B, 0x94, 0x9E)  # secondary text
CODE_GREEN   = RGBColor(0x7E, 0xE7, 0x87)  # code text green
TABLE_HDR_BG = RGBColor(0x16, 0x24, 0x36)  # table header dark blue
DARK_ROW1    = RGBColor(0x16, 0x1B, 0x22)  # table odd row
DARK_ROW2    = RGBColor(0x0D, 0x11, 0x17)  # table even row

# ─── Layout constants ───
SLIDE_W = Cm(33.87)
SLIDE_H = Cm(19.05)
MARGIN_L = Cm(1.8)
MARGIN_R = Cm(1.8)
MARGIN_T = Cm(1.5)
MARGIN_B = Cm(1.2)

FONT_CN = "微软雅黑"
FONT_EN = "Calibri"
FONT_CODE = "Courier New"

def create_prs():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs

def add_slide(prs):
    """Add a blank slide with dark background."""
    layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(layout)
    # set background
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = BG_DARK
    return slide

def _set_font(run, size_pt, color=TEXT_WHITE, bold=False, name=None):
    run.font.size = Pt(size_pt)
    run.font.color.rgb = color
    run.font.bold = bold
    if name:
        run.font.name = name

def add_textbox(slide, left, top, width, height, text, size_pt=14,
                color=TEXT_WHITE, bold=False, align=PP_ALIGN.LEFT,
                font_name=None, line_spacing=None):
    """Add a simple textbox with one paragraph."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    if line_spacing:
        p.line_spacing = Pt(line_spacing)
    run = p.add_run()
    run.text = text
    _set_font(run, size_pt, color, bold, font_name)
    return txBox

def add_multiline_textbox(slide, left, top, width, height, lines, size_pt=14,
                          color=TEXT_WHITE, bold=False, align=PP_ALIGN.LEFT,
                          font_name=None, line_spacing=None):
    """Add a textbox with multiple paragraphs. lines is list of str."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = align
        if line_spacing:
            p.line_spacing = Pt(line_spacing)
        run = p.add_run()
        run.text = line
        _set_font(run, size_pt, color, bold, font_name)
    return txBox

def add_rich_textbox(slide, left, top, width, height, segments, align=PP_ALIGN.LEFT, line_spacing=None):
    """Add a textbox with mixed formatting.
    segments: list of (text, size_pt, color, bold, font_name) tuples, one per paragraph.
    """
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, seg in enumerate(segments):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = align
        if line_spacing:
            p.line_spacing = Pt(line_spacing)
        # seg can be: (text, size, color, bold, font_name) or a list of runs
        if isinstance(seg, list):
            # multiple runs in one paragraph
            for run_data in seg:
                run = p.add_run()
                run.text = run_data[0]
                _set_font(run, run_data[1], run_data[2], run_data[3], run_data[4] if len(run_data) > 4 else None)
        else:
            text, sz, clr, bld = seg[0], seg[1], seg[2], seg[3]
            fn = seg[4] if len(seg) > 4 else None
            run = p.add_run()
            run.text = text
            _set_font(run, sz, clr, bld, fn)
    return txBox

def add_rect(slide, left, top, width, height, fill_color=CARD_BG,
             border_color=BORDER_COLOR, border_width=Pt(0.75)):
    """Add a rounded rectangle."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = border_width
    else:
        shape.line.fill.background()
    return shape

def add_plain_rect(slide, left, top, width, height, fill_color=CARD_BG,
                   border_color=BORDER_COLOR, border_width=Pt(0.75)):
    """Add a plain rectangle."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = border_width
    else:
        shape.line.fill.background()
    return shape

def add_card(slide, left, top, width, height, title, desc_lines,
             title_size=12, desc_size=10, title_color=TEXT_WHITE,
             desc_color=TEXT_GREY, accent_color=None, accent_top=True):
    """Add a card with optional accent bar at top."""
    # background
    card = add_rect(slide, left, top, width, height, CARD_BG, BORDER_COLOR)
    # accent bar
    if accent_color and accent_top:
        bar = add_plain_rect(slide, left, top, width, Cm(0.25), accent_color, None)
    # text
    y_off = Cm(0.4) if accent_top else Cm(0.2)
    add_textbox(slide, left + Cm(0.3), top + y_off, width - Cm(0.6), Cm(0.8),
                title, title_size, title_color, bold=True)
    if desc_lines:
        add_multiline_textbox(slide, left + Cm(0.3), top + y_off + Cm(0.8),
                              width - Cm(0.6), height - y_off - Cm(1.0),
                              desc_lines, desc_size, desc_color)
    return card

def add_arrow(slide, left, top, width, height, color=TEXT_GREY):
    """Add a right-pointing arrow shape."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RIGHT_ARROW, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_down_arrow(slide, left, top, width, height, color=TEXT_GREY):
    """Add a down-pointing arrow shape."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.DOWN_ARROW, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_flow_node(slide, left, top, width, height, text, emoji="",
                  fill_color=CARD_BG, border_color=ACCENT_BLUE,
                  text_color=TEXT_WHITE, text_size=11, bold=True):
    """Add a flow-chart node (rounded rect with centered text)."""
    shape = add_rect(slide, left, top, width, height, fill_color, border_color, Pt(1.5))
    shape.text_frame.word_wrap = True
    p = shape.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    display = f"{emoji} {text}" if emoji else text
    run = p.add_run()
    run.text = display
    _set_font(run, text_size, text_color, bold)
    shape.text_frame.paragraphs[0].space_before = Pt(0)
    shape.text_frame.paragraphs[0].space_after = Pt(0)
    return shape

def add_page_number(slide, page_num):
    """Add page number at bottom-right (skip for covers/section dividers)."""
    add_textbox(slide, SLIDE_W - Cm(2.5), SLIDE_H - Cm(1.0), Cm(2.0), Cm(0.6),
                f"P{page_num}", 10, TEXT_GREY, align=PP_ALIGN.RIGHT)

def add_slide_title(slide, title, top=None):
    """Add a standard slide title."""
    t = top if top is not None else Cm(0.8)
    add_textbox(slide, MARGIN_L, t, SLIDE_W - MARGIN_L - MARGIN_R, Cm(1.5),
                title, 26, ACCENT_BLUE, bold=True)

def add_subtitle(slide, text, top=Cm(2.2), size=14, color=TEXT_GREY):
    add_textbox(slide, MARGIN_L, top, SLIDE_W - MARGIN_L - MARGIN_R, Cm(1.0),
                text, size, color)

def add_big_number(slide, left, top, width, number_text, desc_text,
                   num_color=ACCENT_BLUE, num_size=36, desc_size=12):
    add_textbox(slide, left, top, width, Cm(1.5), number_text,
                num_size, num_color, bold=True, align=PP_ALIGN.CENTER)
    add_textbox(slide, left, top + Cm(1.6), width, Cm(1.0), desc_text,
                desc_size, TEXT_GREY, align=PP_ALIGN.CENTER)

def add_table(slide, left, top, width, col_widths, headers, rows,
              header_bg=TABLE_HDR_BG, row_bgs=(DARK_ROW1, DARK_ROW2)):
    """Add a styled table."""
    n_rows = 1 + len(rows)
    n_cols = len(headers)
    table_shape = slide.shapes.add_table(n_rows, n_cols, left, top, width, Cm(0.8) * n_rows)
    tbl = table_shape.table
    # set column widths
    for i, w in enumerate(col_widths):
        tbl.columns[i].width = w
    # header row
    for i, h in enumerate(headers):
        cell = tbl.cell(0, i)
        cell.text = h
        for p in cell.text_frame.paragraphs:
            p.alignment = PP_ALIGN.CENTER
            for run in p.runs:
                run.font.size = Pt(11)
                run.font.color.rgb = TEXT_WHITE
                run.font.bold = True
                run.font.name = FONT_EN
        cell.fill.solid()
        cell.fill.fore_color.rgb = header_bg
    # data rows
    for r_idx, row in enumerate(rows):
        bg = row_bgs[r_idx % 2]
        for c_idx, val in enumerate(row):
            cell = tbl.cell(r_idx + 1, c_idx)
            cell.text = str(val)
            for p in cell.text_frame.paragraphs:
                p.alignment = PP_ALIGN.LEFT
                for run in p.runs:
                    run.font.size = Pt(10)
                    run.font.color.rgb = TEXT_WHITE
                    run.font.name = FONT_EN
            cell.fill.solid()
            cell.fill.fore_color.rgb = bg
    # set border color for all cells
    for r in range(n_rows):
        for c in range(n_cols):
            cell = tbl.cell(r, c)
            for border_name in ['top', 'bottom', 'left', 'right']:
                pass  # python-pptx table borders are limited, skip for now
    return table_shape

def add_code_block(slide, left, top, width, height, code_lines, font_size=10):
    """Add a code-block style box."""
    bg = add_plain_rect(slide, left, top, width, height, BG_DARK, BORDER_COLOR)
    add_multiline_textbox(slide, left + Cm(0.3), top + Cm(0.2),
                          width - Cm(0.6), height - Cm(0.4),
                          code_lines, font_size, CODE_GREEN, font_name=FONT_CODE,
                          line_spacing=font_size + 4)
    return bg

def add_circle(slide, left, top, size, fill_color, text="", text_size=10, text_color=TEXT_WHITE):
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, size, size)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    if text:
        shape.text_frame.word_wrap = True
        p = shape.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = text
        _set_font(run, text_size, text_color, bold=True)
    return shape

def set_cell_border(cell, color_hex="30363D", width="0.75"):
    """Attempt to set cell border (limited support in python-pptx)."""
    try:
        from lxml import etree
        tc = cell._tc
        tcPr = tc.get_or_add_tcPr()
        for edge in ('lnL', 'lnR', 'lnT', 'lnB'):
            ln = etree.SubElement(tcPr, f'{{{cell._tc.nsmap.get("a", "")}}}{edge}')
            ln.set('w', width)
            solidFill = etree.SubElement(ln, f'{{{cell._tc.nsmap.get("a", "")}}}solidFill')
            srgbClr = etree.SubElement(solidFill, f'{{{cell._tc.nsmap.get("a", "")}}}srgbClr')
            srgbClr.set('val', color_hex)
    except Exception:
        pass
