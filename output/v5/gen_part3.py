"""
Part 3: SDD & Harness Practices (P22-P37)
"""

from pptx.util import Cm, Pt
from pptx.enum.text import PP_ALIGN
from gen_helpers import *

def add_slides(prs):
    # ══════════════════════════════════════════════════════════
    # P22 - Part 3 Section Divider
    # ══════════════════════════════════════════════════════════
    slide = add_slide(prs)
    add_textbox(slide, Cm(3), Cm(5.0), SLIDE_W - Cm(6), Cm(1.2),
                "Part 03", 20, ACCENT_BLUE, align=PP_ALIGN.CENTER)
    add_textbox(slide, Cm(3), Cm(7.0), SLIDE_W - Cm(6), Cm(2.0),
                "\u7cfb\u7edf\u5316\u63d0\u5347", 40, TEXT_WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_textbox(slide, Cm(3), Cm(10.0), SLIDE_W - Cm(6), Cm(1.2),
                "SDD \u4e0e Harness \u5b9e\u8df5", 18, TEXT_GREY, align=PP_ALIGN.CENTER)

    # ══════════════════════════════════════════════════════════
    # P23 - What is SDD
    # ══════════════════════════════════════════════════════════
    slide = add_slide(prs)
    add_slide_title(slide, "\u4ec0\u4e48\u662f SDD\uff1f")

    # Definition card
    add_card(slide, MARGIN_L, Cm(3.0), SLIDE_W - MARGIN_L - MARGIN_R, Cm(2.0),
             "SDD = Specification-Driven Development\uff08\u89c4\u8303\u9a71\u52a8\u5f00\u53d1\uff09",
             [], title_size=16, accent_color=ACCENT_BLUE, title_color=ACCENT_BLUE)

    # Comparison
    add_textbox(slide, MARGIN_L, Cm(5.8), SLIDE_W - MARGIN_L - MARGIN_R, Cm(0.8),
                "\u4f20\u7edf\u5f00\u53d1\uff1a\u89c4\u8303\u662f\u6587\u6863\uff0c\u4eba\u770b\u6587\u6863\u5199\u4ee3\u7801", 13, TEXT_WHITE)
    add_textbox(slide, MARGIN_L, Cm(6.8), SLIDE_W - MARGIN_L - MARGIN_R, Cm(0.8),
                "AI \u7f16\u7a0b\uff1a\u89c4\u8303\u662f\u201cAI \u7684\u884c\u4e3a\u8fb9\u754c\u201d\uff0c\u76f4\u63a5\u51b3\u5b9a\u8f93\u51fa\u8d28\u91cf", 13, TEXT_WHITE)

    # Three analogy lines
    analogies = [
        ("\U0001f635 \u6ca1\u6709\u89c4\u8303 = \u65b0\u5458\u5de5\u81ea\u7531\u53d1\u6325\uff08\u4e0d\u53ef\u63a7\uff09", RED_ORANGE),
        ("\U0001f60a \u6709\u89c4\u8303 = \u65b0\u5458\u5de5\u6309\u7ae0\u529e\u4e8b\uff08\u53ef\u9884\u671f\uff09", GREEN),
        ("\U0001f916 AI \u5c31\u662f\u90a3\u4e2a\u201c\u8d85\u7ea7\u65b0\u5458\u5de5\u201d", ACCENT_BLUE),
    ]
    y = Cm(8.8)
    for text, color in analogies:
        add_textbox(slide, MARGIN_L + Cm(1), y, SLIDE_W - MARGIN_L - MARGIN_R - Cm(2), Cm(0.8),
                    text, 14, color)
        y += Cm(1.4)

    add_page_number(slide, 23)

    # ══════════════════════════════════════════════════════════
    # P24 - SDD Three Layers
    # ══════════════════════════════════════════════════════════
    slide = add_slide(prs)
    add_slide_title(slide, "SDD \u4e09\u5c42\u4f53\u73b0")

    layers = [
        ("\U0001f4cb \u9879\u76ee\u7ea7\u89c4\u8303", "CLAUDE.md\u3001rules/ = \u5458\u5de5\u624b\u518c", "\u6240\u6709\u4f1a\u8bdd\u5171\u4eab\uff0c\u6bcf\u6b21\u81ea\u52a8\u52a0\u8f7d", ACCENT_BLUE),
        ("\U0001f4dd \u4efb\u52a1\u7ea7\u89c4\u8303", "Plan\u3001PRD\u3001Spec = \u4efb\u52a1\u8bf4\u660e\u4e66", "\u7279\u5b9a\u529f\u80fd\u7684\u7ea6\u675f\uff0c\u4e00\u6b21\u6027\u4f7f\u7528", YELLOW),
        ("\U0001f4ac \u4f1a\u8bdd\u7ea7\u5f15\u5bfc", "\u5bf9\u8bdd\u4e2d\u7684\u7ea6\u675f\u548c\u7ea0\u6b63 = \u65e5\u5e38\u5de5\u4f5c\u6307\u5bfc", "\u5b9e\u65f6\u4ea4\u4e92\u4e2d\u8865\u5145\u7684\u4e34\u65f6\u7ea6\u675f", GREEN),
    ]
    y = Cm(3.0)
    for title, desc, note, color in layers:
        card_h = Cm(3.0)
        add_card(slide, MARGIN_L, y, SLIDE_W - MARGIN_L - MARGIN_R, card_h,
                 f"{title}  |  {desc}", [note],
                 title_size=13, desc_size=11, accent_color=color,
                 title_color=color)
        y += card_h + Cm(0.8)

    add_page_number(slide, 24)

    # ══════════════════════════════════════════════════════════
    # P25 - Case: aip-gateway AMG Module
    # ══════════════════════════════════════════════════════════
    slide = add_slide(prs)
    add_slide_title(slide, "\u6848\u4f8b\u9879\u76ee\uff1aaip-gateway \u667a\u80fd\u4f53\u6d88\u606f\u7f51\u5173")

    add_subtitle(slide, "AMG\uff08Agent Message Gateway\uff09\u6838\u5fc3\u6a21\u5757", Cm(2.4), 13)

    # Tech stack
    add_textbox(slide, MARGIN_L, Cm(3.8), SLIDE_W - MARGIN_L - MARGIN_R, Cm(0.8),
                "Apache ShenYu 2.6.1 + Spring WebFlux + R2DBC + Redis Reactive\uff08\u54cd\u5e94\u5f0f\u5168\u94fe\u8def\uff09",
                12, TEXT_GREY)

    # Four data cards
    data = [("209\u6b21", "\u63d0\u4ea4"), ("584\u4e2a", "Java\u6587\u4ef6"), ("18\u4e2a", "\u63d2\u4ef6\u6a21\u5757"), ("4\u4e2a", "\u652f\u6301\u534f\u8bae")]
    card_w = Cm(7)
    gap = Cm(0.6)
    sx = MARGIN_L
    for i, (num, label) in enumerate(data):
        x = sx + i * (card_w + gap)
        add_rect(slide, x, Cm(5.2), card_w, Cm(2.5), CARD_BG, BORDER_COLOR)
        add_textbox(slide, x, Cm(5.4), card_w, Cm(1.2), num, 28, ACCENT_BLUE, bold=True, align=PP_ALIGN.CENTER)
        add_textbox(slide, x, Cm(6.8), card_w, Cm(0.6), label, 12, TEXT_GREY, align=PP_ALIGN.CENTER)

    # Bottom key data
    add_textbox(slide, MARGIN_L, Cm(9.5), SLIDE_W - MARGIN_L - MARGIN_R, Cm(0.8),
                "AMG\u6838\u5fc3\u6a21\u5757\uff1a59\u4e2aJava\u6587\u4ef6 | 2\u5929\u5b8c\u6210 | \u8bbe\u8ba1\u6587\u6863\u4ece v1 \u8fed\u4ee3\u5230 v1.9",
                13, ACCENT_BLUE, align=PP_ALIGN.CENTER)

    add_page_number(slide, 25)

    # ══════════════════════════════════════════════════════════
    # P26 - SDD Practice: PRD -> Spec -> Review -> Code
    # ══════════════════════════════════════════════════════════
    slide = add_slide(prs)
    add_slide_title(slide, "SDD \u5b9e\u8df5\uff1a\u4ece PRD \u5230\u4ee3\u7801\u7684\u5b8c\u6574\u6d41\u7a0b")

    # Four flow nodes
    nodes = [
        ("\U0001f4dd \u7f16\u5199PRD", ACCENT_BLUE, "\u4eba\u5b9a\u4e49\u9700\u6c42\u8fb9\u754c"),
        ("\U0001f4cb \u8f6c\u5316\u4e3aSpec 4796\u884c", ACCENT_BLUE, "AI\u8f85\u52a9\u8f6c\u5316\u4e3a\u6280\u672f\u65b9\u6848"),
        ("\U0001f50d \u591a\u8f6e\u5ba1\u9605 v1\u2192v1.9", GREEN, "\u4eba+AI\u591a\u8f6e\u5ba1\u9605\uff0c\u53d1\u73b0\u5e76\u4fee\u6b63\u95ee\u9898"),
        ("\U0001f916 Agent Team\u7f16\u7801", GREEN, "Agent Team\u6839\u636eSpec\u81ea\u52a8\u7f16\u7801"),
    ]
    node_w = Cm(6.2)
    arrow_w = Cm(1.2)
    sx = Cm(1.2)
    y_node = Cm(3.5)
    for i, (text, color, desc) in enumerate(nodes):
        x = sx + i * (node_w + arrow_w)
        add_flow_node(slide, x, y_node, node_w, Cm(2.0), text, "",
                      fill_color=CARD_BG, border_color=color,
                      text_color=color, text_size=11, bold=True)
        # description below
        add_textbox(slide, x, y_node + Cm(2.3), node_w, Cm(1.0), desc, 10, TEXT_GREY, align=PP_ALIGN.CENTER)
        if i < 3:
            add_arrow(slide, x + node_w + Cm(0.1), y_node + Cm(0.6), Cm(0.9), Cm(0.6), TEXT_GREY)

    # Bottom core message
    add_textbox(slide, Cm(2), Cm(9.5), SLIDE_W - Cm(4), Cm(1.5),
                "\u89c4\u8303\u5148\u884c \u2192 \u5ba1\u9605\u5230\u4f4d \u2192 \u4ee3\u7801\u81ea\u52a8",
                24, ACCENT_BLUE, bold=True, align=PP_ALIGN.CENTER)

    add_page_number(slide, 26)

    # ══════════════════════════════════════════════════════════
    # P27 - Spec Nine Iterations
    # ══════════════════════════════════════════════════════════
    slide = add_slide(prs)
    add_slide_title(slide, "Spec \u7684\u4e5d\u8f6e\u8fed\u4ee3\uff08v1.0 \u2192 v1.9\uff09")

    # Left: timeline (65%)
    tl_items = [
        ("v1.0-v1.2", "\u57fa\u7840\u529f\u80fd + \u6d88\u606f\u6301\u4e45\u5316 + \u80fd\u529b\u58f0\u660e", YELLOW),
        ("v1.3-v1.5", "\u5ba1\u9605\u4fee\u590d + \u67b6\u6784\u5b9a\u4f4d\u8c03\u6574", YELLOW),
        ("v1.6", "\u2605 \u91cd\u5927\u5b9a\u4f4d\u8c03\u6574 \u2014 AMG\u5b9a\u4f4d\u4e3a\u6570\u5b57\u5458\u5de5\u95e8\u6237\u6d88\u606f\u5c42", RED_ORANGE),
        ("v1.7-v1.8", "\u6280\u672f\u6808\u4fee\u6b63\uff08Spring Boot\u7248\u672c\u3001R2DBC API\uff09", YELLOW),
        ("v1.9", "\u4e09\u8f6e\u5ba1\u9605\u7efc\u5408\u4fee\u590d", GREEN),
    ]
    y = Cm(3.0)
    for ver, desc, color in tl_items:
        add_textbox(slide, MARGIN_L, y, Cm(3.5), Cm(0.7), ver, 12, color, bold=True)
        add_textbox(slide, MARGIN_L + Cm(4), y, Cm(17), Cm(0.7), desc, 12, TEXT_WHITE)
        y += Cm(1.5)

    # Right: git commit sequence (code block)
    code_lines = [
        "feat:add-design-v1",
        "feat:add-design-v2",
        "feat:add-design-v3",
        "feat:add-design-v4",
        "feat:add-design-v5",
        "feat:add-design-v6",
        "review: round-1-fix",
        "review: round-2-fix",
        "review: round-3-fix",
        "feat:update-to-v1.9",
        "feat: implement",
    ]
    add_code_block(slide, Cm(22), Cm(3.0), Cm(10), Cm(10.0), code_lines, 9)

    add_page_number(slide, 27)

    # ══════════════════════════════════════════════════════════
    # P28 - Three Review Rounds
    # ══════════════════════════════════════════════════════════
    slide = add_slide(prs)
    add_slide_title(slide, "\u4e09\u8f6e\u5ba1\u9605\uff0c\u53d1\u73b0\u4e86\u4ec0\u4e48\uff1f")

    # Table
    headers = ["\u5ba1\u9605\u7ef4\u5ea6", "\u53d1\u73b0\u7684\u95ee\u9898", "\u4fee\u6b63\u65b9\u6848"]
    rows = [
        ["\u67b6\u6784\u5b9a\u4f4d", "AMG\u88ab\u8bbe\u8ba1\u6210\u72ec\u7acb\u670d\u52a1", "\u91cd\u65b0\u5b9a\u4f4d\u4e3a\u6570\u5b57\u5458\u5de5\u95e8\u6237\u6d88\u606f\u5c42"],
        ["\u6280\u672f\u6808", "\u8bef\u7528 Spring Boot 3.x API", "\u4fee\u6b63\u4e3a 2.7.17\uff0c\u4e0e\u73b0\u6709\u7f51\u5173\u4e00\u81f4"],
        ["\u6570\u636e\u6a21\u578b", "\u4e09\u5c42\u5173\u7cfb\u8868\u8ff0\u4e0d\u6e05", "\u660e\u786e de_catalog \u2192 digital_employee \u2192 ops_channel"],
    ]
    add_table(slide, MARGIN_L, Cm(3.0), SLIDE_W - MARGIN_L - MARGIN_R,
              [Cm(5), Cm(12), Cm(13)], headers, rows)

    # Bottom core message
    add_textbox(slide, Cm(2), Cm(10.0), SLIDE_W - Cm(4), Cm(1.0),
                "\u5ba1\u9605\u4e0d\u662f\u8d70\u5f62\u5f0f\uff0c\u6bcf\u6b21\u5ba1\u9605\u90fd\u5728\u9632\u6b62\u201cAI \u9ad8\u6548\u8dd1\u504f\u201d",
                16, ACCENT_BLUE, bold=True, align=PP_ALIGN.CENTER)
    add_textbox(slide, Cm(2), Cm(11.5), SLIDE_W - Cm(4), Cm(0.8),
                "\u56de\u6263 Part 2 \u2014 \u8fd9\u5c31\u662f\u7cfb\u7edf\u5316\u5730\u9632\u6b62\u201cAI PUA \u4f60\u201d",
                12, YELLOW, align=PP_ALIGN.CENTER)

    add_page_number(slide, 28)

    # ══════════════════════════════════════════════════════════
    # P29 - SDD Results
    # ══════════════════════════════════════════════════════════
    slide = add_slide(prs)
    add_slide_title(slide, "SDD \u6210\u679c")

    # Top three big numbers
    nums = [("59\u4e2a", "Java\u6587\u4ef6"), ("2\u5929", "\u5b8c\u6210"), ("13,093\u884c", "\u4ee3\u7801/71\u4e2a\u6587\u4ef6")]
    card_w = Cm(9)
    gap = Cm(0.5)
    for i, (num, label) in enumerate(nums):
        x = MARGIN_L + i * (card_w + gap)
        add_rect(slide, x, Cm(3.0), card_w, Cm(2.5), CARD_BG, ACCENT_BLUE)
        add_textbox(slide, x, Cm(3.2), card_w, Cm(1.2), num, 28, ACCENT_BLUE, bold=True, align=PP_ALIGN.CENTER)
        add_textbox(slide, x, Cm(4.5), card_w, Cm(0.6), label, 12, TEXT_GREY, align=PP_ALIGN.CENTER)

    # Middle comparison
    lw = Cm(14)
    add_card(slide, MARGIN_L, Cm(6.5), lw, Cm(3.0),
             "aip-server", ["\u9760\u201c\u7ecf\u9a8c+\u8bbe\u8ba1\u5230\u4f4d\u201d", "\u4e2a\u4eba\u80fd\u529b\u9a71\u52a8"],
             title_size=13, accent_color=ACCENT_BLUE, desc_size=11, desc_color=TEXT_GREY)
    rx = Cm(16.5)
    rw = SLIDE_W - rx - MARGIN_R
    add_card(slide, rx, Cm(6.5), rw, Cm(3.0),
             "aip-gateway", ["\u9760\u201c\u7cfb\u7edf\u5316SDD\u6d41\u7a0b\u201d", "\u65b9\u6cd5\u8bba\u9a71\u52a8"],
             title_size=13, accent_color=GREEN, title_color=GREEN, desc_size=11, desc_color=TEXT_GREY)

    add_textbox(slide, Cm(2), Cm(10.2), SLIDE_W - Cm(4), Cm(0.7),
                "\u540c\u6837\u7684\u7ed3\u679c\uff0c\u4f46SDD\u7684\u65b9\u5f0f\u53ef\u590d\u5236\u3001\u53ef\u56e2\u961f\u5316",
                13, ACCENT_BLUE, align=PP_ALIGN.CENTER)

    # Bottom quote
    add_textbox(slide, Cm(2), Cm(11.5), SLIDE_W - Cm(4), Cm(1.2),
                "\u89c4\u8303\u5148\u884c\uff0c\u4ee3\u7801\u81ea\u52a8",
                24, ACCENT_BLUE, bold=True, align=PP_ALIGN.CENTER)

    add_page_number(slide, 29)

    # ══════════════════════════════════════════════════════════
    # P30 - What is Harness Engineering
    # ══════════════════════════════════════════════════════════
    slide = add_slide(prs)
    add_slide_title(slide, "\u4ec0\u4e48\u662f Harness Engineering\uff1f")

    # Definition
    add_textbox(slide, MARGIN_L, Cm(3.0), SLIDE_W - MARGIN_L - MARGIN_R, Cm(1.0),
                "Harness = \u4e3a AI \u642d\u5efa\u7684\u201c\u5de5\u4f5c\u811a\u624b\u67b6\u201d",
                20, ACCENT_BLUE, bold=True)

    # Four elements (2x2)
    elements = [
        ("\U0001f4cb \u89c4\u8303\u6587\u4ef6", "CLAUDE.md\u3001rules/"),
        ("\u2699\ufe0f \u6743\u9650\u914d\u7f6e", "settings.json"),
        ("\U0001f527 \u81ea\u5b9a\u4e49\u6280\u80fd", "Skills"),
        ("\U0001f504 \u5de5\u4f5c\u6d41\u914d\u7f6e", "Plan\u3001Worktree\u7b49"),
    ]
    card_w = Cm(14.5)
    card_h = Cm(2.0)
    positions = [(MARGIN_L, Cm(4.5)), (MARGIN_L + card_w + Cm(0.8), Cm(4.5)),
                 (MARGIN_L, Cm(7.0)), (MARGIN_L + card_w + Cm(0.8), Cm(7.0))]
    for i, (title, desc) in enumerate(elements):
        x, y = positions[i]
        add_card(slide, x, y, card_w, card_h, title, [desc],
                 title_size=13, desc_size=11, accent_color=ACCENT_BLUE)

    # Three pillars
    pillars = [
        ("\U0001f9e0 \u4e0a\u4e0b\u6587\u5de5\u7a0b", "\u7ed9AI\u6b63\u786e\u7684\u4fe1\u606f", ACCENT_BLUE),
        ("\U0001f3d7 \u67b6\u6784\u7ea6\u675f", "\u8ba9AI\u5728\u6846\u67b6\u5185\u5de5\u4f5c", GREEN),
        ("\u267b\ufe0f \u71b5\u7ba1\u7406", "\u9632\u6b62\u6df7\u4e71\u7d2f\u79ef", YELLOW),
    ]
    pillar_w = Cm(9.5)
    gap = Cm(0.5)
    sx = MARGIN_L
    for i, (title, desc, color) in enumerate(pillars):
        x = sx + i * (pillar_w + gap)
        add_card(slide, x, Cm(10.0), pillar_w, Cm(2.2), title, [desc],
                 title_size=12, accent_color=color, title_color=color, desc_size=10)

    add_page_number(slide, 30)

    # ══════════════════════════════════════════════════════════
    # P31 - Case: aip-portal Overview
    # ══════════════════════════════════════════════════════════
    slide = add_slide(prs)
    add_slide_title(slide, "\u6848\u4f8b\u9879\u76ee\uff1aaip-portal \u667a\u80fd\u4f53\u95e8\u6237")

    add_subtitle(slide, "\u4f01\u4e1a\u7ea7 Agent \u7ba1\u7406\u5e73\u53f0", Cm(2.4), 13)

    # Tech stack
    add_textbox(slide, MARGIN_L, Cm(3.8), SLIDE_W - MARGIN_L - MARGIN_R, Cm(0.8),
                "Spring Boot 2.7.18 + MyBatis Plus + PostgreSQL + Java 17 | Maven\u591a\u6a21\u5757 + Feign + \u591a\u6570\u636e\u6e90",
                12, TEXT_GREY)

    # Five data cards
    data = [("259\u6b21", "\u63d0\u4ea4"), ("72\u5929", "\u6d3b\u8dc3\u5f00\u53d1"), ("281\u4e2a", "Java\u6587\u4ef6"), ("57\u4e2a", "\u6587\u6863"), ("4\u4eba", "\u534f\u4f5c")]
    card_w = Cm(5.5)
    gap = Cm(0.5)
    for i, (num, label) in enumerate(data):
        x = MARGIN_L + i * (card_w + gap)
        add_rect(slide, x, Cm(5.2), card_w, Cm(2.5), CARD_BG, BORDER_COLOR)
        add_textbox(slide, x, Cm(5.4), card_w, Cm(1.2), num, 24, ACCENT_BLUE, bold=True, align=PP_ALIGN.CENTER)
        add_textbox(slide, x, Cm(6.8), card_w, Cm(0.6), label, 11, TEXT_GREY, align=PP_ALIGN.CENTER)

    # Bottom note
    add_textbox(slide, MARGIN_L, Cm(9.5), SLIDE_W - MARGIN_L - MARGIN_R, Cm(0.8),
                "\u8fd9\u4e2a\u9879\u76ee\u7ecf\u5386\u4e86 Harness \u4ece\u65e0\u5230\u6709\u7684\u5b8c\u6574\u6f14\u5316",
                14, YELLOW, align=PP_ALIGN.CENTER)

    add_page_number(slide, 31)

    # ══════════════════════════════════════════════════════════
    # P32 - Harness Evolution Overview
    # ══════════════════════════════════════════════════════════
    slide = add_slide(prs)
    add_slide_title(slide, "Rules \u7684\u6f14\u5316\uff1a\u4ece\u6df7\u4e71\u5230\u7cbe\u70bc")

    # Bar chart simulation (5 bars)
    phases = [
        ("Phase 1", 5487, Cm(5.5), YELLOW, "5487\u884c"),
        ("Phase 2", 6000, Cm(6.0), RED_ORANGE, "~6000\u884c"),
        ("Phase 2.5", 0, Cm(3.0), RED_ORANGE, "\u53d1\u73b0\u6839\u56e0\npaths: **.java"),
        ("Phase 3", 3300, Cm(3.3), YELLOW, "~3300\u884c"),
        ("Phase 4", 272, Cm(0.3), GREEN, "272\u884c \u2705"),
    ]
    bar_w = Cm(4.5)
    gap = Cm(1.2)
    max_h = Cm(7.0)
    base_y = Cm(13.0)
    sx = Cm(2.0)

    for i, (name, lines, bar_h, color, label) in enumerate(phases):
        x = sx + i * (bar_w + gap)
        # bar
        add_plain_rect(slide, x, base_y - bar_h, bar_w, bar_h, color, None)
        # label on top
        add_textbox(slide, x - Cm(0.5), base_y - bar_h - Cm(1.5), bar_w + Cm(1), Cm(1.2),
                    label, 10, color, bold=True, align=PP_ALIGN.CENTER)
        # phase name below
        add_textbox(slide, x, base_y + Cm(0.2), bar_w, Cm(0.6), name, 10, TEXT_GREY, align=PP_ALIGN.CENTER)

    add_page_number(slide, 32)

    # ══════════════════════════════════════════════════════════
    # P33 - Phase 1-2: Chaos and Expansion
    # ══════════════════════════════════════════════════════════
    slide = add_slide(prs)
    add_slide_title(slide, "Phase 1-2\uff1a\u6df7\u4e71\u4e0e\u81a8\u80c0")

    # Left: Phase 1
    lw = Cm(14)
    add_card(slide, MARGIN_L, Cm(3.0), lw, Cm(5.5),
             "Phase 1\uff1a5487\u884c\uff0c20\u4e2a\u6587\u4ef6",
             ["\u2022 \u628a\u80fd\u60f3\u5230\u7684\u90fd\u585e\u8fdb\u53bb",
              "\u274c Token\u6d88\u8017\u5de8\u5927\uff0c\u5173\u952e\u89c4\u5219\u88ab\u6df9\u6ca1"],
             title_size=14, desc_size=12, title_color=YELLOW,
             accent_color=YELLOW, desc_color=TEXT_WHITE)

    # Right: Phase 2
    rx = Cm(16.5)
    rw = SLIDE_W - rx - MARGIN_R
    add_card(slide, rx, Cm(3.0), rw, Cm(5.5),
             "Phase 2\uff1a~6000+\u884c\uff0c24\u4e2a\u6587\u4ef6",
             ["\u2022 \u8e29\u5751\u5c31\u52a0\u89c4\u5219\uff0c\u8d8a\u52a0\u8d8a\u591a",
              "\u274c Token\u7206\u70b8\uff0cAI\u7ecf\u5e38\u5ffd\u7565\u5173\u952e\u89c4\u5219"],
             title_size=14, desc_size=12, title_color=RED_ORANGE,
             accent_color=RED_ORANGE, desc_color=TEXT_WHITE)

    # Bottom warning
    add_plain_rect(slide, MARGIN_L, Cm(9.5), SLIDE_W - MARGIN_L - MARGIN_R, Cm(2.0),
                   RGBColor(0x2d, 0x1a, 0x0a), RED_ORANGE)
    add_textbox(slide, MARGIN_L + Cm(0.5), Cm(9.7), SLIDE_W - MARGIN_L - MARGIN_R - Cm(1), Cm(1.5),
                "\u26a0\ufe0f \u6076\u6027\u5faa\u73af \u2014 \u72af\u9519 \u2192 \u52a0\u89c4\u5219 \u2192 \u518d\u72af\u9519 \u2192 \u518d\u52a0\u89c4\u5219\u2026\u2026",
                16, RED_ORANGE, bold=True, align=PP_ALIGN.CENTER)

    add_page_number(slide, 33)

    # ══════════════════════════════════════════════════════════
    # P34 - Phase 3-4: Cleanup and Refinement
    # ══════════════════════════════════════════════════════════
    slide = add_slide(prs)
    add_slide_title(slide, "Phase 3-4\uff1a\u6e05\u7406\u4e0e\u7cbe\u70bc")

    # Left: Phase 3
    lw = Cm(14)
    add_card(slide, MARGIN_L, Cm(3.0), lw, Cm(5.5),
             "Phase 3\uff1a\u5220\u63892714\u884c",
             ["\u2022 \u610f\u8bc6\u5230\uff1a\u89c4\u8303\u4e0d\u662f\u8d8a\u591a\u8d8a\u597d",
              "\u2022 rules/\u662ftoken\u9884\u7b97\u7ba1\u7406\uff0c\u4e0d\u662fwiki",
              "\u2022 \u8be6\u7ec6\u6587\u6863\u79fb\u5230docs/"],
             title_size=14, desc_size=12, title_color=YELLOW,
             accent_color=YELLOW, desc_color=TEXT_WHITE)

    # Right: Phase 4
    rx = Cm(16.5)
    rw = SLIDE_W - rx - MARGIN_R
    add_card(slide, rx, Cm(3.0), rw, Cm(5.5),
             "Phase 4\uff1a\u6700\u7ec8272\u884c\uff08\u538b\u7f29\u738795%\uff09",
             ["\u2022 \u53ea\u4fdd\u7559AI\u5fc5\u987b\u77e5\u9053\u7684\u7ea6\u675f",
              "\u2022 \u8d28\u91cf\u7a33\u5b9a\uff0cToken\u53ef\u63a7"],
             title_size=14, desc_size=12, title_color=GREEN,
             accent_color=GREEN, desc_color=TEXT_WHITE)

    # Bottom timeline
    add_textbox(slide, MARGIN_L, Cm(10.5), SLIDE_W - MARGIN_L - MARGIN_R, Cm(0.8),
                "2026-04-01: \u7cbe\u7b80rules \u2192 04-09: \u89c4\u8303\u79fb\u81f3docs\uff08\u4ec49\u5929\u5b8c\u6210\u91cd\u6784\uff09",
                12, TEXT_GREY, align=PP_ALIGN.CENTER)

    add_page_number(slide, 34)

    # ══════════════════════════════════════════════════════════
    # P35 - Harness Final Form: Three-Layer Knowledge Architecture
    # ══════════════════════════════════════════════════════════
    slide = add_slide(prs)
    add_slide_title(slide, "Harness \u6700\u7ec8\u5f62\u6001")

    # Left: Before (red)
    lw = Cm(14)
    add_card(slide, MARGIN_L, Cm(3.0), lw, Cm(5.0),
             "\u274c Before\uff08Phase 1\uff09",
             ["5487\u884c\uff0c20\u4e2a\u6587\u4ef6",
              "\u6240\u6709\u89c4\u8303\u5806\u5728\u4e00\u4e2a\u76ee\u5f55",
              "\u6bcf\u6b21\u5168\u90e8\u52a0\u8f7d"],
             title_size=14, desc_size=12, title_color=RED_ORANGE,
             accent_color=RED_ORANGE, desc_color=TEXT_WHITE)

    # Right: After (green)
    rx = Cm(16.5)
    rw = SLIDE_W - rx - MARGIN_R
    add_card(slide, rx, Cm(3.0), rw, Cm(5.0),
             "\u2705 After\uff08Phase 4\uff09",
             [".claude/rules/ \u6838\u5fc3\u7ea6\u675f\uff0c\u81ea\u52a8\u52a0\u8f7d",
              "  ARCHITECTURE.md (124\u884c)",
              "  CODE-STYLE.md (14\u884c)",
              "  testing.md (134\u884c)",
              "docs/code-style/ \u8be6\u7ec6\u89c4\u8303\uff0c\u6309\u9700\u5f15\u7528",
              "  01-naming.md ... \u51798\u4e2a\u6587\u6863\uff0c838\u884c"],
             title_size=14, desc_size=10, title_color=GREEN,
             accent_color=GREEN, desc_color=CODE_GREEN)

    # Bottom design principle
    add_textbox(slide, MARGIN_L, Cm(10.0), SLIDE_W - MARGIN_L - MARGIN_R, Cm(0.7),
                "\u6838\u5fc3\u7ea6\u675f \u2192 \u81ea\u52a8\u52a0\u8f7d\uff08\u6bcf\u6b21\u4f1a\u8bdd\uff09  |  \u8be6\u7ec6\u89c4\u8303 \u2192 \u6309\u9700\u52a0\u8f7d\uff08AI\u81ea\u884c\u67e5\u9605\uff09",
                14, ACCENT_BLUE, bold=True, align=PP_ALIGN.CENTER)

    add_page_number(slide, 35)

    # ══════════════════════════════════════════════════════════
    # P36 - Harness Evolution Key Lessons
    # ══════════════════════════════════════════════════════════
    slide = add_slide(prs)
    add_slide_title(slide, "Rules \u6f14\u5316\u7684\u5173\u952e\u7ecf\u9a8c")

    # Three lessons
    lessons = [
        ("1", "rules/ \u662f token \u9884\u7b97\u7ba1\u7406\uff0c\u4e0d\u662f wiki \u2014 \u6bcf\u4e00\u884c\u90fd\u8981\u6709\u4ef7\u503c"),
        ("2", "\u89c4\u8303\u4e0d\u662f\u8d8a\u591a\u8d8a\u597d\uff0c\u800c\u662f\u8d8a\u7cbe\u70bc\u8d8a\u597d \u2014 \u5c11\u5373\u662f\u591a"),
        ("3", "\u6bcf\u6b21\u8e29\u5751\u90fd\u662f\u5b8c\u5584\u89c4\u8303\u7684\u673a\u4f1a \u2014 \u590d\u76d8\u9a71\u52a8\u6539\u8fdb"),
    ]
    y = Cm(3.5)
    for num, text in lessons:
        add_card(slide, MARGIN_L + Cm(1), y, SLIDE_W - MARGIN_L - MARGIN_R - Cm(2), Cm(2.2),
                 f"  {num}. {text}", [],
                 title_size=14, accent_color=ACCENT_BLUE, title_color=ACCENT_BLUE)
        y += Cm(2.8)

    # Bottom timeline
    add_textbox(slide, MARGIN_L, Cm(12.5), SLIDE_W - MARGIN_L - MARGIN_R, Cm(0.7),
                "04-01\u4f18\u5316 \u2192 04-01\u7cbe\u7b80 \u2192 04-02\u6574\u7406 \u2192 04-09\u89c4\u8303\u79fbdocs\uff089\u5929\u8d8b\u4e8e\u7a33\u5b9a\uff09",
                11, TEXT_GREY, align=PP_ALIGN.CENTER)

    add_page_number(slide, 36)

    # ══════════════════════════════════════════════════════════
    # P37 - Two Projects Comparison
    # ══════════════════════════════════════════════════════════
    slide = add_slide(prs)
    add_slide_title(slide, "\u4e24\u79cd\u8def\u5f84\u5bf9\u6bd4")

    # Comparison table
    headers = ["\u7ef4\u5ea6", "aip-gateway\uff08SDD\u9a71\u52a8\uff09", "aip-portal\uff08Harness\u6f14\u5316\uff09"]
    rows = [
        ["\u6838\u5fc3\u7406\u5ff5", "\u89c4\u8303\u5148\u884c\uff0c\u4ee3\u7801\u81ea\u52a8", "\u4ece\u6df7\u4e71\u4e2d\u63d0\u70bc\u79e9\u5e8f"],
        ["\u89c4\u8303\u5efa\u8bbe", "Day 1\u5c31\u5efaSpec\uff08v1\u2192v1.9\uff09", "\u8e29\u5751\u540e\u9010\u6b65\u5efa\u7acb\uff08Phase 1\u21924\uff09"],
        ["\u6df7\u4e71\u671f", "\u57fa\u672c\u6ca1\u6709", "\u7ecf\u5386\u4e86\u7ea62\u4e2a\u6708"],
        ["rules/", "\u96c6\u4e2d\u5728CLAUDE.md\uff08237\u884c\uff09", "\u5206\u5c42 rules/ + docs/\uff08272\u884c+838\u884c\uff09"],
        ["\u5f00\u53d1\u6548\u7387", "AMG\u6838\u5fc3\u6a21\u5757 2\u5929", "259\u6b21\u63d0\u4ea4\uff0c72\u5929"],
    ]
    add_table(slide, MARGIN_L, Cm(3.0), SLIDE_W - MARGIN_L - MARGIN_R,
              [Cm(5), Cm(12.5), Cm(12.5)], headers, rows)

    # Bottom insight
    add_textbox(slide, MARGIN_L, Cm(12.0), SLIDE_W - MARGIN_L - MARGIN_R, Cm(0.8),
                "\u6709\u4e86\u65b9\u6cd5\u8bba\u4e4b\u540e\uff0c\u65b0\u9879\u76ee\u53ef\u4ee5\u5c11\u8d70\u5f88\u591a\u5f2f\u8def",
                14, ACCENT_BLUE, bold=True, align=PP_ALIGN.CENTER)

    add_page_number(slide, 37)
