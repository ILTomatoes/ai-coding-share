#!/usr/bin/env python3
"""生成《Java开发中的 AI 协作编程实战》PPT"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import copy

# ── 配色方案 ──
BG_DARK = RGBColor(0x1A, 0x1A, 0x2E)       # 深蓝黑背景
BG_MID = RGBColor(0x16, 0x21, 0x3E)         # 中深背景
ACCENT = RGBColor(0x00, 0xD2, 0xFF)         # 青蓝强调色
ACCENT2 = RGBColor(0xFF, 0x6B, 0x6B)        # 红色强调
ACCENT3 = RGBColor(0x51, 0xCF, 0x66)        # 绿色强调
ACCENT4 = RGBColor(0xFF, 0xD4, 0x3B)        # 黄色强调
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT = RGBColor(0xCC, 0xCC, 0xCC)
CODE_BG = RGBColor(0x0D, 0x11, 0x17)        # 代码块背景
TABLE_HEADER_BG = RGBColor(0x00, 0x6B, 0x8A)
TABLE_ROW_BG = RGBColor(0x12, 0x1C, 0x32)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)


def set_slide_bg(slide, color=BG_DARK):
    """设置幻灯片背景色"""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, left, top, width, height):
    """添加文本框"""
    return slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))


def add_para(tf, text, size=18, color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, space_after=Pt(6), font_name='微软雅黑'):
    """在 TextFrame 中添加段落"""
    p = tf.paragraphs[-1] if len(tf.paragraphs) == 1 and tf.paragraphs[0].text == '' else tf.add_paragraph()
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    p.space_after = space_after
    return p


def add_code_block(slide, code_text, left=0.8, top=2.0, width=11.7, height=4.5):
    """添加代码块"""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = CODE_BG
    shape.line.color.rgb = RGBColor(0x33, 0x33, 0x44)
    shape.line.width = Pt(1)

    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.3)
    tf.margin_right = Inches(0.3)
    tf.margin_top = Inches(0.2)

    for i, line in enumerate(code_text.split('\n')):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(13)
        p.font.color.rgb = ACCENT
        p.font.name = 'Consolas'
        p.space_after = Pt(2)
    return shape


def add_accent_line(slide, left=0.8, top=1.35, width=3.0):
    """添加强调色横线"""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(0.05)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT
    shape.line.fill.background()
    return shape


def make_title_slide(title, subtitle=''):
    """标题页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_slide_bg(slide)

    # 标题
    tb = add_textbox(slide, 1.0, 2.5, 11.3, 1.5)
    tf = tb.text_frame
    add_para(tf, title, size=40, bold=True, color=WHITE, alignment=PP_ALIGN.LEFT)

    add_accent_line(slide, 1.0, 4.1, 4.0)

    if subtitle:
        tb2 = add_textbox(slide, 1.0, 4.4, 11.3, 1.0)
        tf2 = tb2.text_frame
        add_para(tf2, subtitle, size=20, color=LIGHT)

    return slide


def make_section_slide(part_num, title, subtitle=''):
    """部分标题页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_MID)

    # 部分编号
    tb = add_textbox(slide, 1.0, 2.0, 11.3, 0.8)
    tf = tb.text_frame
    add_para(tf, f'第{part_num}部分', size=18, color=ACCENT, bold=True)

    # 标题
    tb2 = add_textbox(slide, 1.0, 2.8, 11.3, 1.2)
    tf2 = tb2.text_frame
    add_para(tf2, title, size=36, bold=True, color=WHITE)

    add_accent_line(slide, 1.0, 4.2, 5.0)

    if subtitle:
        tb3 = add_textbox(slide, 1.0, 4.5, 11.3, 0.8)
        tf3 = tb3.text_frame
        add_para(tf3, subtitle, size=18, color=LIGHT)

    return slide


def make_content_slide(title, bullets, accent_line=True):
    """内容页：标题 + 列表"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    # 标题
    tb = add_textbox(slide, 0.8, 0.5, 11.7, 0.8)
    tf = tb.text_frame
    add_para(tf, title, size=28, bold=True, color=WHITE)

    if accent_line:
        add_accent_line(slide, 0.8, 1.35, 3.0)

    # 内容
    tb2 = add_textbox(slide, 0.8, 1.6, 11.7, 5.5)
    tf2 = tb2.text_frame
    tf2.word_wrap = True

    for bullet in bullets:
        if isinstance(bullet, tuple):
            # (text, level, color)
            text, level, color = bullet
            prefix = '  ' * level + ('• ' if level == 0 else '◦ ' if level == 1 else '- ')
            add_para(tf2, prefix + text, size=18 if level == 0 else 16, color=color or WHITE)
        elif isinstance(bullet, str):
            if bullet.startswith('##'):
                add_para(tf2, bullet[2:].strip(), size=20, bold=True, color=ACCENT, space_after=Pt(4))
            elif bullet.startswith('>>'):
                add_para(tf2, bullet[2:].strip(), size=18, bold=True, color=ACCENT4, space_after=Pt(4))
            elif bullet.startswith('!!'):
                add_para(tf2, bullet[2:].strip(), size=18, bold=True, color=ACCENT2, space_after=Pt(4))
            else:
                add_para(tf2, '• ' + bullet, size=18, color=WHITE)

    return slide


def make_two_col_slide(title, left_title, left_items, right_title, right_items):
    """两栏对比页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    # 标题
    tb = add_textbox(slide, 0.8, 0.5, 11.7, 0.8)
    tf = tb.text_frame
    add_para(tf, title, size=28, bold=True, color=WHITE)
    add_accent_line(slide, 0.8, 1.35, 3.0)

    # 左栏标题
    col_width = 5.5
    tb_lt = add_textbox(slide, 0.8, 1.7, col_width, 0.5)
    tf_lt = tb_lt.text_frame
    add_para(tf_lt, left_title, size=20, bold=True, color=ACCENT)

    # 左栏内容
    tb_l = add_textbox(slide, 0.8, 2.3, col_width, 4.5)
    tf_l = tb_l.text_frame
    tf_l.word_wrap = True
    for item in left_items:
        add_para(tf_l, '• ' + item, size=16, color=WHITE)

    # 分隔线
    sep = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(6.6), Inches(1.7), Inches(0.02), Inches(5.0)
    )
    sep.fill.solid()
    sep.fill.fore_color.rgb = RGBColor(0x44, 0x44, 0x55)
    sep.line.fill.background()

    # 右栏标题
    tb_rt = add_textbox(slide, 7.0, 1.7, col_width, 0.5)
    tf_rt = tb_rt.text_frame
    add_para(tf_rt, right_title, size=20, bold=True, color=ACCENT)

    # 右栏内容
    tb_r = add_textbox(slide, 7.0, 2.3, col_width, 4.5)
    tf_r = tb_r.text_frame
    tf_r.word_wrap = True
    for item in right_items:
        add_para(tf_r, '• ' + item, size=16, color=WHITE)

    return slide


def make_quote_slide(quote, source=''):
    """金句页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_MID)

    tb = add_textbox(slide, 1.5, 2.5, 10.3, 2.0)
    tf = tb.text_frame
    tf.word_wrap = True
    add_para(tf, f'"{quote}"', size=28, bold=True, color=ACCENT, alignment=PP_ALIGN.CENTER)

    if source:
        tb2 = add_textbox(slide, 1.5, 4.8, 10.3, 0.5)
        tf2 = tb2.text_frame
        add_para(tf2, f'—— {source}', size=16, color=LIGHT, alignment=PP_ALIGN.CENTER)

    return slide


def make_code_slide(title, code_text):
    """代码展示页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    tb = add_textbox(slide, 0.8, 0.5, 11.7, 0.8)
    tf = tb.text_frame
    add_para(tf, title, size=28, bold=True, color=WHITE)
    add_accent_line(slide, 0.8, 1.35, 3.0)

    add_code_block(slide, code_text, top=1.7, height=5.3)

    return slide


def make_table_slide(title, headers, rows):
    """表格页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    tb = add_textbox(slide, 0.8, 0.5, 11.7, 0.8)
    tf = tb.text_frame
    add_para(tf, title, size=28, bold=True, color=WHITE)
    add_accent_line(slide, 0.8, 1.35, 3.0)

    num_cols = len(headers)
    num_rows = len(rows) + 1
    col_width = Inches(11.7 / num_cols)

    table_shape = slide.shapes.add_table(
        num_rows, num_cols,
        Inches(0.8), Inches(1.7),
        Inches(11.7), Inches(0.5 * num_rows)
    )
    table = table_shape.table

    # 表头
    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = h
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(16)
            p.font.bold = True
            p.font.color.rgb = WHITE
            p.font.name = '微软雅黑'
        cell.fill.solid()
        cell.fill.fore_color.rgb = TABLE_HEADER_BG

    # 数据行
    for r_idx, row in enumerate(rows):
        for c_idx, val in enumerate(row):
            cell = table.cell(r_idx + 1, c_idx)
            cell.text = val
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(14)
                p.font.color.rgb = WHITE
                p.font.name = '微软雅黑'
            cell.fill.solid()
            cell.fill.fore_color.rgb = TABLE_ROW_BG if r_idx % 2 == 0 else BG_DARK

    return slide


def make_timeline_slide(title, phases):
    """时间线页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    tb = add_textbox(slide, 0.8, 0.5, 11.7, 0.8)
    tf = tb.text_frame
    add_para(tf, title, size=28, bold=True, color=WHITE)
    add_accent_line(slide, 0.8, 1.35, 3.0)

    # 时间线横线
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(1.0), Inches(3.5), Inches(11.3), Inches(0.04)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT
    line.line.fill.background()

    # 各阶段
    num = len(phases)
    spacing = 11.3 / (num - 1) if num > 1 else 0

    for i, (phase_title, phase_desc, phase_detail) in enumerate(phases):
        x = 1.0 + i * spacing

        # 圆点
        dot = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(x - 0.12), Inches(3.38), Inches(0.24), Inches(0.24)
        )
        dot.fill.solid()
        dot.fill.fore_color.rgb = ACCENT
        dot.line.fill.background()

        # 上方标签
        tb_up = add_textbox(slide, x - 0.8, 2.3, 1.8, 0.8)
        tf_up = tb_up.text_frame
        add_para(tf_up, phase_title, size=14, bold=True, color=ACCENT, alignment=PP_ALIGN.CENTER)
        add_para(tf_up, phase_desc, size=12, color=LIGHT, alignment=PP_ALIGN.CENTER)

        # 下方详情
        tb_dn = add_textbox(slide, x - 1.0, 3.9, 2.2, 2.5)
        tf_dn = tb_dn.text_frame
        tf_dn.word_wrap = True
        add_para(tf_dn, phase_detail, size=13, color=LIGHT, alignment=PP_ALIGN.CENTER)

    return slide


# ════════════════════════════════════════════════════════════
# 开始生成 PPT
# ════════════════════════════════════════════════════════════

# ── PPT 1: 标题页 ──
make_title_slide(
    'Java 开发中的 AI 协作编程实战',
    '如何为 AI 搭建高效工作的环境 —— 来自真实项目的实践经验'
)

# ── PPT 2: 互动问题 ──
make_content_slide('开场小调查', [
    '觉得自己用 AI 写代码效率很高的同事，请举手？',
    '遇到过 AI 生成的代码需要大量返工、反复修改的，请举手？',
    '',
    '>>真实数据：项目中 AI 配置从 819 行单文件 → 精简删除 6000+ 行',
    '>>这个数字背后的故事，就是今天分享的核心'
])

# ── PPT 3: 误区一和误区二 ──
make_two_col_slide(
    'AI 编程的三个常见误区',
    '❌ 误区一："AI 什么都会"',
    [
        'AI 只能看到你给它的上下文',
        '就像新同事没看架构文档就写代码',
        '信息不够时 AI 会自信地"发挥"',
        '而 AI 的方向可能完全不对'
    ],
    '❌ 误区二："告诉 AI 就够了"',
    [
        '缺少约束 → AI 自信地偏离方向',
        '真实教训：让 AI 生成大量代码',
        '没有充分审查直接采纳',
        '对接联调时大量返工、调试',
        '比手写代码更危险——有"已完成"的错觉'
    ]
)

# ── PPT 4: 误区三 + 转折 ──
make_content_slide('误区三 & 根本原因', [
    '❌ 误区三："一次就能做好"',
    ('不可能预测 AI 会犯什么错', 1, LIGHT),
    ('试图一次性把规则写全，往往过度约束', 1, LIGHT),
    ('AI 的输出变得保守、僵化', 1, LIGHT),
    '',
    '##三条误区的共同根本原因：',
    '!!我们没有为 AI 搭建合适的工作环境',
])

# ── PPT 5: Harness 比喻 ──
slide = make_quote_slide(
    'AI 是千里马，Harness 是缰绳。\n缰绳不增加马的力量，但让力量用在对的方向。',
    'Harness Engineering'
)

# ── PPT 6: 三层架构全景图 ──
make_code_slide(
    '三层知识架构——给 AI 传递知识的最佳方式',
    '''┌──────────────────────────────────────────┐
│  全局层 (~/.claude/CLAUDE.md)              │
│  → "我是谁，我怎么工作"                    │
│  → 所有项目共享，长期稳定                   │
├──────────────────────────────────────────┤
│  项目层 (CLAUDE.md + .claude/rules/)       │
│  → "这个项目怎么做"                        │
│  → 每次对话全部加载，每行都在消耗 token     │
├──────────────────────────────────────────┤
│  会话层 (memory/)                          │
│  → "上次踩了什么坑"                        │
│  → 按需检索，不自动全部加载                 │
└──────────────────────────────────────────┘'''
)

# ── PPT 7: 全局层详解 ──
make_content_slide('全局层：跨项目的个人偏好', [
    '位置：~/.claude/CLAUDE.md',
    '只放两样东西：',
    ('所有项目通用的规则', 1, LIGHT),
    ('长期稳定的工作习惯', 1, LIGHT),
    '',
    '##实际内容示例：',
    'Git 提交规范：中文 type: 格式，禁止加 Co-Authored-By 署名',
    'Shell 环境规范：Windows Git Bash 中用 /dev/null 而不是 > nul',
    '',
    '>>原则：这个文件的内容应该几乎不变',
])

# ── PPT 8: 项目层详解 ──
make_content_slide('项目层：最重要的层级', [
    '位置：项目根目录 CLAUDE.md + .claude/rules/',
    '',
    '!!关键特点：每次对话全部加载，每行都在消耗 token',
    '',
    '##核心原则：',
    '>>"rules/ 是给 AI 的参考，不是给人看的 wiki"',
    '',
    '常见误区：把设计文档、执行计划、DB Schema 都塞进 rules/',
    '后果：AI 每次阅读几千行无关内容，效率下降，容易混乱',
    '',
    '正确做法：rules/ 只放 AI 工作时需要参考的内容',
])

# ── PPT 9: 会话层详解 ──
make_content_slide('会话层：跨会话的经验积累', [
    '位置：memory/ 目录',
    '',
    '##与项目层的最大区别：',
    '不自动加载，按需检索',
    'AI 根据当前任务上下文自动查找相关记忆',
    '',
    '##保存的内容：',
    '复盘经验（feedback 类型）——踩过的坑',
    '外部资源索引（reference 类型）',
    '',
    '>>优势：自由积累经验，不担心占用每次对话的 token 预算',
])

# ── PPT 10: 设计原则总结 ──
make_table_slide(
    '三层架构的设计原则',
    ['原则', '说明'],
    [
        ['信息不重复', '同一条规则只存在于一个层级，其他层级引用'],
        ['层级边界清晰', '全局放"我是谁"，项目放"怎么做"，会话放"踩了什么坑"'],
        ['Token 有意识', '项目层每次全部加载须有存在价值；会话层按需检索可更自由'],
    ]
)

# ── PPT 11: Phase 1-2 混乱期 ──
make_two_col_slide(
    'rules/ 演变：Phase 1-2（混乱期）',
    'Phase 1：初始搭建（3.25-26）',
    [
        'CODE-STYLE.md 单文件 819 行',
        '命名、注解、数据库、API 设计全挤在一起',
        '没有测试规范',
        '"把规范写全了 AI 就能写好代码"',
    ],
    'Phase 2：功能驱动扩张（3.27-31）',
    [
        '每做新功能就往 rules/ 加文档',
        '设计文档 1706 行（一个功能）',
        '执行计划动辄 1000+ 行',
        'completed/ 目录堆满旧计划',
        '每次对话加载所有文件包括已完成的',
        'AI 检索效率直线下降',
    ]
)

# ── PPT 12: Phase 3 关键转折 ──
make_content_slide('Phase 3：关键转折（4.1 精简重构）', [
    '##确立核心原则：',
    '!!"rules/ 是给 AI 的参考，不是给人看的 wiki"',
    '',
    '##拆分（+1175 行，结构更清晰）：',
    ('CODE-STYLE.md 819行 → 8 个专题文件 + 1 个索引文件', 1, LIGHT),
    ('每个文件聚焦一个主题：命名、注解、数据库、API 等', 1, LIGHT),
    '',
    '##大清理（-6278 行）：',
    ('删除已完成的执行计划', 1, LIGHT),
    ('设计文档从 rules/ 移除 → 放到 docs/', 1, LIGHT),
    ('删除可再生的数据库 schema', 1, LIGHT),
])

# ── PPT 13: Phase 4 稳定维护 ──
make_code_slide(
    'Phase 4：稳定维护期——最终结构',
    '''rules/
├── ARCHITECTURE.md     # 架构全景（稳定）
├── CODE-STYLE.md       # 规范索引（指向 code-style/）
├── code-style/         # 8个专题文件
│   ├── 01-naming.md
│   ├── 02-annotations.md
│   ├── 03-database.md
│   ├── 04-api-design.md
│   ├── 05-comments-logging.md
│   ├── 06-vo-di.md
│   ├── 07-date-constants.md
│   └── 08-pagination.md
└── testing.md          # 测试规范

编码规范文件几乎不变
仅架构变更时更新 ARCHITECTURE.md'''
)

# ── PPT 14: 时间线与总结 ──
make_timeline_slide(
    '演变时间线：好的结构是"长"出来的',
    [
        ('Day 1', '819 行\n单文件大杂烩', '大胆堆内容\n不知道该放什么'),
        ('Week 1', '混乱膨胀\n+设计文档+计划', '每做功能就加文档\ntoken 大量浪费'),
        ('Week 2', '-6278 行\n大精简重构', '确立核心原则\n拆分+清理'),
        ('Now', '稳定维护\n结构清晰', '几乎不变\n仅架构变更时更新'),
    ]
)

# ── PPT 15: 记忆系统 ──
make_content_slide('记忆系统——把踩过的坑变成制度', [
    '##复盘工作流（/retrospect）：',
    '1. 做对了什么？',
    '2. 做错了什么？',
    '3. 如果重新做会怎么做？',
    '',
    '##过滤标准：',
    ('✗ 一次性的技术细节（已在代码中）', 1, ACCENT2),
    ('✗ 可从代码推导的信息', 1, ACCENT2),
    ('✓ 踩坑原因和预防方法', 1, ACCENT3),
    ('✓ 工作流改进方向', 1, ACCENT3),
    '',
    '>>统一格式：规则 + Why（为什么）+ How to apply（何时用）',
])

# ── PPT 16: 工作流对比 ──
make_two_col_slide(
    '工作流优化：从"直接开干"到"先规划再做"',
    '❌ 早期：直接开干',
    [
        '告诉 AI 要做什么，直接写代码',
        'Git 记录：连续的 fix、fix、fix',
        '边做边改，方向不明',
        '返工频繁',
    ],
    '✅ 后来：先规划再做',
    [
        'Plan → 确认 → 实现 → 验证',
        'Git 记录：清晰的阶段划分',
        '探索代码结构后再动手',
        '方向正确，返工大幅减少',
    ]
)

# ── PPT 17: 推荐工作流 ──
make_content_slide('推荐工作流（四步法）', [
    '##Step 1: Plan 阶段',
    ('让 AI 探索现有代码结构，识别可复用组件和模式', 1, LIGHT),
    '',
    '##Step 2: 确认阶段',
    ('对关键不确定信息先确认：版本号、字段来源、依赖版本等', 1, LIGHT),
    ('这些信息猜错了，后面全是返工', 1, ACCENT2),
    '',
    '##Step 3: 实现阶段',
    ('确认方案后才进入编码', 1, LIGHT),
    '',
    '##Step 4: 验证阶段',
    ('实现后运行编译验证', 1, LIGHT),
])

# ── PPT 18: PRD 质量决定成败 ──
make_quote_slide(
    '原型/PRD 质量高，AI 编程事半功倍。\n反之，模糊不清则飞速跑偏，没法验收、无法判断对错。',
    '实践经验'
)

# ── PPT 19: 小批量推进 + Hook ──
make_content_slide('小批量推进 + Hook 自动化', [
    '##血泪教训：',
    '缺少审查，AI 生成大量代码未验证 → 对接时大量返工',
    '比从头写还花时间',
    '',
    '##两个关键原则：',
    '小批量推进验收，逐渐加快节奏',
    '不要妄图一次性把刀磨好，遇到问题再磨刀',
    '',
    '##Hook 自动化：',
    'PostToolUse 钩子：Java 文件修改后自动执行 mvn compile',
    'PreCommit 钩子：提交前自动验证测试编译',
    '>>把验证变成流程的一部分，不会遗漏',
])

# ── PPT 20: 踩坑沉淀表格 ──
make_table_slide(
    '从错误中提炼规则——踩坑沉淀',
    ['踩坑场景', '沉淀结果', '所在层级'],
    [
        ['Git Bash 中 > nul 创建垃圾文件', 'Shell 重定向规范', '全局 CLAUDE.md'],
        ['Claude 默认加 Co-Authored-By', 'Git 提交规范', '全局 CLAUDE.md'],
        ['SQL 迁移版本号命名错误', '版本号确认规则', 'memory feedback'],
        ['权限控制理解偏差', '架构理解修正', 'feedback + ARCHITECTURE.md'],
    ]
)

# ── PPT 21: 生长循环 ──
make_code_slide(
    'Harness 生长循环',
    '''    ┌─────────────────────────────────────────┐
    │                                         │
    │   ① 让 AI 做任务（尽量少约束）          │
    │          ↓                              │
    │   ② 观察 AI 的输出                      │
    │          ↓                              │
    │   ③ 发现不符合预期的地方                 │
    │          ↓                              │
    │   ④ 分析根因（缺上下文？缺规则？）       │
    │          ↓                              │
    │   ⑤ 添加最小约束（只解决当前问题）       │
    │          ↓                              │
    │   ⑥ 验证效果                            │
    │          ↓                              │
    │   回到 ①，循环往复                      │
    │                                         │
    └─────────────────────────────────────────┘

关键：每次只添加解决当前问题的最小约束'''
)

# ── PPT 22: 错误分类框架 ──
make_table_slide(
    '错误分类框架——不是所有错误都需要加规则',
    ['错误类型', '频率', '应对策略'],
    [
        ['上下文缺失', '几乎每次', '在 CLAUDE.md 写明技术栈'],
        ['规则不明', '前几次', '犯过一次就加规则'],
        ['架构偏离', '项目变大后', '分层文档 + 自动检查'],
        ['偶发错误', '偶尔', '直接修复，不要加规则'],
    ]
)

# ── PPT 23: 黄金法则 ──
make_quote_slide(
    '如果一个错误只发生一次，且不太可能再发生，\n直接修复它，不要加规则。\n\n过多规则带来的阅读成本，比偶发错误本身造成的问题更大。',
    '黄金法则'
)

# ── PPT 24: 演示 A ──
make_content_slide('实战演示 A：从零搭建 Spring Boot 项目的 Harness', [
    '>>视频 1：没有 CLAUDE.md 的裸项目',
    ('让 AI 添加接口 → 业务逻辑在 Controller、无 DTO、无异常处理', 1, LIGHT),
    '',
    '>>视频 2：3 行 CLAUDE.md 的效果',
    ('项目名 + 技术栈 + 编码风格 → 基本分层结构自动出来', 1, LIGHT),
    '',
    '>>视频 3：逐步添加约束的生长过程',
    ('第一次犯错 → 分层架构规则（3行→9行）', 1, LIGHT),
    ('第二次犯错 → DTO 分离规则（9行→15行）', 1, LIGHT),
    ('第三次犯错 → 统一异常处理（15行→30行）', 1, LIGHT),
    '',
    '>>每条规则的来源都是一个具体的"事故"',
])

# ── PPT 25: 演示 B ──
make_content_slide('实战演示 B：多任务并行与闭环工作流', [
    '>>视频 4：worktree 多任务并行',
    ('每个任务一个独立工作目录', 1, LIGHT),
    ('各自有独立的 AI 会话，互不干扰', 1, LIGHT),
    ('两个 AI 实例同时工作', 1, LIGHT),
    '',
    '>>视频 5：设计-开发-测试闭环',
    ('设计阶段：Plan 模式让 AI 先探索、规划', 1, LIGHT),
    ('开发阶段：Claude Code 按方案编码', 1, LIGHT),
    ('测试阶段：Karate 接口测试验证', 1, LIGHT),
    ('每个环节有明确产出物，避免方向跑偏', 1, ACCENT3),
])

# ── PPT 26: 演示 C ──
make_content_slide('实战演示 C：复盘与记忆积累', [
    '>>视频 6：/retrospect 复盘演示',
    '',
    '##复盘流程：',
    '完成任务 → 触发复盘 → 三问反思 → 过滤保存',
    '',
    '##记忆的格式：',
    '规则描述 + Why（为什么）+ How to apply（何时用）',
    '',
    '##记忆的复用：',
    '新会话中 AI 自动检索相关记忆',
    '同样的错误不会犯第二次',
    '',
    '>>"经验制度化"——把"下次注意"变成可执行的知识记录',
])

# ── PPT 27: 演示 D（备用）──
make_content_slide('延伸场景：AI 解读开源项目', [
    '用 AI 分析 Spring Boot 开源项目',
    'AI 阅读源码 → 理解架构 → 输出技术文档',
    '',
    '人工做：1-2 天',
    'AI 做：30 分钟',
    '',
    '##适用场景：',
    '快速了解陌生系统的架构',
    '学习新技术栈',
    '系统调研和技术选型',
])

# ── PPT 28: 完整工作流链 ──
make_code_slide(
    '进阶：Skills 完整工作流链',
    '''需求提出
    ↓
brainstorming          ← 探索需求和设计方案
    ↓
writing-plans          ← 编写实施计划
    ↓
executing-plans        ← 执行计划（带审查检查点）
    ↓
verification           ← 完成前验证（运行测试、编译）
    ↓
code-review            ← 代码审查
    ↓
retrospect             ← 复盘，保存经验

每个环节有明确的检查标准，避免跳步导致的返工'''
)

# ── PPT 29: 工具组合 ──
make_content_slide('配套工具', [
    '##Claude Code',
    '核心编程工具，命令行操作，与项目深度集成',
    '',
    '##playwright-mcp',
    '自动化浏览器操作',
    '系统调研：AI 自动打开系统、操作页面、输出文档',
    '编写用户手册：AI 直接看页面截图描述功能',
    '',
    '##cc-switch',
    '模型切换工具，方便在不同模型之间切换',
])

# ── PPT 30: 模型选择策略 ──
make_content_slide('模型选择策略', [
    '##核心策略：每个模型用在它最擅长的环节',
    '',
    '设计阶段 → 用强模型（Opus）',
    ('架构设计、需求分析的理解力最好', 1, LIGHT),
    '',
    '编码实现 → 用性价比更高的模型',
    ('日常编码不需要最贵的模型', 1, LIGHT),
    '',
    '>>"奢侈但有效"：',
    '如果和 AI 交互了很多次都得不到好方案',
    '直接换用 Opus 反而更省时间',
])

# ── PPT 31: 应用场景矩阵 ──
make_table_slide(
    'AI 协作编程的更多场景',
    ['场景', '说明', '价值'],
    [
        ['解读开源项目', 'AI 阅读源码，输出架构分析文档', '快速理解陌生系统'],
        ['系统调研', 'AI + playwright 深入解读在线系统', '自动化信息收集'],
        ['自学技术', 'AI 讲解新技术 + 代码示例', '加速学习曲线'],
        ['多任务并行', 'worktree + agent teams', '复杂任务拆分并行'],
    ]
)

# ── PPT 32: 多任务并行架构 ──
make_content_slide('多任务并行架构', [
    '##worktree：独立工作目录',
    ('每个任务一个 worktree，互不干扰', 1, LIGHT),
    ('适合多个需求同时推进的场景', 1, LIGHT),
    '',
    '##agent teams：多 Agent 协作',
    ('Explore Agent：只读探索代码', 1, LIGHT),
    ('Plan Agent：编写实施计划', 1, LIGHT),
    ('实现 Agent：编码 + 测试', 1, LIGHT),
    '',
    '>>适用场景：大型功能开发、多模块重构、多方资料调研',
])

# ── PPT 33: 六大经验回顾 ──
make_content_slide('六大关键经验回顾', [
    ('1. rules/ 是 token 预算管理 —— 每个文件都要有存在价值', 0, WHITE),
    ('2. 记忆系统是经验制度化 —— 不复盘等于白做，不保存等于白复盘', 0, WHITE),
    ('3. 三层架构各司其职 —— 全局放偏好，项目放规范，会话放经验', 0, WHITE),
    ('4. 先有再优化 —— 早期大胆堆内容，有积累后再精简', 0, WHITE),
    ('5. 从错误中提炼规则 —— 每次踩坑都是完善知识体系的机会', 0, WHITE),
    ('6. 文档有生命周期 —— 创建→迭代→稳定→过时（更新或删除）', 0, WHITE),
])

# ── PPT 34: 行动清单 ──
make_content_slide('立即可执行的行动清单', [
    '##今天就能做：',
    ('在你的 Java 项目中创建 CLAUDE.md（只要 3-5 行）', 1, ACCENT3),
    ('安装 Claude Code，让它帮你做一个小任务', 1, ACCENT3),
    '',
    '##本周可以做：',
    ('观察 AI 最常犯的 3 个错误', 1, ACCENT4),
    ('为这 3 个错误添加最小约束', 1, ACCENT4),
    ('尝试一次复盘，把踩过的坑记下来', 1, ACCENT4),
    '',
    '##本月可以做：',
    ('完善 rules/ 目录结构，按专题拆分', 1, ACCENT),
    ('配置 Hook 自动化（编译验证等）', 1, ACCENT),
    ('尝试 worktree 多任务并行', 1, ACCENT),
])

# ── PPT 35: 结束页 ──
make_quote_slide(
    'AI 协作编程的核心不是让 AI 变强，\n而是让 AI 在正确的轨道上工作。\n你构建的 Harness，就是这条轨道。',
    ''
)

# 添加底部文字到结束页
last_slide = prs.slides[-1]
tb_end = add_textbox(last_slide, 1.5, 6.0, 10.3, 0.8)
tf_end = tb_end.text_frame
add_para(tf_end, '感谢聆听  |  Q&A 交流时间', size=16, color=LIGHT, alignment=PP_ALIGN.CENTER)


# ── 保存文件 ──
output_path = '/root/ai-coding-share/output/AI协作编程实战.pptx'
prs.save(output_path)
print(f'PPT 已生成: {output_path}')
print(f'共 {len(prs.slides)} 页')
