"""
Part 0: Cover + Opening (P1-P2)
"""

from pptx.util import Cm, Pt
from pptx.enum.text import PP_ALIGN
from gen_helpers import *

def add_slides(prs):
    # ══════════════════════════════════════════════════════════
    # P1 - Cover Page
    # ══════════════════════════════════════════════════════════
    slide = add_slide(prs)

    # Main title - centered, 44pt, bright blue
    add_textbox(slide, Cm(3), Cm(5.0), SLIDE_W - Cm(6), Cm(3.0),
                "Java \u5f00\u53d1\u4e2d\u7684 AI \u534f\u4f5c\u7f16\u7a0b\u5b9e\u8df5",
                44, ACCENT_BLUE, bold=True, align=PP_ALIGN.CENTER)

    # Subtitle - 20pt, grey
    add_textbox(slide, Cm(3), Cm(8.5), SLIDE_W - Cm(6), Cm(1.5),
                "\u4ece\u771f\u5b9e\u9879\u76ee\u8bf4\u8d77",
                20, TEXT_GREY, align=PP_ALIGN.CENTER)

    # Bottom: speaker / date
    add_textbox(slide, Cm(3), Cm(14.0), SLIDE_W - Cm(6), Cm(1.0),
                "\u6f14\u8bb2\u4eba / 2026\u5e744\u6708",
                14, TEXT_GREY, align=PP_ALIGN.CENTER)

    # Bottom-right small text
    add_textbox(slide, SLIDE_W - Cm(6), SLIDE_H - Cm(1.5), Cm(5), Cm(0.8),
                "\u5185\u90e8\u4ea4\u6d41\u5206\u4eab",
                10, TEXT_GREY, align=PP_ALIGN.RIGHT)

    # ══════════════════════════════════════════════════════════
    # P2 - Opening: About this sharing
    # ══════════════════════════════════════════════════════════
    slide = add_slide(prs)
    add_slide_title(slide, "\u5173\u4e8e\u8fd9\u6b21\u5206\u4eab")

    # Three lines with emoji
    items = [
        "\U0001f3d7  \u4e09\u4e2a\u771f\u5b9e Java \u9879\u76ee\u7684 AI \u534f\u4f5c\u7f16\u7a0b\u7ecf\u5386",
        "\U0001f3a2  \u6709\u6210\u529f\u7684\u559c\u60a6\uff0c\u4e5f\u6709\u7ffb\u8f66\u7684\u6559\u8bad",
        "\U0001f3af  \u76ee\u6807\uff1a\u542c\u5b8c\u4e4b\u540e\uff0c\u4f60\u80fd\u5e26\u8d70\u4e00\u4e9b\u53ef\u4ee5\u9a6c\u4e0a\u7528\u7684\u4e1c\u897f"
    ]
    y = Cm(3.5)
    for item in items:
        add_textbox(slide, MARGIN_L, y, SLIDE_W - MARGIN_L - MARGIN_R, Cm(1.5),
                    item, 18, TEXT_WHITE)
        y += Cm(2.2)

    # Bottom yellow note
    add_textbox(slide, MARGIN_L, SLIDE_H - Cm(2.5), SLIDE_W - MARGIN_L - MARGIN_R, Cm(0.8),
                "\u9884\u8ba1 90 \u5206\u949f | \u9f13\u52b1\u968f\u65f6\u63d0\u95ee",
                14, YELLOW, align=PP_ALIGN.CENTER)

    add_page_number(slide, 2)
