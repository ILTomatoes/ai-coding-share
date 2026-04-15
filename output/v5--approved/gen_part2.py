"""
Part 2: Reflections (P15-P21)
"""

from pptx.util import Cm, Pt
from pptx.enum.text import PP_ALIGN
from gen_helpers import *

def add_slides(prs):
    # ══════════════════════════════════════════════════════════
    # P15 - Part 2 Section Divider
    # ══════════════════════════════════════════════════════════
    slide = add_slide(prs)
    add_textbox(slide, Cm(3), Cm(5.0), SLIDE_W - Cm(6), Cm(1.2),
                "Part 02", 20, ACCENT_BLUE, align=PP_ALIGN.CENTER)
    add_textbox(slide, Cm(3), Cm(7.0), SLIDE_W - Cm(6), Cm(2.0),
                "\u53cd\u601d", 40, TEXT_WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_textbox(slide, Cm(3), Cm(10.0), SLIDE_W - Cm(6), Cm(1.2),
                "\u4ece\u6848\u4f8b\u5230\u6d1e\u5bdf", 18, TEXT_GREY, align=PP_ALIGN.CENTER)

    # ══════════════════════════════════════════════════════════
    # P16 - What we saw from cases
    # ══════════════════════════════════════════════════════════
    slide = add_slide(prs)
    add_slide_title(slide, "\u4ece\u6848\u4f8b\u4e2d\u770b\u5230\u4e86\u4ec0\u4e48")

    # Top two comparison lines
    add_textbox(slide, MARGIN_L, Cm(3.0), SLIDE_W - MARGIN_L - MARGIN_R, Cm(0.7),
                "aip-server \u6210\u529f\uff1a\u4eba\u7ed9\u4e86\u6e05\u6670\u7684\u8bbe\u8ba1\u548c\u8fb9\u754c", 13, GREEN)
    add_textbox(slide, MARGIN_L, Cm(4.0), SLIDE_W - MARGIN_L - MARGIN_R, Cm(0.7),
                "A2A \u5931\u8d25\uff1a\u4eba\u6ca1\u6709\u7ed9\u65b9\u5411\uff0cAI \u9ad8\u6548\u8dd1\u504f", 13, RED_ORANGE)

    # Center big text
    add_textbox(slide, Cm(2), Cm(6.5), SLIDE_W - Cm(4), Cm(2.5),
                "AI \u7684\u6548\u7387\u662f\u786e\u5b9a\u7684\uff0c\u4e0d\u786e\u5b9a\u7684\u662f\u4eba\u7684\u5f15\u5bfc",
                26, ACCENT_BLUE, bold=True, align=PP_ALIGN.CENTER)

    # Bottom note
    add_textbox(slide, MARGIN_L, Cm(11.0), SLIDE_W - MARGIN_L - MARGIN_R, Cm(0.8),
                "\u540c\u4e00\u4e2a AI\uff0c\u540c\u4e00\u4e2a\u5f00\u53d1\u8005\uff0c\u7ed3\u679c\u5929\u58e4\u4e4b\u522b",
                14, TEXT_GREY, align=PP_ALIGN.CENTER)

    add_page_number(slide, 16)

    # ══════════════════════════════════════════════════════════
    # P17 - AI's Potential and Boundaries
    # ══════════════════════════════════════════════════════════
    slide = add_slide(prs)
    add_slide_title(slide, "AI \u7684\u6f5c\u529b\u4e0e\u8fb9\u754c")

    # Left: Potential (green)
    lw = Cm(14)
    add_card(slide, MARGIN_L, Cm(3.0), lw, Cm(5.5),
             "\u2705 \u6f5c\u529b",
             ["\u2022 \u7f16\u7801\u6548\u7387\u63d0\u5347 2-4 \u500d",
              "\u2022 \u64c5\u957f\u91cd\u590d\u6027\u5de5\u4f5c\u3001\u6a21\u5f0f\u5316\u4ee3\u7801",
              "\u2022 \u8d8b\u52bf\u660e\u786e\uff0c\u53ea\u4f1a\u8d8a\u6765\u8d8a\u5f3a"],
             title_size=15, desc_size=12, title_color=GREEN,
             accent_color=GREEN, desc_color=TEXT_WHITE)

    # Right: Boundaries (yellow)
    rx = Cm(16.5)
    rw = SLIDE_W - rx - MARGIN_R
    add_card(slide, rx, Cm(3.0), rw, Cm(5.5),
             "\u26a0\ufe0f \u8fb9\u754c",
             ["\u2022 \u4e0d\u4e86\u89e3\u9886\u57df\u80cc\u666f\u65f6\u5bb9\u6613\u8dd1\u504f",
              "\u2022 \u7f3a\u4e4f\u5168\u5c40\u89c6\u89d2\uff0c\u53ef\u80fd\u201c\u5c40\u90e8\u6700\u4f18\u3001\u5168\u5c40\u707e\u96be\u201d",
              "\u2022 \u4e0d\u4f1a\u4e3b\u52a8\u8d28\u7591\u4f60\u7684\u8bbe\u8ba1"],
             title_size=15, desc_size=12, title_color=YELLOW,
             accent_color=YELLOW, desc_color=TEXT_WHITE)

    # Bottom quote
    add_textbox(slide, Cm(2), Cm(10.5), SLIDE_W - Cm(4), Cm(1.5),
                "AI \u662f\u6700\u5f3a\u6267\u884c\u8005\uff0c\u4f46\u4e0d\u662f\u5408\u683c\u7684\u67b6\u6784\u5e08",
                22, ACCENT_BLUE, bold=True, align=PP_ALIGN.CENTER)

    add_page_number(slide, 17)

    # ══════════════════════════════════════════════════════════
    # P18 - Core Insight: Give AI what you're best at
    # ══════════════════════════════════════════════════════════
    slide = add_slide(prs)
    add_slide_title(slide, "\u6838\u5fc3\u6d1e\u5bdf\uff1a\u628a\u6700\u64c5\u957f\u7684\u4e8b\u4ea4\u7ed9 AI")

    # Top comparison
    add_textbox(slide, MARGIN_L, Cm(3.0), SLIDE_W - MARGIN_L - MARGIN_R, Cm(0.8),
                "\u274c \u5e38\u89c1\u60f3\u6cd5\uff1a\u5148\u628a\u4e0d\u64c5\u957f\u7684\u4ea4\u7ed9 AI\uff0c\u81ea\u5df1\u505a\u597d\u64c5\u957f\u7684", 13, RED_ORANGE)
    add_textbox(slide, MARGIN_L, Cm(4.0), SLIDE_W - MARGIN_L - MARGIN_R, Cm(0.8),
                "\u2705 \u5b9e\u9645\u7ecf\u9a8c\uff1a\u5148\u628a\u6700\u64c5\u957f\u7684\u4e8b\u4ea4\u7ed9 AI", 13, GREEN)

    # Three numbered reason cards
    reasons = [
        ("1", "\u64c5\u957f = \u4f60\u80fd\u51c6\u786e\u9a8c\u6536 \u2192 \u8d28\u91cf\u6709\u4fdd\u969c"),
        ("2", "\u8fd9\u4e9b\u5de5\u4f5c\u5bf9\u4e2a\u4eba\u65e0\u6210\u957f\u4ef7\u503c\uff0c\u5c5e\u4e8e\u7eaf\u5de5\u4f5c\u91cf\u5806\u53e0"),
        ("3", "\u91ca\u653e\u65f6\u95f4 \u2192 \u6295\u5165\u5230\u7814\u7a76\u3001\u8bbe\u8ba1\u3001\u67b6\u6784\u7b49\u521b\u9020\u6027\u5de5\u4f5c"),
    ]
    y = Cm(5.5)
    for num, text in reasons:
        add_card(slide, MARGIN_L + Cm(1), y, SLIDE_W - MARGIN_L - MARGIN_R - Cm(2), Cm(1.8),
                 f"  {num}. {text}", [], title_size=13, accent_color=ACCENT_BLUE)
        y += Cm(2.2)

    # Bottom example
    add_textbox(slide, MARGIN_L, Cm(12.5), SLIDE_W - MARGIN_L - MARGIN_R, Cm(0.8),
                "aip-server \u6848\u4f8b\uff1a\u8868\u7ed3\u6784\u662f\u6211\u64c5\u957f\u7684 \u2192 AI \u7f16\u7801\u6211\u80fd\u51c6\u786e\u8bc4\u5ba1",
                11, TEXT_GREY, align=PP_ALIGN.CENTER)

    add_page_number(slide, 18)

    # ══════════════════════════════════════════════════════════
    # P19 - Role Transformation
    # ══════════════════════════════════════════════════════════
    slide = add_slide(prs)
    add_slide_title(slide, "\u89d2\u8272\u8f6c\u53d8")

    # Past flow
    y1 = Cm(3.5)
    add_textbox(slide, MARGIN_L, y1, Cm(3), Cm(0.8), "\u8fc7\u53bb\uff1a", 14, TEXT_GREY, bold=True)
    past_nodes = [
        ("\u8bbe\u8ba1 \u270d\ufe0f", TEXT_GREY),
        ("\u7f16\u7801 \U0001f4bb", TEXT_GREY),
        ("\u6d4b\u8bd5 \U0001f9ea", TEXT_GREY),
    ]
    node_w = Cm(7)
    arrow_w = Cm(1.5)
    sx = Cm(5.5)
    for i, (text, color) in enumerate(past_nodes):
        x = sx + i * (node_w + arrow_w)
        add_flow_node(slide, x, y1, node_w, Cm(1.8), text, "",
                      fill_color=CARD_BG, border_color=BORDER_COLOR,
                      text_color=TEXT_GREY, text_size=12, bold=False)
        if i < 2:
            add_arrow(slide, x + node_w + Cm(0.15), y1 + Cm(0.5), Cm(1.2), Cm(0.6), TEXT_GREY)
    add_textbox(slide, sx + 2 * (node_w + arrow_w) + node_w + Cm(0.3), y1 + Cm(0.3),
                Cm(5), Cm(1.0), "\u5168\u81ea\u5df1\u505a", 11, TEXT_GREY)

    # Now flow
    y2 = Cm(7.0)
    add_textbox(slide, MARGIN_L, y2, Cm(3), Cm(0.8), "\u73b0\u5728\uff1a", 14, ACCENT_BLUE, bold=True)
    now_nodes = [
        ("\u8bbe\u8ba1 \u270d\ufe0f", ACCENT_BLUE),
        ("AI \u7f16\u7801 \U0001f916", GREEN),
        ("\u8bc4\u5ba1\u6d4b\u8bd5 \U0001f50d", ACCENT_BLUE),
    ]
    for i, (text, color) in enumerate(now_nodes):
        x = sx + i * (node_w + arrow_w)
        fill = RGBColor(0x0a, 0x2e, 0x14) if color == GREEN else CARD_BG
        add_flow_node(slide, x, y2, node_w, Cm(1.8), text, "",
                      fill_color=fill, border_color=color,
                      text_color=color, text_size=12, bold=True)
        if i < 2:
            add_arrow(slide, x + node_w + Cm(0.15), y2 + Cm(0.5), Cm(1.2), Cm(0.6), TEXT_GREY)
    add_textbox(slide, sx + 2 * (node_w + arrow_w) + node_w + Cm(0.3), y2 + Cm(0.3),
                Cm(5), Cm(1.0), "AI\u5206\u62c5\u7f16\u7801", 11, GREEN)

    # Bottom quote
    add_textbox(slide, Cm(2), Cm(11.5), SLIDE_W - Cm(4), Cm(1.5),
                "\u4ece\u201c\u5199\u4ee3\u7801\u7684\u4eba\u201d\u53d8\u6210\u201c\u8bbe\u8ba1\u7cfb\u7edf\u7684\u4eba\u201d",
                22, ACCENT_BLUE, bold=True, align=PP_ALIGN.CENTER)

    add_page_number(slide, 19)

    # ══════════════════════════════════════════════════════════
    # P20 - Trap: AI will "PUA" you
    # ══════════════════════════════════════════════════════════
    slide = add_slide(prs)
    add_slide_title(slide, "\u26a0\ufe0f \u9677\u9631\uff1aAI \u4f1a\u201cPUA\u201d\u4f60")

    # Three characteristic cards
    traits = [
        "\u6c38\u8fdc\u81ea\u4fe1\u6ee1\u6ee1\uff0c\u5373\u4f7f\u65b9\u5411\u9519\u8bef",
        "\u8bed\u6c14\u7b03\u5b9a\u3001\u903b\u8f91\u81ea\u6d3d\uff0c\u5f88\u5bb9\u6613\u88ab\u201c\u8bf4\u670d\u201d",
        "\u4e0d\u4f1a\u72b9\u8c6b\uff0c\u4e0d\u4f1a\u8bf4\u201c\u6211\u4e0d\u786e\u5b9a\u201d",
    ]
    card_w = Cm(9.5)
    gap = Cm(0.5)
    for i, trait in enumerate(traits):
        x = MARGIN_L + i * (card_w + gap)
        add_card(slide, x, Cm(3.2), card_w, Cm(3.0), trait, [],
                 title_size=12, title_color=RED_ORANGE, accent_color=RED_ORANGE)

    # Case examples
    add_textbox(slide, MARGIN_L, Cm(7.5), SLIDE_W - MARGIN_L - MARGIN_R, Cm(1.0),
                "\u6848\u4f8b\u56de\u6263\uff1a", 14, TEXT_WHITE, bold=True)
    add_multiline_textbox(slide, MARGIN_L + Cm(1), Cm(8.8), SLIDE_W - MARGIN_L - MARGIN_R - Cm(2), Cm(3.5),
        ["A2A\uff1aAI\u5168\u7a0b\u81ea\u4fe1\u5730\u9020\u4e86\u4e00\u4e2a\u8f6e\u5b50\uff0c46\u6b21\u63d0\u4ea4\uff0c\u6ca1\u6709\u4efb\u4f55\u72b9\u8c6b",
         "\u6743\u9650\u8fc7\u6ee4\uff1aAI\u4fdd\u7559\u8fd9\u4e9b\u903b\u8f91\u65f6\u7684\u7406\u7531\u201c\u5b8c\u5168\u5408\u7406\u201d\uff0c\u4f46\u7f3a\u5931\u4e86\u4e1a\u52a1\u4e0a\u4e0b\u6587"],
        12, TEXT_GREY)

    add_page_number(slide, 20)

    # ══════════════════════════════════════════════════════════
    # P21 - Key Question: How to make AI stable?
    # ══════════════════════════════════════════════════════════
    slide = add_slide(prs)
    add_slide_title(slide, "\u5173\u952e\u95ee\u9898")

    # Progressive chain (5 nodes, arrows down)
    steps = [
        "\u628a\u64c5\u957f\u7684\u4e8b\u4ea4\u7ed9 AI",
        "\u4f46 AI \u4f1a PUA \u4f60",
        "\u9760\u4e2a\u4eba\u5224\u65ad\u529b\u591f\u5417\uff1f",
        "\u5355\u9760\u7ecf\u9a8c \u2192 \u4e0d\u53ef\u6301\u7eed\u3001\u4e0d\u53ef\u590d\u5236",
        "\u56e2\u961f\u534f\u4f5c\uff1f\u9879\u76ee\u590d\u6742\u4e86\u600e\u4e48\u529e\uff1f",
    ]
    node_w = Cm(24)
    node_h = Cm(1.5)
    sx = (SLIDE_W - node_w) // 2
    y = Cm(3.0)
    for i, step in enumerate(steps):
        add_flow_node(slide, sx, y, node_w, node_h, step, "",
                      fill_color=CARD_BG, border_color=BORDER_COLOR,
                      text_color=TEXT_WHITE, text_size=13, bold=False)
        if i < 4:
            add_down_arrow(slide, sx + node_w // 2 - Cm(0.4), y + node_h + Cm(0.05),
                          Cm(0.8), Cm(0.5), TEXT_GREY)
        y += node_h + Cm(0.7)

    # Bottom conclusion
    add_textbox(slide, Cm(2), Cm(13.5), SLIDE_W - Cm(4), Cm(1.5),
                "\u6211\u4eec\u9700\u8981\u4e00\u5957\u7cfb\u7edf\u5316\u7684\u65b9\u6cd5 \u2192",
                24, ACCENT_BLUE, bold=True, align=PP_ALIGN.CENTER)

    add_page_number(slide, 21)
